# SPDX-FileCopyrightText: 2026 DennyWanye
# SPDX-License-Identifier: BUSL-1.1

"""PPT generation tool — python-pptx wrapper with themes + layouts.

Architecture
------------
LLM produces a JSON outline (list of :class:`SlideOutline`); this
module turns it into a real ``.pptx`` file on disk. Three themes and
seven layouts are baked in — the agent picks layouts per slide, never
fiddles with low-level XML.

Themes
~~~~~~
* ``minimal``  — white / blue / charcoal, sans-serif. The default.
* ``dark``     — near-black background, cyan accents, light text.
* ``playful``  — cream background, coral accents, rounded shapes.

Each theme is a dataclass that exposes ``background_rgb``,
``primary_rgb``, ``text_rgb``, ``font_family``. Layouts pick from these
so visuals stay consistent across slides.

Layouts
~~~~~~~
``title`` / ``section`` / ``bullet`` / ``two_column`` / ``image`` /
``quote`` / ``toc``. Each is a small function in this module that takes
a :class:`pptx.Presentation` slide + the outline data and lays it out.

Failure modes
-------------
* ``python-pptx`` not installed → :func:`ppt_create` returns
  ``{"error": ..., "markdown_fallback": "<best-effort .md>"}``. Caller
  (the agent) can still hand the user a readable outline.
* Output path unwritable / invalid → same fallback.
* Image path missing → that slide degrades to a bullet layout with a
  warning embedded as a footnote text.

The whole module is purely synchronous and CPU-cheap — no LLM calls
except optional image generation. Producing the outline is the LLM's job,
layout is ours. By default this stays offline; only slides with
``image_prompt`` trigger synchronous ``image_tools.generate_images`` calls,
and failures degrade to placeholders.
"""
from __future__ import annotations

import asyncio
import copy
import inspect
import json
import logging
import os
import re
import tempfile
import time
from dataclasses import dataclass, field, asdict
from pathlib import Path
from types import SimpleNamespace
from typing import Any, Iterable, Literal, Optional, Sequence

from deskpet.tools import ppt_outline_store
from deskpet.tools.image_tools import generate_images, probe_image_reachable

log = logging.getLogger(__name__)


_PPT_PRO_CTX: dict[str, Any] = {}
_PPT_PRO_TASKS: set[asyncio.Task] = set()
_PPT_PRO_RUNNING: dict[str, asyncio.Task] = {}
_PPT_PRO_RUNNING_TOPIC: dict[str, str] = {}


def _ppt_config_section() -> dict[str, Any]:
    """独立工具进程读取 ``[ppt]`` 配置；读不到时返回空 dict。"""
    merged: dict[str, Any] = {}
    # 兼容旧测试和历史内存注入形态；真实生产配置以后面的 standalone 读取为准。
    try:
        import config as _cfg  # type: ignore[import-not-found]

        raw = getattr(getattr(_cfg, "config", None), "raw", {}) or {}
        val = raw.get("ppt") if isinstance(raw, dict) else {}
        if isinstance(val, dict):
            merged.update(val)
    except Exception:  # noqa: BLE001
        pass
    try:
        from config import standalone_config_section  # type: ignore[import-not-found]

        raw = standalone_config_section("ppt") or {}
        if isinstance(raw, dict):
            merged.update(raw)
    except Exception:  # noqa: BLE001
        pass
    return merged


def _cfg_bool(raw: dict[str, Any], key: str, default: bool) -> bool:
    val = raw.get(key, default)
    if isinstance(val, bool):
        return val
    if isinstance(val, str):
        return val.strip().lower() in {"1", "true", "yes", "on"}
    return bool(val)


def _cfg_int(raw: dict[str, Any], key: str, default: int) -> int:
    try:
        return int(raw.get(key, default))
    except (TypeError, ValueError):
        return default


def _cfg_float(raw: dict[str, Any], key: str, default: float) -> float:
    try:
        return float(raw.get(key, default))
    except (TypeError, ValueError):
        return default


def _cfg_str(raw: dict[str, Any], key: str, default: str) -> str:
    val = raw.get(key, default)
    return str(val).strip() if val is not None else default


def _ppt_pro_cfg() -> SimpleNamespace:
    """读取 PPT Pro 配置，兼容 ``pro_*`` 与短字段，缺失时使用计划默认值。"""
    raw = _ppt_config_section()
    cfg = SimpleNamespace(
        enabled=_cfg_bool(raw, "pro_enabled", _cfg_bool(raw, "enabled", True)),
        default_depth=_cfg_str(raw, "pro_default_depth", _cfg_str(raw, "default_depth", "deep")),
        max_revisions=_cfg_int(raw, "pro_max_revisions", _cfg_int(raw, "max_revisions", 2)),
        research_timeout_s=_cfg_float(raw, "pro_research_timeout_s", _cfg_float(raw, "research_timeout_s", 360.0)),
        confirm_timeout_s=_cfg_float(raw, "pro_confirm_timeout_s", _cfg_float(raw, "confirm_timeout_s", 1800.0)),
        image_probe_timeout_s=_cfg_float(raw, "pro_image_probe_timeout_s", _cfg_float(raw, "image_probe_timeout_s", 8.0)),
        render_timeout_s=_cfg_float(raw, "pro_render_timeout_s", _cfg_float(raw, "render_timeout_s", 0.0)),
        save_research=_cfg_bool(raw, "pro_save_research", _cfg_bool(raw, "save_research", True)),
        outline_history=_cfg_bool(raw, "pro_outline_history", _cfg_bool(raw, "outline_history", True)),
    )
    # 编排骨架里曾使用 pro_* 命名；这里提供别名，避免下一趟接线重复适配。
    cfg.pro_enabled = cfg.enabled
    cfg.pro_default_depth = cfg.default_depth
    cfg.pro_max_revisions = cfg.max_revisions
    cfg.pro_research_timeout_s = cfg.research_timeout_s
    cfg.pro_confirm_timeout_s = cfg.confirm_timeout_s
    cfg.pro_image_probe_timeout_s = cfg.image_probe_timeout_s
    cfg.pro_render_timeout_s = cfg.render_timeout_s
    cfg.pro_save_research = cfg.save_research
    cfg.pro_outline_history = cfg.outline_history
    return cfg


def _in_pytest() -> bool:
    # 预览渲染会真启动 WPS(慢且有副作用),pytest 全程跳过;
    # 专门的 A-5 接线测试会 monkeypatch 本函数放行。
    return bool(os.environ.get("PYTEST_CURRENT_TEST"))


def _ppt_preview_render_enabled() -> bool:
    """config ``[ppt].preview_render``(默认 True)。读法与 _image_model 一致。"""
    try:
        return _cfg_bool(_ppt_config_section(), "preview_render", True)
    except Exception:  # noqa: BLE001
        return True


# lazy 渲染器句柄：生产首次用时才 import(COM 探测有开销);
# 测试 monkeypatch ppt_tools.ppt_render 即可注入假渲染器。
ppt_render = None  # type: ignore[assignment]


def _get_ppt_renderer():
    global ppt_render
    if ppt_render is None:
        from . import ppt_render as _mod  # noqa: PLW0603
        ppt_render = _mod
    return ppt_render


def _maybe_render_preview(result: dict[str, Any]) -> None:
    """Best-effort PPT 视觉预览：每页渲染成 PNG 追加为 image artifacts,
    用户/agent 在聊天里直接看到每页效果。失败只记日志,不影响主结果。"""
    try:
        if _in_pytest():
            return
        if not _ppt_preview_render_enabled():
            return
        if not result.get("ok") or not result.get("path"):
            return

        renderer = _get_ppt_renderer()
        if not renderer.com_render_available():
            return

        pptx = Path(str(result["path"])).expanduser().resolve()
        out_dir = pptx.with_suffix(".preview")
        # 子进程+超时渲染: WPS COM 对某页挂死也不会阻塞调用线程(异步图文
        # PPT 任务靠它返回后才推"做好啦",in-process 挂死会让推回永久丢失)。
        render_fn = getattr(renderer, "render_pptx_to_pngs_safe", None)
        if callable(render_fn):
            pngs = render_fn(str(pptx), str(out_dir), timeout=150.0)
        else:
            pngs = renderer.render_pptx_to_pngs(str(pptx), str(out_dir))
        artifacts = result.setdefault("artifacts", [])
        for idx, png in enumerate(pngs, start=1):
            artifacts.append({
                "kind": "image",
                "path": png,
                "mime": "image/png",
                "title": f"预览 第{idx}页",
            })
    except Exception as exc:  # noqa: BLE001
        log.warning("ppt preview render failed: %s", exc, exc_info=True)


# ---------------------------------------------------------------------
# python-pptx availability — defer the import so callers without the
# dep get a graceful fallback instead of an ImportError at module load.
# ---------------------------------------------------------------------
try:  # pragma: no cover — import probe
    from pptx import Presentation as _Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.dml import MSO_FILL_TYPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    _HAS_PPTX = True
except ImportError:  # pragma: no cover
    _Presentation = None  # type: ignore
    Inches = Pt = Emu = RGBColor = MSO_SHAPE = PP_PLACEHOLDER = PP_ALIGN = MSO_ANCHOR = MSO_FILL_TYPE = MSO_AUTO_SIZE = None  # type: ignore
    CategoryChartData = XL_CHART_TYPE = XL_LEGEND_POSITION = None  # type: ignore
    _HAS_PPTX = False


VALID_LAYOUTS = (
    "title", "section", "bullet", "two_column", "image", "image_full", "quote", "toc", "chart",
)
VALID_THEMES = ("minimal", "dark", "playful")
# 模板源已迁到外部大库 resources/PPT_Template(见 ppt_template_picker)。
# 不再有 git 跟踪的 bundled 模板目录。


# ---------------------------------------------------------------------
# Public dataclasses
# ---------------------------------------------------------------------


@dataclass
class SlideOutline:
    """One slide spec. LLM produces a list of these as JSON."""

    layout: str = "bullet"
    title: str = ""
    subtitle: str = ""
    bullets: list[str] = field(default_factory=list)
    # two_column
    left: list[str] = field(default_factory=list)
    right: list[str] = field(default_factory=list)
    left_title: str = ""
    right_title: str = ""
    # image
    image_path: Optional[str] = None
    # 整页生图：有则 ppt_create 自动生图填 image_path
    image_prompt: Optional[str] = None
    # 整页生图版式变体(版式多样化): cover/split_left/split_right/top/card/
    # quote。空 = 由 _assign_image_layouts 按内容自动分配(相邻不重复)。
    image_variant: str = ""
    # 视觉评审闭环的 shrink_text 动作: 渲染字号 × 此系数(默认 1.0)。
    font_scale: float = 1.0
    caption: str = ""
    # quote
    quote: str = ""
    cite: str = ""
    # toc — uses bullets[]
    # chart: {"type": "bar"|"line"|"pie", "categories": [...],
    #         "series": [{"name": "...", "values": [...]}, ...]}
    chart: Optional[dict] = None
    # generic
    notes: str = ""  # speaker notes

    def normalize(self) -> "SlideOutline":
        """Coerce raw LLM output into the strict schema."""
        layout = (self.layout or "bullet").strip().lower()
        if layout not in VALID_LAYOUTS:
            log.debug("unknown layout %r → falling back to bullet", layout)
            layout = "bullet"
        return SlideOutline(
            layout=layout,
            title=(self.title or "").strip(),
            subtitle=(self.subtitle or "").strip(),
            bullets=[str(b).strip() for b in (self.bullets or []) if str(b).strip()],
            left=[str(b).strip() for b in (self.left or []) if str(b).strip()],
            right=[str(b).strip() for b in (self.right or []) if str(b).strip()],
            left_title=(self.left_title or "").strip(),
            right_title=(self.right_title or "").strip(),
            image_path=(self.image_path or None),
            image_prompt=(self.image_prompt or None),
            image_variant=(self.image_variant or "").strip().lower(),
            font_scale=(
                float(self.font_scale)
                if isinstance(self.font_scale, (int, float)) and 0.5 <= float(self.font_scale) <= 1.5
                else 1.0
            ),
            caption=(self.caption or "").strip(),
            quote=(self.quote or "").strip(),
            cite=(self.cite or "").strip(),
            chart=(self.chart if isinstance(self.chart, dict) else None),
            notes=(self.notes or "").strip(),
        )


# ---------------------------------------------------------------------
# Themes
# ---------------------------------------------------------------------


@dataclass(frozen=True)
class Theme:
    name: str
    background_rgb: tuple[int, int, int]
    dark_bg_rgb: tuple[int, int, int]
    surface_rgb: tuple[int, int, int]
    primary_rgb: tuple[int, int, int]
    secondary_rgb: tuple[int, int, int]
    accent_rgb: tuple[int, int, int]
    highlight_rgb: tuple[int, int, int]
    text_rgb: tuple[int, int, int]
    muted_text_rgb: tuple[int, int, int]
    dark_text_rgb: tuple[int, int, int]
    dark_muted_text_rgb: tuple[int, int, int]
    font_heading: str
    font_body: str


_THEMES: dict[str, Theme] = {
    "minimal": Theme(
        name="minimal",
        background_rgb=(255, 255, 255),
        dark_bg_rgb=(8, 20, 38),       # deep navy
        surface_rgb=(244, 248, 250),   # cool off-white
        primary_rgb=(11, 30, 48),
        secondary_rgb=(16, 123, 136),  # teal
        accent_rgb=(245, 158, 11),     # amber highlight
        highlight_rgb=(20, 184, 166),
        text_rgb=(15, 23, 42),
        muted_text_rgb=(84, 98, 113),
        dark_text_rgb=(248, 250, 252),
        dark_muted_text_rgb=(186, 211, 224),
        font_heading="Microsoft YaHei",
        font_body="Microsoft YaHei",
    ),
    "dark": Theme(
        name="dark",
        background_rgb=(255, 255, 255),
        dark_bg_rgb=(3, 7, 18),
        surface_rgb=(241, 245, 249),
        primary_rgb=(15, 23, 42),
        secondary_rgb=(79, 70, 229),
        accent_rgb=(34, 211, 238),
        highlight_rgb=(251, 191, 36),
        text_rgb=(15, 23, 42),
        muted_text_rgb=(71, 85, 105),
        dark_text_rgb=(248, 250, 252),
        dark_muted_text_rgb=(148, 163, 184),
        font_heading="Microsoft YaHei",
        font_body="Microsoft YaHei",
    ),
    "playful": Theme(
        name="playful",
        background_rgb=(255, 255, 255),
        dark_bg_rgb=(49, 46, 129),     # confident indigo
        surface_rgb=(248, 250, 252),
        primary_rgb=(30, 41, 59),
        secondary_rgb=(20, 184, 166),
        accent_rgb=(239, 68, 68),
        highlight_rgb=(245, 158, 11),
        text_rgb=(30, 41, 59),
        muted_text_rgb=(92, 107, 123),
        dark_text_rgb=(255, 255, 255),
        dark_muted_text_rgb=(221, 214, 254),
        font_heading="Microsoft YaHei",
        font_body="Microsoft YaHei",
    ),
}


def get_theme(name: str) -> Theme:
    return _THEMES.get(name.lower(), _THEMES["minimal"])


# ---------------------------------------------------------------------
# Outline parsing
# ---------------------------------------------------------------------


def parse_outline(raw: Any) -> list[SlideOutline]:
    """Tolerant LLM-output → list[SlideOutline] converter.

    Accepts:
      * JSON string of a list
      * JSON string wrapped in code fences
      * Python list of dicts
      * Single dict (treated as one-slide deck)

    Each slide dict is fed through :class:`SlideOutline` constructor
    (unknown keys are silently ignored). Empty / unparseable input
    returns ``[]``.
    """
    if raw is None:
        return []
    data: Any = raw
    if isinstance(data, str):
        text = data.strip()
        if not text:
            return []
        # Strip code fences
        if text.startswith("```"):
            nl = text.find("\n")
            if nl != -1:
                text = text[nl + 1:]
            if text.endswith("```"):
                text = text[:-3]
            text = text.strip()
        # Try to extract first JSON array / object
        lb_arr, rb_arr = text.find("["), text.rfind("]")
        lb_obj, rb_obj = text.find("{"), text.rfind("}")
        slice_arr = text[lb_arr:rb_arr + 1] if 0 <= lb_arr < rb_arr else None
        slice_obj = text[lb_obj:rb_obj + 1] if 0 <= lb_obj < rb_obj else None
        for candidate in (slice_arr, slice_obj, text):
            if not candidate:
                continue
            try:
                data = json.loads(candidate)
                break
            except json.JSONDecodeError:
                continue
        else:
            return []
    if isinstance(data, dict):
        data = [data]
    if not isinstance(data, list):
        return []
    out: list[SlideOutline] = []
    for item in data:
        if not isinstance(item, dict):
            continue
        try:
            so = SlideOutline(**{
                k: v for k, v in item.items()
                if k in SlideOutline.__dataclass_fields__
            })
        except (TypeError, ValueError):
            continue
        out.append(so.normalize())
    return out


# ---------------------------------------------------------------------
# PPT Pro content helpers (WI-1 / WI-2)
# ---------------------------------------------------------------------


_PPT_PRO_RESEARCH_CTX_CHARS = 6000


def _is_config_or_auth_error(e: Exception) -> bool:
    """区分必须上抛的配置/认证硬错误，避免静默降级成无来源 PPT。"""
    status = getattr(e, "status_code", None) or getattr(e, "status", None)
    response = getattr(e, "response", None)
    if status is None and response is not None:
        status = getattr(response, "status_code", None) or getattr(response, "status", None)
    try:
        if int(status) in {401, 403}:
            return True
    except (TypeError, ValueError):
        pass

    cls_name = type(e).__name__.lower()
    msg = str(e).lower()
    text = f"{cls_name} {msg}"
    hard_markers = (
        "401",
        "403",
        "unauthorized",
        "forbidden",
        "authentication",
        "auth",
        "api key",
        "apikey",
        "invalid key",
        "missing key",
        "no llm provider configured",
        "provider modules unavailable",
        "provider not configured",
        "no provider configured",
        "not configured",
        "未配置",
        "认证",
        "鉴权",
        "授权",
        "密钥",
        "api_key",
    )
    return any(marker in text for marker in hard_markers)


async def _research_topic_for_ppt(topic: str, *, depth: str, timeout_s: float):
    """为 PPT Pro 跑 DeepResearch；软失败返回 None，认证/配置错误上抛。"""
    from deskpet.tools.research_tools import (  # type: ignore[import-not-found]
        _DEPTH_PRESETS,
        _resolve_default_llm_call,
        deepresearch,
    )

    chosen_depth = (depth or "standard").strip().lower()
    d_subq, d_urls, d_pass, d_rounds = _DEPTH_PRESETS.get(
        chosen_depth, _DEPTH_PRESETS.get("standard", (5, 4, 12, 1))
    )

    async def _run():
        llm_call = await _resolve_default_llm_call()
        return await deepresearch(
            topic,
            llm_call=llm_call,
            max_sub_questions=d_subq,
            max_urls_per_query=d_urls,
            max_total_passages=d_pass,
            max_rounds=d_rounds,
            mode=chosen_depth,
            user_request=f"为制作PPT调研：{topic}",
        )

    try:
        return await asyncio.wait_for(_run(), timeout=timeout_s)
    except asyncio.TimeoutError:
        log.warning("ppt_pro research timeout -> degrade to no-research outline")
        return None
    except Exception as exc:  # noqa: BLE001
        if _is_config_or_auth_error(exc):
            raise
        log.warning("ppt_pro research failed: %s -> degrade to no-research outline", exc)
        return None


def _save_and_index_research(topic: str, report: Any):
    """Best-effort 落盘调研报告并更新 DeepResearch 索引；失败只记录日志。"""
    try:
        from deskpet.tools.research_tools import (  # type: ignore[import-not-found]
            _save_report,
            _update_deepresearch_index,
        )

        saved = _save_report(topic, report)
        if not saved:
            return None

        async def _update() -> None:
            await _update_deepresearch_index(saved, topic, report)

        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            asyncio.run(_update())
        else:
            task = loop.create_task(_update())

            def _log_done(t: asyncio.Task) -> None:
                try:
                    t.result()
                except Exception as exc:  # noqa: BLE001
                    log.debug("ppt_pro research index update skipped: %s", exc)

            task.add_done_callback(_log_done)
        return saved
    except Exception as exc:  # noqa: BLE001
        log.debug("ppt_pro research save skipped: %s", exc)
        return None


async def _draft_outline_from_research(
    topic: str,
    report: Any,
    *,
    pages: int,
    theme: str,
    image_mode: bool,
    llm_call,
    feedback: str = "",
    prev_slides: Optional[Sequence[SlideOutline]] = None,
) -> list[SlideOutline]:
    """基于调研报告拟制 SlideOutline；LLM 输出异常时用本地最小大纲兜底。"""
    _ = theme
    research_md = (getattr(report, "report_md", "") if report else "")[:_PPT_PRO_RESEARCH_CTX_CHARS]
    prompt = _build_outline_prompt(
        topic,
        research_md,
        pages=pages,
        image_mode=image_mode,
        feedback=feedback,
        prev_md=(_outline_to_markdown(prev_slides) if prev_slides else ""),
    )
    raw = await llm_call(prompt)
    slides = parse_outline(raw)
    if not slides:
        slides = _fallback_minimal_outline(topic, pages, image_mode)
    return slides


def _build_outline_prompt(
    topic: str,
    research_md: str,
    *,
    pages: int,
    image_mode: bool,
    feedback: str = "",
    prev_md: str = "",
) -> str:
    """构造兼容 SlideOutline JSON 数组的拟纲提示词。"""
    pages = max(1, int(pages or 1))
    mode_hint = (
        "优先使用 image_full 布局；每页仍必须有 title、3-5 条 bullets 和英文 image_prompt。"
        if image_mode
        else "优先使用 title、section、bullet、two_column 等可编辑布局；每页仍必须有 image_prompt 供后续可选生图。"
    )
    research_block = research_md.strip() or "调研报告为空：必须明确标注待核验，不得伪造来源、数据或引用。"
    revise_block = ""
    if prev_md.strip():
        revise_block = (
            "\n这是当前大纲和用户修改意见。请做增量修改：只改与反馈相关的页，保留其余页的标题、顺序和核心要点。\n"
            f"\n当前大纲：\n{prev_md.strip()}\n"
        )
    feedback_block = f"\n用户修改意见：{feedback.strip()}\n" if feedback.strip() else ""

    return (
        "你是严谨的 PPT 内容策划助手。请只输出 JSON 数组，不要 Markdown、解释或代码围栏。\n"
        f"主题：{topic}\n"
        f"页数：{pages}\n"
        f"视觉模式：{'AI 整页图文' if image_mode else '可编辑内容页'}。{mode_hint}\n"
        "输出必须兼容 SlideOutline：每个对象可含 layout/title/subtitle/bullets/image_prompt/notes。"
        "每页必须同时包含 title、bullets、image_prompt；bullets 为 3-5 条充实要点，避免空泛短词。\n"
        "内容硬约束：必须基于调研报告组织论点，优先引用报告中的数据、事实、趋势和来源编号；"
        "不得编造报告没有支持的数字、机构、结论或引用。没有报告时要写成待核验假设。\n"
        "image_prompt 用英文，要求 cinematic, negative space, editorial composition，并包含 no text, no watermark, no logo, no UI labels。"
        "不要让图片里出现文字、商标、水印或界面字。\n"
        f"{revise_block}{feedback_block}"
        "\n调研报告：\n"
        f"{research_block}\n"
        "\nJSON 示例：\n"
        "[{\"layout\":\"image_full\",\"title\":\"标题\",\"bullets\":[\"要点一\",\"要点二\",\"要点三\"],"
        "\"image_prompt\":\"cinematic scene, negative space, no text, no watermark, no logo\"}]\n"
    )


def _outline_to_markdown(slides: Sequence[SlideOutline]) -> str:
    """把 SlideOutline 转成确认卡可读 Markdown。"""
    lines: list[str] = []
    for i, slide in enumerate(slides or [], start=1):
        title = slide.title or "(无标题)"
        lines.append(f"第 {i} 页：{title}")
        if slide.subtitle:
            lines.append(f"  _{slide.subtitle}_")
        for bullet in slide.bullets or []:
            lines.append(f"  - {bullet}")
        if slide.left_title or slide.left:
            lines.append(f"  {slide.left_title or '左栏'}")
            for bullet in slide.left:
                lines.append(f"  - {bullet}")
        if slide.right_title or slide.right:
            lines.append(f"  {slide.right_title or '右栏'}")
            for bullet in slide.right:
                lines.append(f"  - {bullet}")
        if slide.quote:
            lines.append(f"  > {slide.quote}")
        if slide.cite:
            lines.append(f"  -- {slide.cite}")
        if i != len(slides):
            lines.append("")
    return "\n".join(lines).strip()


def _fallback_minimal_outline(topic: str, pages: int, image_mode: bool) -> list[SlideOutline]:
    """LLM 不可用时的本地最小大纲：封面 + 内容页占位。"""
    total = max(1, int(pages or 1))
    slides: list[SlideOutline] = [
        SlideOutline(
            layout="title",
            title=topic or "PPT Pro",
            subtitle="基于当前可用信息生成的初稿",
            image_prompt=(
                f"cinematic keynote cover about {topic}, negative space, no text, no watermark, no logo"
                if image_mode else None
            ),
        )
    ]
    for i in range(2, total + 1):
        slides.append(
            SlideOutline(
                layout="image_full" if image_mode else "bullet",
                title=f"{topic}：关键问题 {i - 1}" if topic else f"关键问题 {i - 1}",
                bullets=[
                    "补充来自调研报告的核心事实",
                    "提炼对用户目标有用的判断",
                    "标注仍需核验的数据和来源",
                ],
                image_prompt=(
                    f"cinematic editorial visual about {topic}, negative space, no text, no watermark, no logo"
                    if image_mode else None
                ),
            )
        )
    return slides


# ---------------------------------------------------------------------
# Markdown fallback (when python-pptx unavailable)
# ---------------------------------------------------------------------


def render_markdown_fallback(
    outline: Sequence[SlideOutline], *, title: str = "", author: str = "",
) -> str:
    """Render a slide outline as a Markdown document for the
    no-python-pptx fallback path. Used by callers as the safe ``markdown_fallback``
    field in the error result.
    """
    lines: list[str] = []
    if title:
        lines.append(f"# {title}")
    if author:
        lines.append(f"_作者: {author}_\n")
    for i, slide in enumerate(outline, start=1):
        lines.append(f"\n## 第 {i} 页 — {slide.title or '(无标题)'}")
        if slide.subtitle:
            lines.append(f"_{slide.subtitle}_")
        if slide.layout == "two_column":
            if slide.left_title or slide.left:
                lines.append(f"\n**左 — {slide.left_title}**")
                for b in slide.left:
                    lines.append(f"- {b}")
            if slide.right_title or slide.right:
                lines.append(f"\n**右 — {slide.right_title}**")
                for b in slide.right:
                    lines.append(f"- {b}")
        elif slide.layout == "quote":
            lines.append(f"\n> {slide.quote}")
            if slide.cite:
                lines.append(f"\n— {slide.cite}")
        elif slide.layout == "image":
            if slide.image_path:
                lines.append(f"\n![{slide.caption}]({slide.image_path})")
            if slide.caption:
                lines.append(f"\n_{slide.caption}_")
        else:
            for b in slide.bullets:
                lines.append(f"- {b}")
        if slide.notes:
            lines.append(f"\n> 备注: {slide.notes}")
    return "\n".join(lines).strip() + "\n"


# ---------------------------------------------------------------------
# Builder
# ---------------------------------------------------------------------


# 16:9 default canvas. EMU = English Metric Units; 914400 = 1 inch.
_SLIDE_WIDTH = 9144000   # 10 inches
_SLIDE_HEIGHT = 5143500  # 5.625 inches → 16:9


def _rgb(triple: tuple[int, int, int]):
    return RGBColor(*triple)


def _fill_slide_bg(
    slide, theme: Theme, color: tuple[int, int, int] | None = None,
) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color or theme.background_rgb)


def _estimate_capacity(width: int, height: int, font_size: int) -> int:
    width_in = max(width / 914400, 0.5)
    height_in = max(height / 914400, 0.25)
    chars_per_line = max(int(width_in * 12.5 * 16 / max(font_size, 8)), 8)
    lines = max(int(height_in * 72 / max(font_size * 1.3, 1)), 1)
    return max(chars_per_line * lines, 12)


def _shorten(text: str, max_chars: int) -> str:
    text = str(text or "").strip()
    if len(text) <= max_chars:
        return text
    return text[:max(max_chars - 3, 1)].rstrip() + "..."


def _fit_font_size(text: str, width: int, height: int, desired: int, minimum: int = 10) -> int:
    size = desired
    while size > minimum and len(str(text)) > _estimate_capacity(width, height, size):
        size -= 1
    return size


def _add_shape(
    slide,
    shape_type,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    fill: tuple[int, int, int],
    line: tuple[int, int, int] | None = None,
):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb(line)
    return shape


def _set_fill_alpha(shape, alpha_pct: int) -> None:
    """给已 solid 填充的形状加透明度(0=全透,100=不透)。python-pptx 无原生
    alpha API,直接往 srgbClr 注入 <a:alpha>。半透明深色遮罩让底图透出来
    (引用页/卡片可读又不死黑)。从不抛异常。"""
    try:
        from pptx.oxml.ns import qn

        srgb = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
        if srgb is None:
            return
        for old in srgb.findall(qn("a:alpha")):
            srgb.remove(old)
        a = srgb.makeelement(qn("a:alpha"), {"val": str(int(max(0, min(100, alpha_pct)) * 1000))})
        srgb.append(a)
    except Exception:  # noqa: BLE001
        pass


def _add_text(
    slide,
    text: str,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    font_size: int,
    bold: bool = False,
    color: tuple[int, int, int] = (0, 0, 0),
    font_name: str = "Microsoft YaHei UI",
    align: str = "left",
    anchor: str = "top",
    margin: float = 0.05,
):
    font_size = _fit_font_size(text, width, height, font_size)
    text = _shorten(text, _estimate_capacity(width, height, font_size))
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(anchor, MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font_name
    f.size = Pt(font_size)
    f.bold = bold
    f.color.rgb = _rgb(color)
    return tb


def _add_bullet_text(
    slide,
    bullets: Iterable[str],
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    font_size: int,
    color: tuple[int, int, int],
    font_name: str,
    accent_color: tuple[int, int, int] | None = None,
):
    """Render bullets with a coloured dot prefix.

    python-pptx doesn't expose first-class bullet formatting reliably
    across themes, so we draw the bullet character ourselves (●)
    coloured with the accent colour.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    bullets = list(bullets)
    if not bullets:
        # Empty — leave an invisible placeholder so layout stays stable.
        tf.paragraphs[0].text = ""
        return tb
    max_each = max(_estimate_capacity(width, height // max(len(bullets), 1), font_size) - 4, 18)
    for i, line in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(8)
        dot = para.add_run()
        dot.text = "● "
        dot.font.size = Pt(font_size)
        dot.font.name = font_name
        dot.font.color.rgb = _rgb(accent_color or color)
        body = para.add_run()
        body.text = _shorten(str(line), max_each)
        body.font.size = Pt(font_size)
        body.font.name = font_name
        body.font.color.rgb = _rgb(color)
    return tb


def _add_title_block(slide, title: str, subtitle: str, theme: Theme) -> None:
    _add_text(
        slide, title,
        left=Inches(0.62), top=Inches(0.36),
        width=Inches(7.4), height=Inches(0.68),
        font_size=36, bold=True,
        color=theme.primary_rgb,
        font_name=theme.font_heading,
    )
    if subtitle:
        _add_text(
            slide, subtitle,
            left=Inches(0.64), top=Inches(0.95),
            width=Inches(7.7), height=Inches(0.36),
            font_size=12,
            color=theme.muted_text_rgb,
            font_name=theme.font_body,
        )


def _add_corner_motif(slide, theme: Theme, *, dark: bool = False) -> None:
    base = theme.secondary_rgb if dark else theme.surface_rgb
    _add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(7.78), top=Inches(0.34),
        width=Inches(1.45), height=Inches(0.82),
        fill=base,
    )
    _add_shape(
        slide, MSO_SHAPE.OVAL,
        left=Inches(8.92), top=Inches(0.70),
        width=Inches(0.34), height=Inches(0.34),
        fill=theme.accent_rgb,
    )


def _add_icon_badge(
    slide,
    label: str,
    *,
    left: int,
    top: int,
    size: int,
    theme: Theme,
    fill: tuple[int, int, int] | None = None,
    text_color: tuple[int, int, int] | None = None,
):
    _add_shape(
        slide, MSO_SHAPE.OVAL,
        left=left, top=top, width=size, height=size,
        fill=fill or theme.secondary_rgb,
    )
    return _add_text(
        slide, label,
        left=left, top=top + Inches(0.01),
        width=size, height=size,
        font_size=10, bold=True,
        color=text_color or (255, 255, 255),
        font_name="Calibri",
        align="center", anchor="middle",
        margin=0.0,
    )


_NUMBER_RE = re.compile(r"(?<![\w.])(\d+(?:\.\d+)?\s?(?:%|x|X|倍|万|亿|年|小时|分钟|人|页)?)")


def _extract_callouts(items: Sequence[str], limit: int = 2) -> tuple[list[tuple[str, str]], list[str]]:
    callouts: list[tuple[str, str]] = []
    remaining: list[str] = []
    for item in items:
        text = str(item).strip()
        match = _NUMBER_RE.search(text)
        if match and len(callouts) < limit:
            value = match.group(1).replace(" ", "")
            label = (text[:match.start()] + text[match.end():]).strip(" :-：，,。")
            callouts.append((value, label or text))
        else:
            remaining.append(text)
    return callouts, remaining


def _add_callout_card(
    slide,
    value: str,
    label: str,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    theme: Theme,
):
    _add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left=left, top=top, width=width, height=height,
        fill=theme.primary_rgb,
    )
    _add_shape(
        slide, MSO_SHAPE.OVAL,
        left=left + width - Inches(0.58), top=top + Inches(0.18),
        width=Inches(0.32), height=Inches(0.32),
        fill=theme.accent_rgb,
    )
    _add_text(
        slide, value,
        left=left + Inches(0.28), top=top + Inches(0.18),
        width=width - Inches(0.55), height=Inches(0.72),
        font_size=66, bold=True,
        color=theme.dark_text_rgb,
        font_name="Calibri",
    )
    _add_text(
        slide, label,
        left=left + Inches(0.32), top=top + Inches(1.02),
        width=width - Inches(0.64), height=height - Inches(1.12),
        font_size=13,
        color=theme.dark_muted_text_rgb,
        font_name=theme.font_body,
    )


def _add_item_card(
    slide,
    title: str,
    body: str,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    theme: Theme,
    badge: str,
):
    _add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left=left, top=top, width=width, height=height,
        fill=theme.surface_rgb,
        line=(226, 232, 240),
    )
    _add_icon_badge(
        slide, badge,
        left=left + Inches(0.22), top=top + Inches(0.22),
        size=Inches(0.42), theme=theme,
    )
    _add_text(
        slide, title,
        left=left + Inches(0.78), top=top + Inches(0.20),
        width=width - Inches(1.0), height=Inches(0.34),
        font_size=16, bold=True,
        color=theme.primary_rgb,
        font_name=theme.font_heading,
    )
    _add_text(
        slide, body,
        left=left + Inches(0.28), top=top + Inches(0.76),
        width=width - Inches(0.56), height=height - Inches(0.88),
        font_size=14,
        color=theme.text_rgb,
        font_name=theme.font_body,
    )


def _add_accent_bar(slide, theme: Theme, *, top: int = 0) -> None:
    """Decorative left-edge accent bar for non-title slides."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Emu(top),
        Inches(0.18), Emu(_SLIDE_HEIGHT - top),
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(theme.accent_rgb)


# ---- Per-layout renderers ------------------------------------------


def _render_title(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    # Big horizontal accent bar at top-third
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(2.2),
        Inches(10), Inches(0.08),
    )
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(theme.accent_rgb)
    _add_text(
        slide, outline.title or "(无标题)",
        left=Inches(0.6), top=Inches(1.3),
        width=Inches(8.8), height=Inches(0.9),
        font_size=44, bold=True,
        color=theme.primary_rgb,
        font_name=theme.font_heading,
        align="left", anchor="top",
    )
    if outline.subtitle:
        _add_text(
            slide, outline.subtitle,
            left=Inches(0.6), top=Inches(2.4),
            width=Inches(8.8), height=Inches(0.6),
            font_size=20,
            color=theme.muted_text_rgb,
            font_name=theme.font_body,
            align="left",
        )


def _render_section(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    # Centered large title
    _add_text(
        slide, outline.title or "(章节)",
        left=Inches(0.5), top=Inches(2.0),
        width=Inches(9), height=Inches(1.0),
        font_size=40, bold=True,
        color=theme.accent_rgb,
        font_name=theme.font_heading,
        align="center", anchor="middle",
    )
    if outline.subtitle:
        _add_text(
            slide, outline.subtitle,
            left=Inches(0.5), top=Inches(3.0),
            width=Inches(9), height=Inches(0.6),
            font_size=18,
            color=theme.muted_text_rgb,
            font_name=theme.font_body,
            align="center",
        )


def _render_bullet(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    _add_accent_bar(slide, theme)
    _add_text(
        slide, outline.title or "(无标题)",
        left=Inches(0.6), top=Inches(0.4),
        width=Inches(8.8), height=Inches(0.7),
        font_size=28, bold=True,
        color=theme.primary_rgb,
        font_name=theme.font_heading,
    )
    if outline.subtitle:
        _add_text(
            slide, outline.subtitle,
            left=Inches(0.6), top=Inches(1.05),
            width=Inches(8.8), height=Inches(0.4),
            font_size=14,
            color=theme.muted_text_rgb,
            font_name=theme.font_body,
        )
    _add_bullet_text(
        slide, outline.bullets,
        left=Inches(0.7), top=Inches(1.6),
        width=Inches(8.6), height=Inches(3.8),
        font_size=18,
        color=theme.text_rgb,
        font_name=theme.font_body,
        accent_color=theme.accent_rgb,
    )


def _render_two_column(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    _add_accent_bar(slide, theme)
    _add_text(
        slide, outline.title or "(对比)",
        left=Inches(0.6), top=Inches(0.4),
        width=Inches(8.8), height=Inches(0.7),
        font_size=28, bold=True,
        color=theme.primary_rgb,
        font_name=theme.font_heading,
    )
    # Two equal columns
    col_w = Inches(4.2)
    col_top = Inches(1.5)
    col_h = Inches(3.8)
    # Left column header
    _add_text(
        slide, outline.left_title or "左",
        left=Inches(0.6), top=col_top,
        width=col_w, height=Inches(0.5),
        font_size=18, bold=True,
        color=theme.accent_rgb,
        font_name=theme.font_heading,
    )
    _add_bullet_text(
        slide, outline.left,
        left=Inches(0.6), top=Inches(2.0),
        width=col_w, height=col_h,
        font_size=16,
        color=theme.text_rgb,
        font_name=theme.font_body,
        accent_color=theme.accent_rgb,
    )
    # Right column header
    _add_text(
        slide, outline.right_title or "右",
        left=Inches(5.2), top=col_top,
        width=col_w, height=Inches(0.5),
        font_size=18, bold=True,
        color=theme.accent_rgb,
        font_name=theme.font_heading,
    )
    _add_bullet_text(
        slide, outline.right,
        left=Inches(5.2), top=Inches(2.0),
        width=col_w, height=col_h,
        font_size=16,
        color=theme.text_rgb,
        font_name=theme.font_body,
        accent_color=theme.accent_rgb,
    )


def _render_image(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    _add_accent_bar(slide, theme)
    _add_text(
        slide, outline.title or "",
        left=Inches(0.6), top=Inches(0.4),
        width=Inches(8.8), height=Inches(0.7),
        font_size=24, bold=True,
        color=theme.primary_rgb,
        font_name=theme.font_heading,
    )
    img = outline.image_path
    if img and Path(img).is_file():
        try:
            slide.shapes.add_picture(
                img,
                Inches(2.0), Inches(1.3),
                width=Inches(6.0),
            )
            if outline.caption:
                _add_text(
                    slide, outline.caption,
                    left=Inches(0.6), top=Inches(4.8),
                    width=Inches(8.8), height=Inches(0.4),
                    font_size=12,
                    color=theme.muted_text_rgb,
                    font_name=theme.font_body,
                    align="center",
                )
            return
        except Exception as exc:  # noqa: BLE001
            log.debug("image insert failed (%s); falling back to bullet", exc)
    # Fallback to bullet layout with caption + warning marker
    fallback_bullets = (
        [f"[image missing: {img}]"] if img else []
    ) + ([outline.caption] if outline.caption else [])
    _add_bullet_text(
        slide, fallback_bullets or ["(图片缺失)"],
        left=Inches(0.7), top=Inches(1.6),
        width=Inches(8.6), height=Inches(3.5),
        font_size=18,
        color=theme.muted_text_rgb,
        font_name=theme.font_body,
        accent_color=theme.accent_rgb,
    )


def _render_quote(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    quote_text = outline.quote or outline.title or "(无引述)"
    # Big left quote mark
    _add_text(
        slide, "“",
        left=Inches(0.6), top=Inches(0.7),
        width=Inches(1.2), height=Inches(1.5),
        font_size=110, bold=True,
        color=theme.accent_rgb,
        font_name=theme.font_heading,
    )
    _add_text(
        slide, quote_text,
        left=Inches(1.4), top=Inches(1.4),
        width=Inches(7.8), height=Inches(2.8),
        font_size=24,
        color=theme.text_rgb,
        font_name=theme.font_body,
        anchor="middle",
    )
    if outline.cite:
        _add_text(
            slide, f"— {outline.cite}",
            left=Inches(1.4), top=Inches(4.4),
            width=Inches(7.8), height=Inches(0.5),
            font_size=14,
            color=theme.muted_text_rgb,
            font_name=theme.font_body,
        )


def _render_toc(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    _add_accent_bar(slide, theme)
    _add_text(
        slide, outline.title or "目录",
        left=Inches(0.6), top=Inches(0.4),
        width=Inches(8.8), height=Inches(0.7),
        font_size=32, bold=True,
        color=theme.primary_rgb,
        font_name=theme.font_heading,
    )
    # Numbered list rendered manually for visual weight
    tb = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(8.6), Inches(3.8),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(outline.bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(10)
        num = para.add_run()
        num.text = f"{i + 1:02d}.  "
        num.font.size = Pt(22)
        num.font.bold = True
        num.font.name = theme.font_heading
        num.font.color.rgb = _rgb(theme.accent_rgb)
        body = para.add_run()
        body.text = str(item)
        body.font.size = Pt(20)
        body.font.name = theme.font_body
        body.font.color.rgb = _rgb(theme.text_rgb)


def _render_chart(slide, outline: SlideOutline, theme: Theme) -> None:
    """Native PowerPoint chart slide (bar / line / pie).

    ``outline.chart`` shape::

        {"type": "bar"|"line"|"pie",
         "categories": ["Q1", "Q2", ...],
         "series": [{"name": "营收", "values": [10, 20, ...]}, ...]}

    Missing / malformed data degrades to a bullet layout so the deck
    never aborts on one bad slide.
    """
    _fill_slide_bg(slide, theme)
    _add_accent_bar(slide, theme)
    _add_text(
        slide, outline.title or "(图表)",
        left=Inches(0.6), top=Inches(0.4),
        width=Inches(8.8), height=Inches(0.7),
        font_size=28, bold=True,
        color=theme.primary_rgb,
        font_name=theme.font_heading,
    )

    spec = outline.chart or {}
    categories = spec.get("categories") or []
    series = [s for s in (spec.get("series") or []) if isinstance(s, dict)]
    if not categories or not series:
        _add_bullet_text(
            slide, outline.bullets or ["（图表数据缺失，已降级为要点）"],
            left=Inches(0.7), top=Inches(1.6),
            width=Inches(8.6), height=Inches(3.8),
            font_size=18, color=theme.text_rgb,
            font_name=theme.font_body, accent_color=theme.accent_rgb,
        )
        return

    ctype = str(spec.get("type") or "bar").lower()
    xl_type = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in categories]
    for s in series:
        name = str(s.get("name") or "系列")
        raw_values = s.get("values") or []
        values: list[float] = []
        for v in raw_values:
            try:
                values.append(float(v))
            except (TypeError, ValueError):
                values.append(0.0)
        chart_data.add_series(name, tuple(values))

    try:
        gframe = slide.shapes.add_chart(
            xl_type,
            Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.5),
            chart_data,
        )
        chart = gframe.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    except Exception as exc:  # noqa: BLE001 — degrade rather than abort
        log.warning("chart render failed, degrading to bullets: %s", exc)
        _add_bullet_text(
            slide, outline.bullets or [f"{s.get('name')}" for s in series],
            left=Inches(0.7), top=Inches(1.6),
            width=Inches(8.6), height=Inches(3.8),
            font_size=18, color=theme.text_rgb,
            font_name=theme.font_body, accent_color=theme.accent_rgb,
        )


def _render_title_v2(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme, theme.dark_bg_rgb)
    _add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(6.95), top=Inches(0.58),
        width=Inches(2.25), height=Inches(1.28),
        fill=theme.secondary_rgb,
    )
    _add_shape(
        slide, MSO_SHAPE.OVAL,
        left=Inches(7.82), top=Inches(3.62),
        width=Inches(1.15), height=Inches(1.15),
        fill=theme.accent_rgb,
    )
    _add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(0.66), top=Inches(0.72),
        width=Inches(1.20), height=Inches(0.34),
        fill=theme.highlight_rgb,
    )
    _add_text(
        slide, outline.title or "(鏃犳爣棰?",
        left=Inches(0.62), top=Inches(1.32),
        width=Inches(6.8), height=Inches(1.38),
        font_size=44, bold=True,
        color=theme.dark_text_rgb,
        font_name=theme.font_heading,
        align="left", anchor="top",
    )
    if outline.subtitle:
        _add_text(
            slide, outline.subtitle,
            left=Inches(0.66), top=Inches(2.92),
            width=Inches(6.6), height=Inches(0.56),
            font_size=19,
            color=theme.dark_muted_text_rgb,
            font_name=theme.font_body,
            align="left",
        )
    _add_text(
        slide, "DeskPet",
        left=Inches(0.66), top=Inches(4.82),
        width=Inches(1.6), height=Inches(0.28),
        font_size=11, bold=True,
        color=theme.dark_muted_text_rgb,
        font_name="Calibri",
    )


def _render_section_v2(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    _add_corner_motif(slide, theme)
    _add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(0.80), top=Inches(1.28),
        width=Inches(8.40), height=Inches(3.10),
        fill=theme.surface_rgb,
        line=(226, 232, 240),
    )
    _add_icon_badge(
        slide, "S",
        left=Inches(1.22), top=Inches(1.70),
        size=Inches(0.62),
        theme=theme,
        fill=theme.accent_rgb,
        text_color=theme.dark_text_rgb,
    )
    _add_text(
        slide, outline.title or "(绔犺妭)",
        left=Inches(1.16), top=Inches(2.34),
        width=Inches(7.75), height=Inches(0.88),
        font_size=38, bold=True,
        color=theme.primary_rgb,
        font_name=theme.font_heading,
        align="center", anchor="middle",
    )
    if outline.subtitle:
        _add_text(
            slide, outline.subtitle,
            left=Inches(1.30), top=Inches(3.22),
            width=Inches(7.40), height=Inches(0.42),
            font_size=15,
            color=theme.muted_text_rgb,
            font_name=theme.font_body,
            align="center",
        )


def _render_bullet_v2(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    _add_title_block(slide, outline.title or "(鏃犳爣棰?", outline.subtitle, theme)
    _add_corner_motif(slide, theme)
    callouts, rest = _extract_callouts(outline.bullets)
    if callouts:
        card_h = Inches(1.46) if len(callouts) > 1 else Inches(1.78)
        for idx, (value, label) in enumerate(callouts):
            _add_callout_card(
                slide, value, label,
                left=Inches(0.68), top=Inches(1.55 + idx * 1.66),
                width=Inches(3.12), height=card_h,
                theme=theme,
            )
        items = rest or [b for b in outline.bullets if b]
        for idx, item in enumerate(items[:4]):
            _add_item_card(
                slide, f"要点 {idx + 1}", item,
                left=Inches(4.10), top=Inches(1.50 + idx * 0.92),
                width=Inches(5.18), height=Inches(0.78),
                theme=theme, badge=str(idx + 1),
            )
        return

    items = list(outline.bullets or [])
    if not items:
        items = [outline.subtitle or ""]
    rows = 2 if len(items) <= 4 else 3
    card_w = Inches(4.16)
    card_h = Inches(1.14 if rows == 3 else 1.48)
    start_top = Inches(1.46)
    for idx, item in enumerate(items[:6]):
        row, col = divmod(idx, 2)
        _add_item_card(
            slide, f"要点 {idx + 1}", item,
            left=Inches(0.70 + col * 4.62),
            top=start_top + row * (card_h + Inches(0.22)),
            width=card_w, height=card_h,
            theme=theme, badge=str(idx + 1),
        )


def _render_two_column_v2(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    _add_title_block(slide, outline.title or "(瀵规瘮)", outline.subtitle, theme)
    _add_corner_motif(slide, theme)
    for col, (x, heading, items, badge) in enumerate((
        (Inches(0.70), outline.left_title or "方案 A", outline.left, "A"),
        (Inches(5.18), outline.right_title or "方案 B", outline.right, "B"),
    )):
        fill = theme.primary_rgb if col == 0 else theme.surface_rgb
        text_color = theme.dark_text_rgb if col == 0 else theme.text_rgb
        body_color = theme.dark_muted_text_rgb if col == 0 else theme.text_rgb
        _add_shape(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            left=x, top=Inches(1.42),
            width=Inches(4.12), height=Inches(3.65),
            fill=fill,
            line=None if col == 0 else (226, 232, 240),
        )
        _add_icon_badge(
            slide, badge,
            left=x + Inches(0.30), top=Inches(1.74),
            size=Inches(0.50),
            theme=theme,
            fill=theme.accent_rgb if col == 0 else theme.secondary_rgb,
            text_color=theme.dark_text_rgb,
        )
        _add_text(
            slide, heading,
            left=x + Inches(0.92), top=Inches(1.74),
            width=Inches(2.85), height=Inches(0.46),
            font_size=21, bold=True,
            color=text_color,
            font_name=theme.font_heading,
        )
        body = "\n".join(f"- {item}" for item in items[:5])
        _add_text(
            slide, body,
            left=x + Inches(0.36), top=Inches(2.52),
            width=Inches(3.42), height=Inches(2.08),
            font_size=15,
            color=body_color,
            font_name=theme.font_body,
        )


def _render_image_v2(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    _add_title_block(slide, outline.title or "", outline.subtitle, theme)
    _add_corner_motif(slide, theme)
    img = outline.image_path
    frame_left, frame_top = Inches(0.72), Inches(1.42)
    frame_w, frame_h = Inches(5.64), Inches(3.46)
    _add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left=frame_left, top=frame_top, width=frame_w, height=frame_h,
        fill=theme.surface_rgb, line=(226, 232, 240),
    )
    if img and Path(img).is_file():
        try:
            slide.shapes.add_picture(
                img,
                frame_left + Inches(0.14), frame_top + Inches(0.14),
                width=frame_w - Inches(0.28),
            )
        except Exception as exc:  # noqa: BLE001
            log.debug("image insert failed (%s); falling back to placeholder", exc)
            img = None
    if not img or not Path(str(img)).is_file():
        _add_icon_badge(
            slide, "IMG",
            left=frame_left + Inches(2.44), top=frame_top + Inches(1.20),
            size=Inches(0.72), theme=theme, fill=theme.secondary_rgb,
        )
        _add_text(
            slide, f"image missing: {img}" if img else "image placeholder",
            left=frame_left + Inches(0.46), top=frame_top + Inches(2.05),
            width=frame_w - Inches(0.92), height=Inches(0.36),
            font_size=12,
            color=theme.muted_text_rgb,
            font_name=theme.font_body,
            align="center",
        )
    _add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(6.62), top=Inches(1.42),
        width=Inches(2.72), height=Inches(3.46),
        fill=theme.primary_rgb,
    )
    _add_text(
        slide, outline.caption or "视觉说明",
        left=Inches(6.96), top=Inches(2.06),
        width=Inches(2.06), height=Inches(1.32),
        font_size=17, bold=True,
        color=theme.dark_text_rgb,
        font_name=theme.font_heading,
        align="left", anchor="middle",
    )


# ── 整页生图 版式多样化(调研: 一套 deck 在 6 种版式里轮换才不单调) ──
# 变体: cover(全幅封面/转场) split_left/split_right(图一侧+文字面板,左右
# 交替) top(上图下文) card(全幅背景+半透明深色卡片) quote(大图引用)。
IMAGE_VARIANTS = ("cover", "split_left", "split_right", "top", "card", "quote")


def _place_cover(slide, img: str, *, left: int, top: int, width: int, height: int) -> bool:
    """object-fit: cover —— 图按【槽位比例中心裁切】后铺满槽位,零拉伸变形。

    用户实测痛点: 把 16:9 横图硬塞进竖/方槽 → 水平挤压。正解不是拉伸,
    是按槽位宽高比中心裁掉多余边(python-pptx picture.crop_*),放置框仍是
    槽位尺寸,但可见部分比例 = 槽位比例 → 不变形。PIL 不可用时退化为直接
    铺(可能轻微拉伸,但有兜底)。从不抛异常。
    """
    try:
        pic = slide.shapes.add_picture(img, Emu(left), Emu(top),
                                       width=Emu(width), height=Emu(height))
        try:
            from PIL import Image  # type: ignore

            with Image.open(img) as im:
                iw, ih = im.size
            if iw <= 0 or ih <= 0:
                return True
            img_ar = iw / ih
            slot_ar = width / height if height else img_ar
            if img_ar > slot_ar:
                # 图太宽 → 裁左右
                frac = (1.0 - slot_ar / img_ar) / 2.0
                pic.crop_left = frac
                pic.crop_right = frac
            elif img_ar < slot_ar:
                # 图太高 → 裁上下
                frac = (1.0 - img_ar / slot_ar) / 2.0
                pic.crop_top = frac
                pic.crop_bottom = frac
        except Exception:  # noqa: BLE001
            pass  # PIL 不可用 → 保持铺满(退化)
        return True
    except Exception as exc:  # noqa: BLE001
        log.debug("place cover failed: %s", exc)
        return False


def _fullbleed_picture(slide, img: str) -> bool:
    return _place_cover(slide, img, left=0, top=0,
                        width=_SLIDE_WIDTH, height=_SLIDE_HEIGHT)


def _side_picture(slide, img: str, *, left_emu: int, width_emu: int) -> bool:
    return _place_cover(slide, img, left=left_emu, top=0,
                        width=width_emu, height=_SLIDE_HEIGHT)


def _panel_text(slide, outline: SlideOutline, theme: Theme, *,
                px: float, pw: float, title_pt: int = 26) -> None:
    """深色面板内: accent 竖条 + 标题 + 要点/caption。px/pw 为英寸。
    字号 × outline.font_scale(视觉评审 shrink_text 动作)。"""
    fs = outline.font_scale or 1.0
    _add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        left=Inches(px), top=Inches(0.7),
        width=Inches(0.06), height=Inches(0.62),
        fill=theme.accent_rgb, line=None,
    )
    _add_text(
        slide, outline.title or "",
        left=Inches(px + 0.2), top=Inches(0.62),
        width=Inches(pw - 0.4), height=Inches(0.92),
        font_size=max(14, int(title_pt * fs)), bold=True,
        color=theme.dark_text_rgb, font_name=theme.font_heading,
        align="left", anchor="middle",
    )
    bullets = [b for b in (outline.bullets or []) if str(b).strip()]
    if bullets:
        _add_bullet_text(
            slide, bullets,
            left=Inches(px + 0.2), top=Inches(1.78),
            width=Inches(pw - 0.5), height=Inches(3.2),
            font_size=max(10, int(15 * fs)), color=theme.dark_text_rgb,
            font_name=theme.font_body, accent_color=theme.accent_rgb,
        )
    elif outline.caption:
        _add_text(
            slide, outline.caption,
            left=Inches(px + 0.2), top=Inches(1.78),
            width=Inches(pw - 0.5), height=Inches(2.6),
            font_size=max(11, int(16 * fs)), color=theme.dark_muted_text_rgb,
            font_name=theme.font_body, align="left",
        )


def _render_var_cover(slide, outline: SlideOutline, theme: Theme) -> None:
    """A 全幅封面/转场: 全幅图 + 底部暗带(标题 + caption)。"""
    band_top = Emu(int(_SLIDE_HEIGHT * 0.76))
    _add_shape(slide, MSO_SHAPE.RECTANGLE, left=Emu(0), top=band_top,
               width=Emu(_SLIDE_WIDTH), height=Emu(int(_SLIDE_HEIGHT * 0.24)),
               fill=theme.dark_bg_rgb, line=None)
    _add_text(slide, outline.title or "", left=Inches(0.62),
              top=band_top + Inches(0.2), width=Inches(8.76), height=Inches(0.6),
              font_size=30, bold=True, color=theme.dark_text_rgb,
              font_name=theme.font_heading, align="left")
    sub = outline.caption or outline.subtitle
    if sub:
        _add_text(slide, sub, left=Inches(0.64), top=band_top + Inches(0.92),
                  width=Inches(8.72), height=Inches(0.36), font_size=14,
                  color=theme.dark_muted_text_rgb, font_name=theme.font_body,
                  align="left")


def _render_var_split(slide, outline: SlideOutline, theme: Theme, *,
                      image_right: bool) -> None:
    """B 图一侧 + 文字面板另一侧(左右交替防呆板)。"""
    panel_w = int(_SLIDE_WIDTH * 0.46)
    img = outline.image_path
    if image_right:
        # 面板在左,图在右
        if img and Path(img).is_file():
            _side_picture(slide, img, left_emu=panel_w,
                          width_emu=_SLIDE_WIDTH - panel_w)
        _add_shape(slide, MSO_SHAPE.RECTANGLE, left=Emu(0), top=Emu(0),
                   width=Emu(panel_w), height=Emu(_SLIDE_HEIGHT),
                   fill=theme.dark_bg_rgb, line=None)
        _panel_text(slide, outline, theme, px=0.62, pw=4.0)
    else:
        # 图在左,面板在右
        if img and Path(img).is_file():
            _side_picture(slide, img, left_emu=0,
                          width_emu=_SLIDE_WIDTH - panel_w)
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   left=Emu(_SLIDE_WIDTH - panel_w), top=Emu(0),
                   width=Emu(panel_w), height=Emu(_SLIDE_HEIGHT),
                   fill=theme.dark_bg_rgb, line=None)
        _panel_text(slide, outline, theme, px=6.2, pw=3.6)


def _render_var_top(slide, outline: SlideOutline, theme: Theme) -> None:
    """C 上图下文: 图占上 60%, 下方标题 + 横排要点(浅底)。"""
    img = outline.image_path
    img_h = int(_SLIDE_HEIGHT * 0.6)
    if img and Path(img).is_file():
        _place_cover(slide, img, left=0, top=0, width=_SLIDE_WIDTH, height=img_h)
    _add_shape(slide, MSO_SHAPE.RECTANGLE, left=Emu(0), top=Emu(img_h),
               width=Emu(_SLIDE_WIDTH), height=Emu(_SLIDE_HEIGHT - img_h),
               fill=theme.background_rgb, line=None)
    _add_text(slide, outline.title or "", left=Inches(0.62),
              top=Emu(img_h) + Inches(0.16), width=Inches(8.76), height=Inches(0.5),
              font_size=24, bold=True, color=theme.primary_rgb,
              font_name=theme.font_heading, align="left")
    bullets = [b for b in (outline.bullets or []) if str(b).strip()]
    if bullets:
        _add_bullet_text(slide, bullets, left=Inches(0.66),
                         top=Emu(img_h) + Inches(0.82), width=Inches(8.7),
                         height=Inches(1.4), font_size=14, color=theme.text_rgb,
                         font_name=theme.font_body, accent_color=theme.accent_rgb)


def _render_var_card(slide, outline: SlideOutline, theme: Theme) -> None:
    """D 全幅背景图 + 一个局部半透明深色卡片(标题 + 要点)。"""
    img = outline.image_path
    if img and Path(img).is_file():
        _fullbleed_picture(slide, img)
    # 右下角深色卡片(半透明,留出大片氛围图,底图微透更精致)
    cx, cy, cw, ch = 5.2, 1.0, 4.4, 3.7
    _card = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left=Inches(cx), top=Inches(cy),
                       width=Inches(cw), height=Inches(ch), fill=(10, 16, 28), line=None)
    _set_fill_alpha(_card, 82)
    _panel_text(slide, outline, theme, px=cx + 0.3, pw=cw - 0.5, title_pt=23)


def _render_var_quote(slide, outline: SlideOutline, theme: Theme) -> None:
    """F 大图引用: 全幅(暗)图 + 居中大字短句。"""
    img = outline.image_path
    if img and Path(img).is_file():
        _fullbleed_picture(slide, img)
    # 半透明压暗整页(alpha 55%)提升可读,但底图仍透出来(不死黑)
    _ov = _add_shape(slide, MSO_SHAPE.RECTANGLE, left=Emu(0), top=Emu(0),
                     width=Emu(_SLIDE_WIDTH), height=Emu(_SLIDE_HEIGHT),
                     fill=(8, 12, 22), line=None)
    _set_fill_alpha(_ov, 58)
    text = outline.quote or outline.title or ""
    _add_text(slide, f"“{text}”" if text else "", left=Inches(1.2),
              top=Inches(2.0), width=Inches(7.6), height=Inches(1.6),
              font_size=34, bold=True, color=theme.dark_text_rgb,
              font_name=theme.font_heading, align="center", anchor="middle")
    sub = outline.cite or outline.caption or outline.subtitle
    if sub:
        _add_text(slide, sub, left=Inches(1.2), top=Inches(3.7), width=Inches(7.6),
                  height=Inches(0.5), font_size=15, color=theme.dark_muted_text_rgb,
                  font_name=theme.font_body, align="center")


def _resolve_image_variant(outline: SlideOutline) -> str:
    """决定本页用哪个版式: 显式 image_variant 优先;否则按内容兜底
    (有要点→split_right; 有 quote→quote; 否则→cover)。"""
    v = (outline.image_variant or "").strip().lower()
    if v in IMAGE_VARIANTS:
        return v
    if outline.quote:
        return "quote"
    if [b for b in (outline.bullets or []) if str(b).strip()]:
        return "split_right"
    return "cover"


def _render_image_full_v2(slide, outline: SlideOutline, theme: Theme) -> None:
    img = outline.image_path
    if img and Path(img).is_file():
        try:
            variant = _resolve_image_variant(outline)
            if variant == "cover":
                _fullbleed_picture(slide, img)
                _render_var_cover(slide, outline, theme)
            elif variant == "split_left":
                _render_var_split(slide, outline, theme, image_right=False)
            elif variant == "split_right":
                _render_var_split(slide, outline, theme, image_right=True)
            elif variant == "top":
                _render_var_top(slide, outline, theme)
            elif variant == "card":
                _render_var_card(slide, outline, theme)
            elif variant == "quote":
                _render_var_quote(slide, outline, theme)
            else:
                _fullbleed_picture(slide, img)
                _render_var_cover(slide, outline, theme)
            return
        except Exception as exc:  # noqa: BLE001
            log.debug("full image insert failed (%s); falling back to placeholder", exc)

    _fill_slide_bg(slide, theme)
    _add_title_block(slide, outline.title or "", outline.subtitle, theme)
    _add_corner_motif(slide, theme)
    _add_icon_badge(
        slide, "IMG",
        left=Inches(4.64), top=Inches(2.22),
        size=Inches(0.72), theme=theme, fill=theme.secondary_rgb,
    )
    _add_text(
        slide, f"image missing: {img}" if img else "image placeholder",
        left=Inches(2.12), top=Inches(3.04),
        width=Inches(5.76), height=Inches(0.42),
        font_size=13,
        color=theme.muted_text_rgb,
        font_name=theme.font_body,
        align="center",
    )
    if outline.caption:
        _add_text(
            slide, outline.caption,
            left=Inches(2.12), top=Inches(3.54),
            width=Inches(5.76), height=Inches(0.38),
            font_size=12,
            color=theme.muted_text_rgb,
            font_name=theme.font_body,
            align="center",
        )


def _render_quote_v2(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    _add_corner_motif(slide, theme)
    _add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(0.78), top=Inches(1.02),
        width=Inches(8.44), height=Inches(3.58),
        fill=theme.surface_rgb,
        line=(226, 232, 240),
    )
    _add_text(
        slide, "\"",
        left=Inches(1.06), top=Inches(0.90),
        width=Inches(0.92), height=Inches(1.10),
        font_size=92, bold=True,
        color=theme.accent_rgb,
        font_name="Calibri",
    )
    _add_text(
        slide, outline.quote or outline.title or "(鏃犲紩杩?",
        left=Inches(1.78), top=Inches(1.64),
        width=Inches(6.96), height=Inches(1.44),
        font_size=24,
        color=theme.text_rgb,
        font_name=theme.font_body,
        anchor="middle",
    )
    if outline.cite:
        _add_text(
            slide, outline.cite,
            left=Inches(1.84), top=Inches(3.50),
            width=Inches(6.7), height=Inches(0.32),
            font_size=12,
            color=theme.muted_text_rgb,
            font_name=theme.font_body,
            align="right",
        )


def _render_toc_v2(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    _add_title_block(slide, outline.title or "鐩綍", outline.subtitle, theme)
    _add_corner_motif(slide, theme)
    items = outline.bullets or []
    rows = 2 if len(items) <= 4 else 3
    card_w = Inches(4.10)
    card_h = Inches(1.18 if rows == 3 else 1.44)
    for idx, item in enumerate(items[:6]):
        row, col = divmod(idx, 2)
        _add_item_card(
            slide, f"{idx + 1:02d}", item,
            left=Inches(0.74 + col * 4.55),
            top=Inches(1.48) + row * (card_h + Inches(0.22)),
            width=card_w, height=card_h,
            theme=theme, badge=f"{idx + 1}",
        )


def _render_chart_v2(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme)
    _add_title_block(slide, outline.title or "(鍥捐〃)", outline.subtitle, theme)
    _add_corner_motif(slide, theme)
    spec = outline.chart or {}
    categories = spec.get("categories") or []
    series = [s for s in (spec.get("series") or []) if isinstance(s, dict)]
    if not categories or not series:
        _render_bullet_v2(slide, outline, theme)
        return

    values: list[float] = []
    for s in series:
        for v in s.get("values") or []:
            try:
                values.append(float(v))
            except (TypeError, ValueError):
                values.append(0.0)
    if values:
        top_value = max(values)
        _add_callout_card(
            slide, f"{top_value:g}", "峰值指标",
            left=Inches(0.72), top=Inches(1.72),
            width=Inches(2.54), height=Inches(1.76),
            theme=theme,
        )
    ctype = str(spec.get("type") or "bar").lower()
    xl_type = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)
    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in categories]
    for s in series:
        raw_values = s.get("values") or []
        clean_values: list[float] = []
        for v in raw_values:
            try:
                clean_values.append(float(v))
            except (TypeError, ValueError):
                clean_values.append(0.0)
        chart_data.add_series(str(s.get("name") or "绯诲垪"), tuple(clean_values))
    try:
        gframe = slide.shapes.add_chart(
            xl_type,
            Inches(3.55), Inches(1.52), Inches(5.70), Inches(3.28),
            chart_data,
        )
        chart = gframe.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    except Exception as exc:  # noqa: BLE001
        log.warning("chart render failed, degrading to cards: %s", exc)
        _render_bullet_v2(slide, outline, theme)


def _render_conclusion_v2(slide, outline: SlideOutline, theme: Theme) -> None:
    _fill_slide_bg(slide, theme, theme.dark_bg_rgb)
    _add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(7.18), top=Inches(0.62),
        width=Inches(1.74), height=Inches(1.04),
        fill=theme.secondary_rgb,
    )
    _add_shape(
        slide, MSO_SHAPE.OVAL,
        left=Inches(0.82), top=Inches(4.12),
        width=Inches(0.64), height=Inches(0.64),
        fill=theme.accent_rgb,
    )
    _add_text(
        slide, outline.title or "Conclusion",
        left=Inches(0.86), top=Inches(1.30),
        width=Inches(7.70), height=Inches(1.05),
        font_size=42, bold=True,
        color=theme.dark_text_rgb,
        font_name=theme.font_heading,
        align="center", anchor="middle",
    )
    if outline.subtitle:
        _add_text(
            slide, outline.subtitle,
            left=Inches(1.34), top=Inches(2.44),
            width=Inches(6.72), height=Inches(0.46),
            font_size=17,
            color=theme.dark_muted_text_rgb,
            font_name=theme.font_body,
            align="center",
        )
    bullets = outline.bullets or outline.left + outline.right
    for idx, item in enumerate(bullets[:3]):
        _add_text(
            slide, item,
            left=Inches(1.45 + idx * 2.45), top=Inches(3.42),
            width=Inches(2.08), height=Inches(0.58),
            font_size=13,
            color=theme.dark_muted_text_rgb,
            font_name=theme.font_body,
            align="center",
        )


_RENDERERS = {
    "title": _render_title_v2,
    "section": _render_section_v2,
    "bullet": _render_bullet_v2,
    "two_column": _render_two_column_v2,
    "image": _render_image_v2,
    "image_full": _render_image_full_v2,
    "quote": _render_quote_v2,
    "toc": _render_toc_v2,
    "chart": _render_chart_v2,
}


def _is_conclusion_slide(slide: SlideOutline, index: int, total: int) -> bool:
    if index != total or slide.layout == "title":
        return False
    title = (slide.title or "").lower()
    keywords = ("结论", "总结", "收尾", "conclusion", "summary", "wrap-up", "wrap up")
    return slide.layout == "section" or any(k in title for k in keywords)


def _add_footer(
    slide, theme: Theme, page_number: int, total: int, *, dark: bool = False,
) -> None:
    """Small page marker, without decorative bars."""
    color = theme.dark_muted_text_rgb if dark else theme.muted_text_rgb
    _add_text(
        slide, f"{page_number} / {total}",
        left=Inches(8.4), top=Inches(5.25),
        width=Inches(1.0), height=Inches(0.3),
        font_size=10,
        color=color,
        font_name=theme.font_body,
        align="right",
    )


def _crop_image_to_169(path: str) -> str:
    """把生成图中心裁切成 16:9 后另存,返回裁切版路径。

    image_full 全幅铺图会把图拉到 16:9 画布 —— 方图/3:2 直接拉伸会明显
    变形。中心裁切零变形。PIL 不可用/失败 → 返回原路径(退化为拉伸)。
    """
    try:
        from PIL import Image  # type: ignore

        src = Path(path)
        with Image.open(src) as im:
            w, h = im.size
            target_h = int(w * 9 / 16)
            if target_h > h:
                # 图太「方」,反向裁宽
                target_w = int(h * 16 / 9)
                left = max(0, (w - target_w) // 2)
                box = (left, 0, left + target_w, h)
            else:
                top = max(0, (h - target_h) // 2)
                box = (0, top, w, top + target_h)
            out = src.with_name(src.stem + "_169" + src.suffix)
            im.crop(box).save(out)
        return str(out)
    except Exception as exc:  # noqa: BLE001
        log.debug("16:9 crop failed (%s); using original image", exc)
        return path


# 每个版式对应的「负空间」构图指令(调研: prompt 主动留白给文字才不盖字)。
_VARIANT_NEG_SPACE = {
    "cover": "Composition: keep the bottom third darker, simpler and uncluttered as clean negative space for a title overlay.",
    "split_right": "Composition: place the main subject on the RIGHT two-thirds; keep the LEFT third clean, simple and darker as negative space for text.",
    "split_left": "Composition: place the main subject on the LEFT two-thirds; keep the RIGHT third clean, simple and darker as negative space for text.",
    "top": "Composition: the main scene fills the upper two-thirds; keep the lower third simple and uncluttered for a caption strip.",
    "card": "Composition: atmospheric wide scene; keep the right side calmer and slightly darker so a text card can sit there.",
    "quote": "Composition: dark, moody, cinematic wide scene with a calm low-contrast center for a short centered quote; subtle vignette.",
}
_NO_TEXT_SUFFIX = " No text, no letters, no words, no typography, no watermark, no logo."


def _assign_image_layouts(slides: list[SlideOutline]) -> None:
    """给整页生图(image_full)的页【自动分配版式变体】,让一套 deck 在多种
    版式里轮换(调研结论: 不是一种打天下)。规则:
      - 封面(第一张 image_full) → cover
      - 结尾(最后一张,含 quote 或 结论意味) → quote
      - 中间内容页 → 在 [split_right, top, card, split_left] 里轮换
        (split 左右交替防呆板);无要点的纯氛围页 → cover。
    已显式写了合法 image_variant 的页保留不动(LLM 可手动指定)。从不抛异常。
    """
    try:
        idxs = [i for i, so in enumerate(slides) if so.layout == "image_full"]
        if not idxs:
            return
        rotation = ["split_right", "top", "card", "split_left"]
        rot_i = 0
        for pos, i in enumerate(idxs):
            so = slides[i]
            if (so.image_variant or "").strip().lower() in IMAGE_VARIANTS:
                continue  # LLM 显式指定,尊重
            has_bullets = bool([b for b in (so.bullets or []) if str(b).strip()])
            if pos == 0:
                so.image_variant = "cover"
            elif so.quote or (pos == len(idxs) - 1 and not has_bullets):
                so.image_variant = "quote"
            elif not has_bullets:
                so.image_variant = "cover"
            else:
                so.image_variant = rotation[rot_i % len(rotation)]
                rot_i += 1
    except Exception as exc:  # noqa: BLE001
        log.debug("assign image layouts failed: %s", exc)


def _autofill_image_prompts(slides: list[SlideOutline]) -> None:
    # 先定版式 → 据版式给每张图的 prompt 追加对应负空间指令 + 禁字后缀。
    _assign_image_layouts(slides)

    pending = [
        (idx, so.image_prompt)
        for idx, so in enumerate(slides)
        if so.image_prompt and not so.image_path
    ]
    if not pending:
        return

    try:
        try:
            from .image_tools import generate_images
        except Exception as exc:  # noqa: BLE001
            log.warning("image prompt autofill unavailable: %s", exc)
            return

        # 按版式选生成尺寸(最佳实践: 生成贴近槽位比例的图,cover 裁切损失
        # 最小)。split 侧栏≈方 → 1024x1024;其余(全幅/上图/封面)→ 横版
        # 1536x1024。prompt 追加版式负空间指令 + 禁字后缀。同尺寸分一组批量。
        from collections import OrderedDict
        groups: "OrderedDict[str, list[tuple[int, str]]]" = OrderedDict()
        for idx, prompt in pending:
            variant = _resolve_image_variant(slides[idx])
            neg = _VARIANT_NEG_SPACE.get(variant, "")
            full = f"{prompt} {neg}{_NO_TEXT_SUFFIX}".strip()
            size = "1024x1024" if variant in ("split_left", "split_right") else "1536x1024"
            groups.setdefault(size, []).append((idx, full))
        for size, items in groups.items():
            results = generate_images([p for _, p in items], size=size)
            for (idx, _p), result in zip(items, results):
                path = result.get("path") if isinstance(result, dict) else None
                if path:
                    slides[idx].image_path = str(path)
                else:
                    err = result.get("error") if isinstance(result, dict) else None
                    log.debug("image gen failed idx=%d err=%r", idx, err)
    except Exception as exc:  # noqa: BLE001
        log.warning("image prompt autofill failed: %s", exc, exc_info=True)


def _default_template() -> Optional[str]:
    """无显式 template 时的默认模板(大类名或 .pptx 路径)，来自环境变量
    ``DESKPET_PPT_DEFAULT_TEMPLATE``。不设 → None(走 from-scratch)。
    可设成一个大类名(如「高级色」)让生成默认套模板并按预览图选具体设计。
    """
    try:
        return (os.environ.get("DESKPET_PPT_DEFAULT_TEMPLATE", "") or "").strip() or None
    except Exception:  # noqa: BLE001
        return None


def _degrade_to_template(slides: Sequence[SlideOutline]) -> list[SlideOutline]:
    """Return template-friendly slide copies and clear all generated-image fields."""
    degraded: list[SlideOutline] = []
    for slide in slides:
        so = copy.deepcopy(slide).normalize()
        if so.layout == "image_full" and so.image_prompt:
            so.layout = "bullet" if so.bullets else "section"
        so.image_prompt = None
        so.image_path = None
        degraded.append(so)
    return degraded


def _should_fallback(res: dict) -> bool:
    return res.get("error_kind") in {"connectivity", "model_unavailable"}


def _image_prompt_payload(slide: SlideOutline) -> tuple[str, str]:
    variant = _resolve_image_variant(slide)
    neg = _VARIANT_NEG_SPACE.get(variant, "")
    prompt = f"{slide.image_prompt or ''} {neg}{_NO_TEXT_SUFFIX}".strip()
    size = "1024x1024" if variant in ("split_left", "split_right") else "1536x1024"
    return prompt, size


def _autofill_with_connectivity_gate(slides: list[SlideOutline]) -> tuple[bool, int]:
    """Generate image prompts with a first-image connectivity gate."""
    _assign_image_layouts(slides)
    pending = [
        (idx, so)
        for idx, so in enumerate(slides)
        if so.image_prompt and not so.image_path
    ]
    if not pending:
        return True, 0

    i0, s0 = pending[0]
    prompt0, size0 = _image_prompt_payload(s0)
    try:
        first = generate_images([prompt0], size=size0)
    except Exception as exc:  # noqa: BLE001
        log.warning("ppt_pro first image generation failed: %s", exc)
        return False, 0
    res0 = first[0] if first else {}
    n_ok = 0
    path0 = res0.get("path") if isinstance(res0, dict) else None
    if path0:
        slides[i0].image_path = str(path0)
        n_ok += 1
    elif isinstance(res0, dict) and _should_fallback(res0):
        return False, 0

    for idx, slide in pending[1:]:
        prompt, size = _image_prompt_payload(slide)
        try:
            results = generate_images([prompt], size=size)
        except Exception as exc:  # noqa: BLE001
            log.debug("ppt_pro image generation failed idx=%d err=%s", idx, exc)
            continue
        res = results[0] if results else {}
        path = res.get("path") if isinstance(res, dict) else None
        if path:
            slide.image_path = str(path)
            n_ok += 1
    if n_ok == 0:
        return False, 0
    return True, n_ok


def _drive_free_bytes(path: Path) -> Optional[int]:
    """Free bytes on the volume containing ``path``, walking up to the first
    existing parent (the dir may not exist yet on a fresh install). Returns
    None if it can't be determined. Never raises."""
    import shutil

    p = path
    for _ in range(8):
        try:
            return shutil.disk_usage(str(p)).free
        except Exception:  # noqa: BLE001
            if p.parent == p:
                return None
            p = p.parent
    return None


def _disk_preflight(*, image_mode: bool, pages: int) -> Optional[str]:
    """Return a user-facing error message if the user-data volume is too low
    on free space to safely generate images + write the ``.pptx``, else None.

    Conservative estimate so we fail *early* with a clear ask instead of
    crashing mid-render with a ``disk I/O error`` (observed 2026-06-24 when
    C: hit 0 bytes). Never raises — on any uncertainty it returns None so a
    space check can never itself block a valid render."""
    try:
        from paths import user_data_dir  # type: ignore[import-not-found]

        base = user_data_dir()
    except Exception:  # noqa: BLE001
        return None
    free = _drive_free_bytes(base)
    if free is None:
        return None
    per_page_mb = 8 if image_mode else 1  # seedream PNG ~1-3MB/page + headroom
    need_mb = 60 + max(1, int(pages)) * per_page_mb
    free_mb = int(free / (1024 * 1024))
    if free_mb >= need_mb:
        return None
    try:
        drive = base.anchor or str(base)
    except Exception:  # noqa: BLE001
        drive = str(base)
    return (
        f"磁盘空间不足：{drive} 仅剩 {free_mb}MB，生成这份 PPT 约需 "
        f"{need_mb}MB，请清理磁盘后重试。"
    )


def _render_pro(
    slides: list[SlideOutline],
    *,
    theme: str,
    title: str,
    author: str,
    output_path: Optional[str],
    image_mode: bool,
    probe_timeout_s: float,
    notify,
) -> dict[str, Any]:
    use_template = not image_mode
    working = [copy.deepcopy(s).normalize() for s in slides]
    log.info("ppt_pro _render_pro start image_mode=%s pages=%d", image_mode, len(working))

    disk_err = _disk_preflight(image_mode=image_mode, pages=len(working))
    if disk_err:
        log.warning("ppt_pro disk preflight blocked render: %s", disk_err)
        return {"ok": False, "error": disk_err}

    if image_mode and not probe_image_reachable(timeout_s=probe_timeout_s):
        use_template = True
        notify("AI 配图暂时连不上，已切换模板生成。")
    if image_mode and not use_template:
        reachable, _n_ok = _autofill_with_connectivity_gate(working)
        log.info("ppt_pro gate reachable=%s n_ok=%s", reachable, _n_ok)
        if not reachable:
            use_template = True
            notify("AI 配图暂时用不了，已切换模板生成。")
        else:
            # Partial-failure visibility: the gate skips images that fail
            # after the first one succeeds (those slides keep their prompt
            # but no image_path). Tell the user instead of silently shipping
            # a deck with blank image slots.
            n_want = sum(1 for s in working if s.image_prompt)
            n_got = sum(1 for s in working if s.image_prompt and s.image_path)
            if 0 < n_got < n_want:
                notify(
                    f"有 {n_want - n_got}/{n_want} 张配图没生成成功，"
                    f"这些页改用纯色版式（其余 {n_got} 张已配图）。"
                )

    if use_template:
        # bug#2 修：回退**不走大类名 vision 选图**（外部 2.8GB 库 PIL 拼 contact-sheet
        # 在 executor 线程会阻塞 → deck 永不落盘）。改用确定性 bundled 模板**直传路径**
        # （_resolve_template_path 走路径分支，跳过 pick_template_by_preview）。
        tpl = _fallback_template_path()
        log.info("ppt_pro render template path=%s tpl=%r", "template", tpl)
        # bug#2-b 修：ppt_create 的 parse_outline 只认 JSON/dict 列表，不认 SlideOutline
        # 实例 → 必须 asdict 转 dict 再传，否则 "outline parse failed or empty"。
        result = ppt_create(
            [asdict(s) for s in _degrade_to_template(working)],
            theme=theme,
            title=title,
            author=author,
            output_path=output_path,
            template=tpl,            # 直传 .pptx 路径或 None(→fromscratch)，绝不传大类名
            skip_image_gen=True,
        )
        log.info("ppt_pro render done(template) ok=%s path=%s", result.get("ok"), result.get("path"))
        return result
    log.info("ppt_pro render path=fromscratch(惊艳)")
    result = ppt_create(
        [asdict(s) for s in working],
        theme=theme,
        title=title,
        author=author,
        output_path=output_path,
        skip_image_gen=True,
    )
    log.info("ppt_pro render done(fromscratch) ok=%s path=%s", result.get("ok"), result.get("path"))
    return result


def _fallback_template_path() -> Optional[str]:
    """回退用**确定性 bundled 模板直传路径**（跳过外部大库 + vision 选图，避免
    pick_template_by_preview 在 executor 线程做 90 张大预览图 PIL 拼图阻塞）。
    返回 bundled 通用商务 第一套 .pptx 的绝对路径；无则 None（上层 fromscratch）。从不抛。"""
    try:
        bundled = Path(__file__).parent / "ppt_templates" / "通用商务"
        cands = sorted(bundled.glob("*.pptx"))
        return str(cands[0]) if cands else None
    except Exception:  # noqa: BLE001
        return None


def _user_template_roots() -> list[Path]:
    """额外模板根目录(递归搜索),来自 env ``DESKPET_PPT_TEMPLATE_ROOTS``
    (``;`` 分隔)。库主根由 ppt_template_picker.template_library_root() 提供。
    从不抛异常。
    """
    roots: list[Path] = []
    try:
        env = (os.environ.get("DESKPET_PPT_TEMPLATE_ROOTS", "") or "").strip()
        if env:
            roots += [Path(p.strip()) for p in env.split(";") if p.strip()]
    except Exception:  # noqa: BLE001
        pass
    return [r for r in roots if r.is_dir()]


def _resolve_template_path(template: Optional[str]) -> Optional[str]:
    """解析【直传 .pptx 路径】或【库内模板 stem】→ 绝对路径。从不抛异常。

    注意:这里【不】处理大类名 —— 大类名要靠预览图视觉选具体文件,见
    :func:`_resolve_template_for_render`(它有 topic 上下文)。
    """
    try:
        from . import ppt_template_picker as _tpl

        raw = str(template or "").strip()
        if not raw:
            return None

        direct = Path(raw).expanduser()
        if direct.is_file():
            return str(direct.resolve())

        name = Path(raw).name
        if name.lower().endswith(".pptx"):
            name = name[:-5]
        if not name:
            return None

        # 库主根 + env 额外根,按 stem 递归匹配
        roots: list[Path] = []
        lib = _tpl.template_library_root()
        if lib is not None:
            roots.append(lib)
        roots += _user_template_roots()
        for root in roots:
            try:
                hit = next(
                    (p for p in root.rglob(f"{name}.pptx") if p.is_file()), None
                )
                if hit is not None:
                    return str(hit.resolve())
            except Exception:  # noqa: BLE001
                continue
    except Exception:  # noqa: BLE001
        return None
    return None


def _resolve_template_for_render(
    template: Optional[str], *, topic: str = ""
) -> Optional[str]:
    """把 ``template``(大类名 / .pptx 路径 / 库内 stem)解析成具体 .pptx 路径。

    大类名 → 预览图视觉选一套;否则按路径/stem 解析。从不抛异常。
    """
    try:
        from . import ppt_template_picker as _tpl

        raw = str(template or "").strip()
        if not raw:
            return None
        if _tpl.is_category(raw):
            picked = _tpl.pick_template_by_preview(raw, topic)
            return str(picked) if picked is not None else None
        return _resolve_template_path(raw)
    except Exception:  # noqa: BLE001
        return None


def _pick_template_layout(prs, deskpet_layout: str):
    """把 DeskPet 布局名映射到模板的 slide layout。

    先按关键词(中英别名)匹配 layout 名 —— 国内下载的「高级感」模板布局名
    多为中文('标题幻灯片'/'标题和内容'/'节标题'/'两栏内容'/'图片与标题'),
    纯英文关键词匹配不到会回退标准 idx,遇到布局重排的模板就选错。补中文
    别名后绝大多数中文模板能正确映射。匹配不到再回退标准 PowerPoint idx。
    """
    mapping = {
        "title": (["title slide", "标题幻灯片", "封面", "首页"], 0),
        "section": (["section header", "section", "节标题", "章节", "过渡"], 2),
        "bullet": (["title and content", "标题和内容", "内容与标题", "正文"], 1),
        "two_column": (["two content", "两栏内容", "两栏", "双栏", "comparison", "比较"], 3),
        "image": (["picture with caption", "picture", "图片与标题", "图文", "图片"], 8),
        "image_full": (["picture with caption", "picture", "图片与标题", "图文", "图片"], 8),
        "quote": (["section header", "section", "节标题", "引用"], 2),
        "toc": (["title and content", "标题和内容", "目录", "内容与标题"], 1),
        "chart": (["title and content", "标题和内容", "图表", "内容与标题"], 1),
    }
    keywords, fallback_idx = mapping.get(
        deskpet_layout, (["title and content", "标题和内容"], 1)
    )

    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ValueError("template has no slide layouts")
    # 按关键词优先级逐个扫 layout 名(中英大小写不敏感),命中即返回
    for kw in keywords:
        kw_l = kw.lower()
        for layout in layouts:
            if kw_l in (layout.name or "").lower():
                return layout
    # 回退标准 idx(clamp 防越界)
    idx = min(max(fallback_idx, 0), len(layouts) - 1)
    try:
        return layouts[idx]
    except Exception:  # noqa: BLE001
        return layouts[0]


def _template_placeholders(slide, *types: Any) -> list[Any]:
    placeholders = []
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            if shape.placeholder_format.type in types:
                placeholders.append(shape)
        except Exception:  # noqa: BLE001
            continue
    return placeholders


def _set_placeholder_text(ph, text: str) -> None:
    if not text:
        return
    try:
        ph.text_frame.text = text
    except Exception:  # noqa: BLE001
        try:
            ph.text = text
        except Exception:  # noqa: BLE001
            pass


def _fill_placeholder_lines(ph, lines: Sequence[str], *, bold_first: bool = False) -> None:
    clean = [str(line).strip() for line in lines if str(line).strip()]
    if not clean:
        return
    try:
        tf = ph.text_frame
        tf.text = clean[0]
        if bold_first and tf.paragraphs and tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].font.bold = True
        for line in clean[1:]:
            p = tf.add_paragraph()
            p.text = line
    except Exception:  # noqa: BLE001
        pass


def _iter_design_text_shapes(shapes) -> Iterable[Any]:
    for shape in shapes:
        try:
            if getattr(shape, "shape_type", None) == 6:
                yield from _iter_design_text_shapes(shape.shapes)
                continue
            if getattr(shape, "has_text_frame", False) and (shape.text or "").strip():
                yield shape
        except Exception:  # noqa: BLE001
            continue


def _shape_text_max_pt(shape) -> float:
    try:
        max_pt = 0.0
        for paragraph in shape.text_frame.paragraphs:
            try:
                size = paragraph.font.size
                if size is not None:
                    max_pt = max(max_pt, float(size.pt))
            except Exception:  # noqa: BLE001
                pass
            for run in paragraph.runs:
                try:
                    size = run.font.size
                    if size is not None:
                        max_pt = max(max_pt, float(size.pt))
                except Exception:  # noqa: BLE001
                    continue
        return max_pt or 18.0
    except Exception:  # noqa: BLE001
        return 18.0


def _shape_text_len(shape) -> int:
    try:
        return len((shape.text or "").strip())
    except Exception:  # noqa: BLE001
        return 0


def _shape_reading_key(shape) -> tuple[int, int]:
    try:
        return (int(shape.top), int(shape.left))
    except Exception:  # noqa: BLE001
        return (0, 0)


def _shape_mostly_in_canvas(shape, slide_w: int, slide_h: int) -> bool:
    """形状面积是否 ≥70% 落在画布内。出血/越界的艺术装饰字(如页缘的
    QUOTE / 手写体大字)塞中文必被裁切,不能当 title/subtitle 槽。
    取不到几何时保守返回 True(不误杀)。"""
    try:
        left, top = int(shape.left), int(shape.top)
        w, h = int(shape.width), int(shape.height)
        if w <= 0 or h <= 0:
            return True
        vis_w = min(left + w, slide_w) - max(left, 0)
        vis_h = min(top + h, slide_h) - max(top, 0)
        if vis_w <= 0 or vis_h <= 0:
            return False
        return (vis_w * vis_h) >= 0.7 * (w * h)
    except Exception:  # noqa: BLE001
        return True


def _fit_font_pt(shape, lines: list[str], orig_pt: float) -> float:
    """按形状宽高估算能放下 lines 的字号(中文全角宽≈1em)。
    从原字号 0.9 步进往下缩,下限 24pt;估不出时返回原字号。"""
    try:
        width, height = int(shape.width), int(shape.height)
        if width <= 0 or height <= 0 or not lines:
            return orig_pt
        if len(lines) == 1:
            # 单行标题：必不折行(折行孤字最丑,如「核心能/力」)。
            # 可用宽 = 框宽 - 内边距(默认左右各~0.1in,取 0.28in 保险),
            # 中文全角字宽≈1em,再乘 0.92 留呼吸感,解出单行最大字号。
            n = max(len(lines[0]), 1)
            usable_pt = (width - 360000) / 12700.0
            pt_single = usable_pt / n * 0.92
            return max(24.0, min(float(orig_pt), pt_single))
        pt = float(orig_pt)
        while pt > 24.0:
            per_line = max(1, int(width / (pt * 12700 * 0.95)))
            needed = sum(
                max(1, -(-len(line) // per_line))  # ceil
                for line in lines
            )
            max_lines = max(1, int(height / (pt * 12700 * 1.25)))
            if needed <= max_lines:
                return pt
            pt *= 0.9
        return 24.0
    except Exception:  # noqa: BLE001
        return orig_pt


def _ensure_readable_text(shape) -> None:
    """标题若继承了「背景装饰字」的极淡显式 RGB(亮度>0.72) → 加深为深灰,
    保证主标题可读。主题色继承(非显式 RGB)不动。"""
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                color = run.font.color
                if getattr(color, "type", None) is None:
                    continue
                rgb = getattr(color, "rgb", None)
                if rgb is None:
                    continue
                r, g, b = rgb[0], rgb[1], rgb[2]
                luma = (r * 299 + g * 587 + b * 114) / 255000.0
                if luma > 0.72:
                    run.font.color.rgb = RGBColor(0x20, 0x29, 0x33)
    except Exception:  # noqa: BLE001
        pass


def _analyze_design_page(slide, slide_w: int = 0, slide_h: int = 0) -> dict[str, Any]:
    """Detect editable text slots on a designed example slide."""
    try:
        records: list[dict[str, Any]] = []
        for shape in _iter_design_text_shapes(slide.shapes):
            try:
                text_len = _shape_text_len(shape)
                records.append({
                    "shape": shape,
                    "max_pt": _shape_text_max_pt(shape),
                    "text_len": text_len,
                    "key": _shape_reading_key(shape),
                    "in_canvas": (
                        _shape_mostly_in_canvas(shape, slide_w, slide_h)
                        if slide_w and slide_h else True
                    ),
                })
            except Exception:  # noqa: BLE001
                continue

        # 出血/越界形状不当 title/subtitle(塞中文必裁切);它们若是 >60pt
        # 大字会自然落入 decor 被清空。
        title_candidates = [
            item for item in records
            if item["text_len"] <= 60 and item.get("in_canvas", True)
        ]
        title_item = max(title_candidates, key=lambda item: item["max_pt"], default=None)
        title_shape = title_item["shape"] if title_item else None

        remaining = [item for item in records if item.get("shape") is not title_shape]
        # subtitle 必须是小字(<=60pt)：>60pt 的非 title 大字是「双层艺术叠排」
        # 的另一层(如封面背景大字 + 前景手写体)，把 subtitle 填进去会和中文
        # 标题重叠混乱 —— 它们归 decor，填充后清空。
        subtitle_item = max(
            (
                item for item in remaining
                if item["max_pt"] <= 60 and item.get("in_canvas", True)
            ),
            key=lambda item: item["max_pt"],
            default=None,
        )
        subtitle_shape = subtitle_item["shape"] if subtitle_item else None
        # decor = 非 title/subtitle 的艺术大字(>60pt)。纯数字/超短文本
        # (如步骤编号 '1' '2' '3')保留当装饰(decor_nums 单列,若与已填
        # 内容重叠会在填充后被清),其余(如 'About Me' 'graphic designer')
        # 残留英文穿帮 → 进 decor 待清空。
        decor = []
        decor_nums = []
        for item in remaining:
            if item["max_pt"] <= 60 or item["shape"] is subtitle_shape:
                continue
            try:
                _txt = (item["shape"].text_frame.text or "").strip()
            except Exception:  # noqa: BLE001
                _txt = ""
            if len(_txt) <= 2 or _txt.isdigit():
                decor_nums.append(item["shape"])
                continue
            decor.append(item["shape"])
        bodies = sorted(
            (
                item["shape"]
                for item in records
                if item["max_pt"] <= 32 and item["text_len"] >= 25
            ),
            key=_shape_reading_key,
        )
        labels = sorted(
            (
                item["shape"]
                for item in records
                if item["max_pt"] <= 32 and item["text_len"] < 25
            ),
            key=_shape_reading_key,
        )
        title_pt = float(title_item["max_pt"]) if title_item else 0.0
        return {
            "title": title_shape,
            "subtitle": subtitle_shape,
            "bodies": bodies,
            "labels": labels,
            "decor": decor,
            "decor_nums": decor_nums,
            "n_body": len(bodies),
            "title_pt": title_pt,
            "slide_w": slide_w,
            "slide_h": slide_h,
        }
    except Exception:  # noqa: BLE001
        return {
            "title": None,
            "subtitle": None,
            "bodies": [],
            "labels": [],
            "decor": [],
            "decor_nums": [],
            "n_body": 0,
            "title_pt": 0.0,
            "slide_w": slide_w,
            "slide_h": slide_h,
        }


def _set_text_keep_style(shape, lines: list[str], *, fit_pt: bool = False) -> None:
    """Replace text while preserving the first run/paragraph style.

    fit_pt=True(title/subtitle 大字槽用): 原字号是为原英文词定制的
    (如 162pt×'SARAH AMELIA'),中文长内容塞进去会撑爆折行孤字 —— 按
    形状宽高自适应缩字号后再替换。
    """
    try:
        clean = [str(line).strip()[:80] for line in (lines or []) if str(line).strip()]
        tf = shape.text_frame
        if not clean:
            tf.text = ""
            return

        p0 = tf.paragraphs[0] if tf.paragraphs else None
        if p0 is None or not p0.runs:
            tf.text = "\n".join(clean)
            return

        r0 = p0.runs[0]
        if fit_pt:
            try:
                orig_size = r0.font.size
                orig_pt = float(orig_size.pt) if orig_size else 0.0
                if orig_pt > 24.0:
                    new_pt = _fit_font_pt(shape, clean, orig_pt)
                    if new_pt < orig_pt:
                        r0.font.size = Pt(int(round(new_pt)))
            except Exception:  # noqa: BLE001
                pass
        r0.text = clean[0]
        for run in list(p0.runs)[1:]:
            try:
                run._r.getparent().remove(run._r)
            except Exception:  # noqa: BLE001
                pass

        tx_body = tf._txBody
        for paragraph in list(tf.paragraphs)[1:]:
            try:
                tx_body.remove(paragraph._p)
            except Exception:  # noqa: BLE001
                pass

        for line in clean[1:]:
            new_p = copy.deepcopy(p0._p)
            first_text = None
            for node in new_p.iter():
                if node.tag.endswith("}t"):
                    if first_text is None:
                        first_text = node
                        node.text = line
                    else:
                        node.text = ""
            tx_body.append(new_p)
    except Exception:  # noqa: BLE001
        try:
            shape.text_frame.text = "\n".join(str(line).strip()[:80] for line in lines if str(line).strip())
        except Exception:  # noqa: BLE001
            pass


def _shapes_v_overlap(a, b) -> bool:
    """两形状垂直方向是否大面积重叠(>50% 较小高度)。判断失败保守视为重叠。"""
    try:
        top_a, bot_a = int(a.top), int(a.top) + int(a.height)
        top_b, bot_b = int(b.top), int(b.top) + int(b.height)
        inter = min(bot_a, bot_b) - max(top_a, top_b)
        return inter > 0.5 * min(bot_a - top_a, bot_b - top_b)
    except Exception:  # noqa: BLE001
        return True


def _take_subtitle_from_decor(info: dict[str, Any]):
    """subtitle 没有小字槽可放时，从 decor 里挑一个与 title 不重叠的形状
    用来放 subtitle(双层艺术叠排里位置错开的那层)。挑中的从 decor 移除，
    位置重叠的(会和中文标题打架)留在 decor 待清空。"""
    title_shape = info.get("title")
    slide_w = int(info.get("slide_w") or 0)
    slide_h = int(info.get("slide_h") or 0)
    for shape in list(info.get("decor", [])):
        if slide_w and slide_h and not _shape_mostly_in_canvas(shape, slide_w, slide_h):
            continue  # 出血装饰层塞中文必裁切,留 decor 清空
        if title_shape is None or not _shapes_v_overlap(shape, title_shape):
            info["decor"].remove(shape)
            return shape
    return None


def _clear_design_body_slots(info: dict[str, Any], *, keep: Sequence[Any] = ()) -> None:
    keep_ids = {id(shape) for shape in keep if shape is not None}
    for shape in info.get("bodies", []):
        if id(shape) not in keep_ids:
            _set_text_keep_style(shape, [])


def _best_design_content_page(
    pages: Sequence[dict[str, Any]], used: set[int], wanted: int,
) -> Optional[dict[str, Any]]:
    candidates = [page for page in pages if page["index"] not in used]
    if not candidates:
        return None
    return min(
        candidates,
        key=lambda page: (
            page["info"]["n_body"] < wanted,
            abs(page["info"]["n_body"] - wanted),
            -page.get("score", 0.0),  # 槽数相当时挑更规整的页
            page["index"],
        ),
    )


def _first_unused_page(pages: Sequence[dict[str, Any]], used: set[int]) -> Optional[dict[str, Any]]:
    for page in pages:
        if page["index"] not in used:
            return page
    return None


def _page_quality_score(info: dict[str, Any]) -> float:
    """设计页「规整度」打分(越高越好,给选页排序用)。

    真机看图发现: 竖排标题页(标题框窄高,中文被裁成"城乡与区")、图标网格
    页(一堆小 label)、装饰过多页填充易穿帮;而标题横排 + 几个对称正文槽的
    页填充干净。据此打分,选页优先挑高分页。从不抛异常。
    """
    try:
        score = 0.0
        title = info.get("title")
        if title is not None:
            try:
                w, h = int(title.width), int(title.height)
                if w >= h:
                    score += 3.0           # 横排标题 = 好
                else:
                    score -= 5.0           # 竖排/窄高标题 = 中文必裁,重罚
            except Exception:  # noqa: BLE001
                pass
        # 正文槽适量(1~4)最佳;过多 = 图标网格类,易乱
        nb = int(info.get("n_body", 0))
        if 1 <= nb <= 4:
            score += 2.0
        elif nb > 6:
            score -= 2.0
        # 装饰/标签噪声越多越易穿帮
        score -= 0.4 * len(info.get("labels", []))
        score -= 0.6 * len(info.get("decor", []))
        return score
    except Exception:  # noqa: BLE001
        return 0.0


def _select_design_page(
    so: SlideOutline,
    pages: Sequence[dict[str, Any]],
    section_pages: Sequence[dict[str, Any]],
    content_pages: Sequence[dict[str, Any]],
    used: set[int],
) -> Optional[dict[str, Any]]:
    if so.layout == "title" and 0 not in used:
        return pages[0]

    if so.layout in {"section", "quote"}:
        return (
            _first_unused_page(section_pages, used)
            or _first_unused_page(content_pages, used)
            or _first_unused_page(pages, used)
        )

    if so.layout == "two_column":
        two_slot_pages = [
            page for page in content_pages
            if page["info"]["n_body"] >= 2
        ]
        return (
            _first_unused_page(two_slot_pages, used)
            or _best_design_content_page(content_pages, used, 2)
            or _first_unused_page(pages, used)
        )

    if so.layout in {"bullet", "toc", "chart", "image", "image_full"}:
        wanted = len(so.bullets) or (1 if so.caption else 0) or 1
        return (
            _best_design_content_page(content_pages, used, wanted)
            or _first_unused_page(pages, used)
        )

    return _first_unused_page(pages, used)


def _enable_text_autofit(shape) -> None:
    """开启文本框「自动换行 + 收缩字号填充」(WPS/PowerPoint 渲染时把超长
    中文压进框)。模板正文槽是为短英文 lorem 设计的,长中文会溢出/挤成
    一团 —— normAutofit 让渲染器自动缩字号到放得下。从不抛异常。
    """
    try:
        tf = shape.text_frame
        tf.word_wrap = True
        if MSO_AUTO_SIZE is not None:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:  # noqa: BLE001
        pass


def _fit_text_to_shape(shape, *, min_pt: int = 10, default_pt: float = 18.0) -> None:
    """主动按槽位高度缩字号 —— 真机实测 WPS 渲染【不执行】normAutofit
    声明,长中文照样溢出框底(压页边/装饰条)。这里在填充后用估算把
    所有 run 字号显式缩到「估算总高 ≤ 槽高」,确定性根治垂直溢出。

    估算: 每段行数 = ceil(字符宽/可用宽)(中文全角≈1em,保守按全角),
    段高 = 行数 × 字号 × 1.36(行距)。从不抛异常。
    """
    try:
        tf = shape.text_frame
        width = int(shape.width)
        height = int(shape.height)
        if width <= 0 or height <= 0:
            return

        paras = list(tf.paragraphs)
        if not paras:
            return

        def _para_pt(p) -> float:
            for r in p.runs:
                if r.font.size is not None:
                    return float(r.font.size.pt)
            if p.font.size is not None:
                return float(p.font.size.pt)
            return float(default_pt)

        def _est_height(scale: float) -> float:
            total = 0.0
            for p in paras:
                txt = "".join(r.text or "" for r in p.runs) or (p.text or "")
                n = max(len(txt.strip()), 1)
                pt = max(float(min_pt), _para_pt(p) * scale)
                per_line = max(1, int(width / (pt * 12700 * 0.98)))
                lines = -(-n // per_line)  # ceil
                total += lines * pt * 12700 * 1.36
            return total

        usable_h = height * 0.96
        if _est_height(1.0) <= usable_h:
            return  # 放得下,不动

        # 二分/步进找能放下的最大缩放(0.5 下限)
        scale = 1.0
        while scale > 0.5 and _est_height(scale) > usable_h:
            scale -= 0.06
        scale = max(scale, 0.5)
        for p in paras:
            base = _para_pt(p)
            new_pt = max(min_pt, int(base * scale))
            for r in p.runs:
                r.font.size = Pt(new_pt)
            # 段级字号也写(无 run 的空段/继承段)
            try:
                p.font.size = Pt(new_pt)
            except Exception:  # noqa: BLE001
                pass
        log.debug("fit_text_to_shape scale=%.2f", scale)
    except Exception:  # noqa: BLE001
        pass


import re as _re_ppt

# 残留占位文 = 模板自带的示例/提示文字,填充后没被覆盖会穿帮。中英都覆盖。
_PLACEHOLDER_PATTERNS = (
    "输入标题", "输入内容", "请输入", "在此输入", "点击输入", "点击此处",
    "标题内容", "您的内容", "替换文字", "此处添加", "添加标题", "添加文字",
    "输入文本", "输入您的", "输入相关", "请替换", "示例文字", "正文内容",
    "presentations are communication", "click to edit", "lorem ipsum",
    "your text here", "add text", "sample text", "ipsum lorem",
)


def _is_residual_placeholder(text: str) -> bool:
    """文本是否是模板残留占位文(中英)。命中 → 填充后清空,防穿帮。

    判定: 归一化后命中占位短语 OR 是大段英文 lorem(≥6 个连续英文单词且
    无中文 —— 模板正文槽的英文示例段)。真实中文内容不会命中。
    """
    t = (text or "").strip()
    if not t:
        return False
    low = t.lower()
    for pat in _PLACEHOLDER_PATTERNS:
        # 只有当文本【基本就是占位词本身】(命中且全文不比占位词长太多)才
        # 算残留 —— 防误伤含占位词子串的真实内容(如 bullet「替换文字槽…」)。
        if pat in low and len(t) <= len(pat) + 6:
            return True
    # 大段纯英文(无中文)且词数多 → 模板英文 lorem 正文示例(整段替换文)
    if not _re_ppt.search(r"[一-鿿]", t):
        words = _re_ppt.findall(r"[A-Za-z]{2,}", t)
        if len(words) >= 6:
            return True
    return False


def _sweep_residual_placeholders(slide) -> None:
    """填充后兜底清扫: 把本页所有【仍是残留占位文】的文本框清空。

    通用兜底 —— design-fill 只能识别/填充规整槽位,模板里复杂版式
    (图标网格/竖排/装饰副标题)的占位文常漏清。真实内容刚被填进去不会
    命中占位模式,所以这一扫只杀穿帮、不误伤。从不抛异常。
    """
    try:
        for shape in _iter_design_text_shapes(slide.shapes):
            try:
                if _is_residual_placeholder(shape.text_frame.text):
                    _set_text_keep_style(shape, [])
            except Exception:  # noqa: BLE001
                continue
    except Exception:  # noqa: BLE001
        pass


def _is_english_semantic_label(text: str) -> bool:
    """是否是「该清掉的英文语义小标签」(EXPERIENCE/AWARDS/BA PRODUCT DESIGN)。

    判定: 含 ≥3 个连续英文字母(英文单词) 且 不含中文。装饰编号 ONE/TWO/
    THREE/FOUR 也是英文单词,但它们是「位置编号」装饰、配在正文上方读起来
    自然,豁免保留。纯数字/符号天然不命中。
    """
    import re

    t = text.strip()
    if not t:
        return False
    if re.search(r"[一-鿿]", t):  # 有中文 → 不是纯英文标签
        return False
    if not re.search(r"[A-Za-z]{3,}", t):  # 没有英文单词(纯数字/符号)
        return False
    DECOR_NUMS = {"one", "two", "three", "four", "five", "six",
                  "seven", "eight", "nine", "ten"}
    if t.lower() in DECOR_NUMS:
        return False
    return True


def _fill_design_bullets(info: dict[str, Any], bullets: Sequence[str]) -> None:
    clean = [str(item).strip() for item in bullets if str(item).strip()]
    bodies = list(info.get("bodies", []))
    if not bodies:
        return
    if not clean:
        _clear_design_body_slots(info)
        return
    if len(bodies) >= len(clean):
        for shape, bullet in zip(bodies, clean):
            _set_text_keep_style(shape, [bullet])
            _enable_text_autofit(shape)
        for shape in bodies[len(clean):]:
            _set_text_keep_style(shape, [])
        return
    for shape, bullet in zip(bodies[:-1], clean[: max(len(bodies) - 1, 0)]):
        _set_text_keep_style(shape, [bullet])
        _enable_text_autofit(shape)
    _set_text_keep_style(bodies[-1], clean[len(bodies) - 1:])
    _enable_text_autofit(bodies[-1])


def _insert_design_picture(slide, so: SlideOutline) -> None:
    try:
        if not so.image_path or not Path(str(so.image_path)).is_file():
            return
        picture_ph = _template_placeholders(slide, PP_PLACEHOLDER.PICTURE)
        if picture_ph:
            try:
                picture_ph[0].insert_picture(str(so.image_path))
            except Exception as exc:  # noqa: BLE001
                log.debug("design picture insert failed: %s", exc)
    except Exception:  # noqa: BLE001
        pass


def _swap_design_picture(slide, image_path: str) -> bool:
    """把设计页里【最大的嵌入图片】的图源换成 AI 生成图(保留原图位置/裁切/
    边框样式) —— 这是「模板丰富内容 + 定制 AI 视觉」的关键: 用模板的多段
    文字排版 + 专业设计,但配图换成贴合主题的 AI 图,而非模板通用库存照。

    技术: 找最大 Picture shape → 经 a:blip r:embed 定位 image part →
    替换其 _blob。AI 图(PNG)转 JPEG 以匹配多数模板的 jpeg part(跨
    PowerPoint/WPS 更稳)。成功返回 True。从不抛异常。
    """
    try:
        src = Path(str(image_path))
        if not src.is_file():
            return False
        from pptx.oxml.ns import qn

        # 选面积最大的 Picture(通常是主视觉位)
        best = None
        best_area = -1
        for sh in slide.shapes:
            if getattr(sh, "shape_type", None) == 13:  # PICTURE
                try:
                    area = int(sh.width) * int(sh.height)
                except Exception:  # noqa: BLE001
                    area = 0
                if area > best_area:
                    best_area, best = area, sh
        if best is None:
            return False

        blips = best._element.findall(".//" + qn("a:blip"))
        if not blips:
            return False
        r_id = blips[0].get(qn("r:embed"))
        if not r_id:
            return False
        image_part = slide.part.related_part(r_id)

        # AI 图转 JPEG(白底), 匹配模板 part 的 content-type, 避免某些
        # 阅读器对 ext/content-type 不一致挑剔。
        new_bytes = src.read_bytes()
        try:
            from io import BytesIO
            from PIL import Image  # type: ignore

            with Image.open(src) as im:
                if im.mode in ("RGBA", "P", "LA"):
                    bg = Image.new("RGB", im.size, (255, 255, 255))
                    bg.paste(im.convert("RGBA"), mask=im.convert("RGBA").split()[-1])
                    im = bg
                else:
                    im = im.convert("RGB")
                buf = BytesIO()
                im.save(buf, format="JPEG", quality=88)
                new_bytes = buf.getvalue()
        except Exception:  # noqa: BLE001
            pass  # PIL 不可用 → 直接塞原字节(WPS 实测可读)

        image_part._blob = new_bytes
        return True
    except Exception as exc:  # noqa: BLE001
        log.debug("design picture swap failed: %s", exc)
        return False


def _rects_overlap(a, b, *, min_frac: float = 0.2) -> bool:
    """两形状矩形相交面积 > min_frac × 较小形状面积 → 视为重叠。
    判断失败保守返回 False(不误清装饰)。"""
    try:
        ax1, ay1 = int(a.left), int(a.top)
        ax2, ay2 = ax1 + int(a.width), ay1 + int(a.height)
        bx1, by1 = int(b.left), int(b.top)
        bx2, by2 = bx1 + int(b.width), by1 + int(b.height)
        iw = min(ax2, bx2) - max(ax1, bx1)
        ih = min(ay2, by2) - max(ay1, by1)
        if iw <= 0 or ih <= 0:
            return False
        inter = iw * ih
        amin = min((ax2 - ax1) * (ay2 - ay1), (bx2 - bx1) * (by2 - by1))
        return amin > 0 and inter > min_frac * amin
    except Exception:  # noqa: BLE001
        return False


def _scale_shape_runs(shape, scale: float) -> None:
    """把形状内所有 run 的字号 × scale(视觉评审 shrink_text 用于模板页)。
    没显式字号的 run 跳过(继承模板样式)。从不抛异常。"""
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                sz = run.font.size
                if sz is not None:
                    run.font.size = Pt(max(9, int(sz.pt * scale)))
    except Exception:  # noqa: BLE001
        pass


def _fill_design_page(slide, so: SlideOutline, info: dict[str, Any]) -> None:
    title_shape = info.get("title")
    subtitle_shape = info.get("subtitle")
    bodies = list(info.get("bodies", []))
    labels = list(info.get("labels", []))
    _divider_target = None  # title/section 页填了副标题的形状(保留它,清其余)

    if title_shape is not None:
        title_text = so.quote if so.layout == "quote" else so.title
        _set_text_keep_style(title_shape, [title_text], fit_pt=True)
        _ensure_readable_text(title_shape)

    if so.layout == "title":
        target = subtitle_shape or (bodies[0] if bodies else None)
        if target is None and so.subtitle:
            target = _take_subtitle_from_decor(info)
        if target is not None and so.subtitle:
            _set_text_keep_style(target, [so.subtitle], fit_pt=True)
            _ensure_readable_text(target)
            _clear_design_body_slots(info, keep=[target])
        else:
            _clear_design_body_slots(info)
        _divider_target = target
    elif so.layout == "section":
        target = bodies[0] if bodies else subtitle_shape
        if target is None and so.subtitle:
            target = _take_subtitle_from_decor(info)
        if target is not None and so.subtitle:
            _set_text_keep_style(target, [so.subtitle], fit_pt=True)
            _ensure_readable_text(target)
            _clear_design_body_slots(info, keep=[target])
        else:
            _clear_design_body_slots(info)
        _divider_target = target
    elif so.layout == "quote":
        cite_target = labels[0] if labels else (bodies[0] if bodies else None)
        if cite_target is not None and so.cite:
            _set_text_keep_style(cite_target, [so.cite])
            _clear_design_body_slots(info, keep=[cite_target])
        else:
            _clear_design_body_slots(info)
    elif so.layout == "two_column":
        if bodies:
            _set_text_keep_style(
                bodies[0],
                ([so.left_title] if so.left_title else []) + list(so.left),
            )
            _enable_text_autofit(bodies[0])
        if len(bodies) > 1:
            _set_text_keep_style(
                bodies[1],
                ([so.right_title] if so.right_title else []) + list(so.right),
            )
            _enable_text_autofit(bodies[1])
        for shape in bodies[2:]:
            _set_text_keep_style(shape, [])
    elif so.layout in {"image", "image_full"}:
        _insert_design_picture(slide, so)
        lines = list(so.bullets)
        if so.caption:
            lines = lines + [so.caption]
        _fill_design_bullets(info, lines)
    else:
        _fill_design_bullets(info, so.bullets)

    # 清掉残留的英文艺术大字(decor)：原模板的主题词叠排(About Me /
    # graphic designer 等)不清会和中文内容穿帮混排。形状本身保留
    # (底色/位置仍是设计的一部分)，只清文本。
    for shape in info.get("decor", []):
        _set_text_keep_style(shape, [])

    # 清掉英文语义小标签(EXPERIENCE/AWARDS/BA PRODUCT DESIGN 等)：中文
    # deck 里残留英文小标题会穿帮。纯数字/罗马序号(ONE 这类装饰编号靠
    # _is_decor_label 豁免)保留。
    for shape in info.get("labels", []):
        try:
            txt = (shape.text_frame.text or "").strip()
        except Exception:  # noqa: BLE001
            txt = ""
        if txt and _is_english_semantic_label(txt):
            _set_text_keep_style(shape, [])

    # 装饰编号(1/2/3 大数字、ONE/TWO/THREE 标签)若与【已填内容】矩形重叠
    # → 清掉。模板设计里它们是衬在英文短词后面的背景装饰,中文长句填进
    # 槽后必撞(用户实测「目录大数字与文字重叠」的根因)。不重叠的保留。
    _filled = [s for s in [title_shape, subtitle_shape, _divider_target] if s is not None]
    _filled += [b for b in bodies if b is not None]
    _DECOR_WORDS = {"one", "two", "three", "four", "five", "six",
                    "seven", "eight", "nine", "ten"}
    _candidates = list(info.get("decor_nums", []))
    for shape in info.get("labels", []):
        try:
            t = (shape.text_frame.text or "").strip().lower()
        except Exception:  # noqa: BLE001
            t = ""
        if t in _DECOR_WORDS or (t and (t.isdigit() or len(t) <= 2)):
            _candidates.append(shape)
    for shape in _candidates:
        try:
            if any(_rects_overlap(shape, f) for f in _filled):
                _set_text_keep_style(shape, [])
        except Exception:  # noqa: BLE001
            continue

    # 章节/标题/引用页是「分隔页」: 本就只该有标题 + 副标题,模板自带的
    # 装饰副标题/标语(如"健康教育/心理辅导/重拾信心" —— 真实中文,占位
    # 清扫和英文标签都抓不到)留着必穿帮。清掉除标题/已填副标题外的所有
    # 文字。从不抛异常。
    if so.layout in {"title", "section", "quote"}:
        # 身份用底层 lxml 元素(python-pptx 每次迭代生成新 shape 包装,
        # id(shape) 不稳定;_element 是同一个节点)。
        _keep_ids = set()
        for s in (title_shape, _divider_target):
            if s is not None:
                try:
                    _keep_ids.add(id(s._element))
                except Exception:  # noqa: BLE001
                    pass
        try:
            for shape in _iter_design_text_shapes(slide.shapes):
                try:
                    if id(shape._element) in _keep_ids:
                        continue
                    if (shape.text_frame.text or "").strip():
                        _set_text_keep_style(shape, [])
                except Exception:  # noqa: BLE001
                    continue
        except Exception:  # noqa: BLE001
            pass

    # 有 AI 生成图(image_prompt 出图后的 image_path) → 换进本页最大图片位,
    # 把模板通用库存照替成贴合主题的定制视觉(模板丰富内容 + AI 视觉合一)。
    if so.image_path and Path(str(so.image_path)).is_file():
        _swap_design_picture(slide, str(so.image_path))

    # 兜底清扫所有残留占位文(中英) —— 复杂版式漏清的"输入标题内容"/
    # 英文 lorem 在这里统一清掉,防穿帮。
    _sweep_residual_placeholders(slide)

    # 视觉评审 shrink_text(模板页): 缩本页已填文字槽的字号。
    if so.font_scale and so.font_scale < 1.0:
        for shape in [info.get("title"), info.get("subtitle")] + list(info.get("bodies", [])):
            if shape is not None:
                _scale_shape_runs(shape, so.font_scale)

    # 垂直溢出根治: 已填正文槽按槽高主动缩字号(WPS 不执行 normAutofit
    # 声明,长中文必溢出框底压页边 —— 用户真机三页实测)。
    for shape in list(info.get("bodies", [])) + (
        [_divider_target] if _divider_target is not None else []
    ):
        if shape is not None:
            _fit_text_to_shape(shape)

    if so.notes:
        try:
            slide.notes_slide.notes_text_frame.text = so.notes
        except Exception:  # noqa: BLE001
            pass


def _drop_slide_id(prs, sld_id) -> None:
    _R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    try:
        r_id = sld_id.get(_R_ID)
        if r_id:
            try:
                prs.part.drop_rel(r_id)
            except Exception:  # noqa: BLE001
                pass
        prs.slides._sldIdLst.remove(sld_id)
    except Exception:  # noqa: BLE001
        pass


def _render_with_design_pages(
    slides: list[SlideOutline], template_path: str, *, title: str, author: str, out_path: Path,
    banned_pages: Optional[dict[int, set[int]]] = None,
) -> Optional[dict[str, Any]]:
    """banned_pages: {outline 页序 i: 禁用的设计页 index 集合} —— 视觉评审
    change_page 动作把上一轮该页用过的设计页 ban 掉,本轮换一张重填。"""
    try:
        prs = _Presentation(template_path)
        if len(prs.slides) < 3:
            return None

        pages = [
            {"index": idx, "slide": slide, "info": _analyze_design_page(
                slide, int(prs.slide_width or 0), int(prs.slide_height or 0))}
            for idx, slide in enumerate(prs.slides)
        ]
        # 按规整度给每页打分,选页时优先挑高分页(避开竖排标题/图标网格
        # 等易穿帮版式) —— 真机看图驱动的修复。
        for page in pages:
            page["score"] = _page_quality_score(page["info"])
        section_pages = sorted(
            [p for p in pages[1:]
             if p["info"]["title_pt"] >= 80 and p["info"]["n_body"] <= 1],
            key=lambda p: (-p["score"], p["index"]),
        )
        content_pages = sorted(
            [p for p in pages[1:] if p["info"]["n_body"] >= 1],
            key=lambda p: (-p["score"], p["index"]),
        )

        used: set[int] = set()
        selected: list[tuple[dict[str, Any], SlideOutline]] = []
        page_map: list[int] = []  # 第 k 个成品页 ← 设计页 index(给闭环 ban 用)
        for i, so in enumerate(slides):
            # 该 outline 页被 ban 的设计页先记成 used 再选(选完恢复,不影响后页)
            ban = (banned_pages or {}).get(i, set())
            eff_used = used | set(ban)
            page = _select_design_page(so, pages, section_pages, content_pages, eff_used)
            if page is None and ban:
                # ban 后无页可选 → 放开 ban 兜底(宁可重复风格也别丢页)
                page = _select_design_page(so, pages, section_pages, content_pages, used)
            if page is None:
                log.debug("design page exhausted, skipping slide: %s", so.title)
                continue
            used.add(page["index"])
            selected.append((page, so))
            page_map.append(page["index"])
        if not selected:
            return None

        for page, so in selected:
            _fill_design_page(page["slide"], so, page["info"])

        xml_slides = prs.slides._sldIdLst
        sld_ids = list(xml_slides)
        selected_ids = [sld_ids[page["index"]] for page, _ in selected]
        selected_id_set = {id(sld_id) for sld_id in selected_ids}
        for sld_id in selected_ids:
            xml_slides.append(sld_id)
        for sld_id in list(xml_slides):
            if id(sld_id) not in selected_id_set:
                _drop_slide_id(prs, sld_id)

        try:
            cp = prs.core_properties
            if title:
                cp.title = title
            if author:
                cp.author = author
                cp.last_modified_by = author
        except Exception:  # noqa: BLE001
            pass

        prs.save(out_path)
        return {
            "ok": True,
            "path": str(out_path),
            "slide_count": len(selected),
            "theme": "template-design",
            "page_map": page_map,
            "artifacts": [{
                "kind": "file",
                "path": str(out_path),
                "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "title": Path(str(out_path)).name,
            }],
        }
    except Exception as exc:  # noqa: BLE001
        log.warning("template design render failed: %s", exc, exc_info=True)
        return None


def _render_with_template(
    slides: list[SlideOutline], template_path: str, *, title: str, author: str, out_path: Path,
) -> Optional[dict[str, Any]]:
    try:
        prs = _Presentation(template_path)

        # 模板文件常带示例页；这里只保留 master/layout/theme，避免用户示例页混入输出。
        # 关键：只从 _sldIdLst remove(sldId) 仅删「引用」，底层 ppt/slides/slideN.xml
        # 部件仍残留 → 新加 slide 复用 slide1.xml 等名字 → zip 内 Duplicate name →
        # WPS/PowerPoint 会渲染【旧示例页】而非填充内容（实测「graphic designer」串图
        # 根因）。必须同时 drop_rel 真正解除关系，让残留部件不被写回。
        _R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        xml_slides = prs.slides._sldIdLst
        for sldId in list(xml_slides):
            rId = sldId.get(_R_ID)
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except Exception:  # noqa: BLE001
                    pass
            xml_slides.remove(sldId)

        for so in slides:
            layout = _pick_template_layout(prs, so.layout)
            slide = prs.slides.add_slide(layout)

            title_ph = _template_placeholders(
                slide, PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
            )
            if title_ph:
                _set_placeholder_text(title_ph[0], so.title)

            subtitle_ph = _template_placeholders(slide, PP_PLACEHOLDER.SUBTITLE)
            if subtitle_ph:
                _set_placeholder_text(subtitle_ph[0], so.subtitle)

            body_ph = _template_placeholders(slide, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
            object_ph = _template_placeholders(slide, PP_PLACEHOLDER.OBJECT)

            if so.layout == "two_column":
                targets = object_ph if len(object_ph) >= 2 else body_ph
                if targets:
                    _fill_placeholder_lines(
                        targets[0],
                        ([so.left_title] if so.left_title else []) + so.left,
                        bold_first=bool(so.left_title),
                    )
                if len(targets) > 1:
                    _fill_placeholder_lines(
                        targets[1],
                        ([so.right_title] if so.right_title else []) + so.right,
                        bold_first=bool(so.right_title),
                    )
            elif so.layout in {"section", "quote"}:
                lines = [so.subtitle] if so.layout == "section" else [so.quote, so.cite]
                if body_ph:
                    _fill_placeholder_lines(body_ph[0], lines)
            elif so.layout not in {"title", "image", "image_full"} and body_ph:
                _fill_placeholder_lines(body_ph[0], so.bullets)

            if so.layout in {"image", "image_full"}:
                picture_ph = _template_placeholders(slide, PP_PLACEHOLDER.PICTURE)
                if picture_ph and so.image_path and Path(str(so.image_path)).is_file():
                    try:
                        picture_ph[0].insert_picture(str(so.image_path))
                    except Exception as exc:  # noqa: BLE001
                        log.debug("template picture insert failed: %s", exc)
                caption_targets = [
                    ph for ph in body_ph
                    if ph not in picture_ph and ph not in title_ph and ph not in subtitle_ph
                ]
                if caption_targets and so.caption:
                    _set_placeholder_text(caption_targets[0], so.caption)

            if so.notes:
                slide.notes_slide.notes_text_frame.text = so.notes

        try:
            cp = prs.core_properties
            if title:
                cp.title = title
            if author:
                cp.author = author
                cp.last_modified_by = author
        except Exception:  # noqa: BLE001
            pass

        prs.save(out_path)
        return {
            "ok": True,
            "path": str(out_path),
            "slide_count": len(slides),
            "theme": "template",
            "artifacts": [{
                "kind": "file",
                "path": str(out_path),
                "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "title": Path(str(out_path)).name,
            }],
        }
    except Exception as exc:  # noqa: BLE001
        log.warning("template render failed: %s", exc, exc_info=True)
        return None


# ---------------------------------------------------------------------
# Public entrypoint
# ---------------------------------------------------------------------


def ppt_create(
    outline: Any,
    *,
    theme: str = "minimal",
    title: str = "",
    author: str = "DeskPet",
    output_path: Optional[str] = None,
    template: Optional[str] = None,
    dry_run: bool = False,
    skip_image_gen: bool = False,
) -> dict[str, Any]:
    """Render an outline into a ``.pptx`` file on disk.

    Parameters
    ----------
    outline:
        Either a JSON string of a list of slide dicts, or a Python list
        of dicts, or a single dict. See :class:`SlideOutline` for shape.
    theme:
        One of ``minimal`` / ``dark`` / ``playful``. Falls back to
        ``minimal`` on typo.
    title:
        Optional deck title. If the first slide is ``layout=title`` the
        outline's title wins; otherwise this is used as the document
        property only.
    author:
        Set on the document core properties.
    output_path:
        Absolute path. Defaults to ``<tempdir>/deskpet-ppt-<ts>.pptx``.
    template:
        Optional bundled template name or .pptx template path. When
        valid, DeskPet uses the template slide layouts and fills
        placeholders so formatting is inherited from the template.

    Returns
    -------
    dict
        On success: ``{"ok": True, "path": str, "slide_count": int, "theme": str}``.
        On failure: ``{"ok": False, "error": str, "markdown_fallback": str}``.
    """
    slides = parse_outline(outline)
    if not slides:
        return {
            "ok": False,
            "error": "outline parse failed or empty",
            "markdown_fallback": render_markdown_fallback(
                slides, title=title, author=author,
            ),
        }

    # WI-T1.6: dry_run 预览模式（PRD §3 D9）。返回 outline markdown 作为
    # text artifact，不写 .pptx — 用户/LLM 可先看大纲再决定是否真生成。
    # 显式 emit artifacts[] 走 D1 一等公民路径。
    if dry_run:
        md = render_markdown_fallback(slides, title=title, author=author)
        return {
            "ok": True,
            "dry_run": True,
            "slide_count": len(slides),
            "artifacts": [{
                "kind": "text",
                "title": (title or "outline") + " (preview)",
                "preview": md,
            }],
        }

    if not _HAS_PPTX:
        return {
            "ok": False,
            "error": "python-pptx not installed; install with `pip install python-pptx`",
            "markdown_fallback": render_markdown_fallback(
                slides, title=title, author=author,
            ),
        }

    out_path = _resolve_output_path(output_path)

    # 视觉风格仲裁(防 AI 整页生图 与 模板模式 互踩 + 防白烧钱)：
    # - 显式 template= 永远优先(用户明确要模板)。
    # - 否则若有 image_full 页带 image_prompt → 用户要 AI 整页铺图,不让
    #   env 默认模板(DESKPET_PPT_DEFAULT_TEMPLATE)劫持。
    wants_fullbleed = any(
        so.layout == "image_full" and so.image_prompt for so in slides
    )
    chosen_template = template or (None if wants_fullbleed else _default_template())
    # 生图消费者: ①模板模式 → AI 图换进设计页图片位(_swap_design_picture,
    # 丰富内容+定制视觉) ②image_full/image 页 → 全幅铺图/插图。两者都不沾
    # 才跳过生图省钱($0.15/张)。
    has_prompts = any(so.image_prompt for so in slides)
    has_img_layout = any(
        so.layout in {"image_full", "image"} and so.image_prompt for so in slides
    )
    if not skip_image_gen and has_prompts and (chosen_template or has_img_layout):
        _autofill_image_prompts(slides)
    else:
        # 没走 autofill(图已预设/无 prompt)也要给 image_full 页分配版式 ——
        # 版式是渲染属性,与是否现场生图无关。
        _assign_image_layouts(slides)
    if chosen_template:
        # 大类名 → 按预览图视觉选一套;路径/stem → 直接解析。topic 给 picker 选最贴的设计。
        topic = title or (slides[0].title if slides else "")
        resolved = _resolve_template_for_render(chosen_template, topic=topic)
        if resolved:
            result = _render_with_design_pages(
                slides, resolved, title=title, author=author, out_path=out_path,
            )
            if result is not None:
                log.debug("ppt template rendered with design pages: %s", resolved)
                # 视觉评估闭环(模板版): 看每页 → 缩字号/换设计页 → 重渲染
                _visual_review_loop_template(
                    slides, resolved, out_path,
                    title=title, author=author, result=result,
                )
                _maybe_render_preview(result)
                return result
            result = _render_with_template(
                slides, resolved, title=title, author=author, out_path=out_path,
            )
            if result is not None:
                log.debug("ppt template rendered with layouts: %s", resolved)
                _maybe_render_preview(result)
                return result
            log.warning("template render failed, falling back to from-scratch engine")
        else:
            log.warning("template not found: %s - falling back to from-scratch engine", chosen_template)

    theme_obj = get_theme(theme)

    try:
        result = _render_fromscratch(
            slides, theme_obj, out_path, title=title, author=author,
        )
        # 视觉评估闭环(问题1): 桌宠「亲眼看」每页渲染图 → 评审 → 自动修
        # (换版式/缩字号) → 重渲染。仅 AI 图文 deck(有 image_full 页)走。
        _visual_review_loop(
            slides, theme_obj, out_path, title=title, author=author,
            result=result,
        )
        _maybe_render_preview(result)
        return result
    except Exception as exc:  # noqa: BLE001
        log.warning("ppt_create failed: %s", exc, exc_info=True)
        return {
            "ok": False,
            "error": f"render failed: {exc}",
            "markdown_fallback": render_markdown_fallback(
                slides, title=title, author=author,
            ),
        }


def _render_fromscratch(
    slides: list[SlideOutline],
    theme_obj: Theme,
    out_path: Path,
    *,
    title: str,
    author: str,
) -> dict[str, Any]:
    """from-scratch 引擎渲染并保存(可重入 — 视觉闭环改完 slides 再调一遍)。"""
    prs = _Presentation()
    prs.slide_width = _SLIDE_WIDTH
    prs.slide_height = _SLIDE_HEIGHT
    total = len(slides)
    blank_layout = prs.slide_layouts[6]  # 6 = "Blank" in default template
    for i, so in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        is_conclusion = _is_conclusion_slide(so, i, total)
        if is_conclusion:
            _render_conclusion_v2(slide, so, theme_obj)
        else:
            renderer = _RENDERERS.get(so.layout, _render_bullet_v2)
            renderer(slide, so, theme_obj)
        # Footer everywhere except the very first title slide for breathing room.
        if so.layout != "title":
            _add_footer(slide, theme_obj, i, total, dark=is_conclusion or so.layout == "image_full")
        # Speaker notes
        if so.notes:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = so.notes
    # Document core properties
    try:
        cp = prs.core_properties
        if title:
            cp.title = title
        if author:
            cp.author = author
            cp.last_modified_by = author
    except Exception:  # noqa: BLE001
        pass
    prs.save(out_path)
    # WI-T1.2 D1：显式 emit artifacts[]（一等公民路径，前端按 kind=file
    # 渲染 ArtifactCard；保留 path 字段保 BC）。
    return {
        "ok": True,
        "path": str(out_path),
        "slide_count": total,
        "theme": theme_obj.name,
        "artifacts": [{
            "kind": "file",
            "path": str(out_path),
            "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "title": Path(str(out_path)).name,
        }],
    }


def _ppt_visual_review_enabled() -> bool:
    """config ``[ppt].visual_review``(默认 True)。"""
    try:
        return _cfg_bool(_ppt_config_section(), "visual_review", True)
    except Exception:  # noqa: BLE001
        return True


def _apply_review_actions(
    slides: list[SlideOutline], reviews: list[dict[str, Any]],
) -> bool:
    """把视觉评审动作应用到 slides。返回是否有实际改动(需要重渲染)。"""
    changed = False
    for r in reviews:
        try:
            idx = int(r.get("page") or 0) - 1
            if idx < 0 or idx >= len(slides) or r.get("ok", True):
                continue
            so = slides[idx]
            action = str(r.get("action") or "ok")
            if action == "change_variant" and so.layout == "image_full":
                v = str(r.get("variant") or "").strip().lower()
                if v in IMAGE_VARIANTS and v != so.image_variant:
                    log.info(
                        "visual_review fix page=%d variant %s→%s issues=%s",
                        idx + 1, so.image_variant, v, r.get("issues"),
                    )
                    so.image_variant = v
                    changed = True
            elif action == "shrink_text":
                if so.font_scale > 0.72:
                    log.info(
                        "visual_review fix page=%d shrink_text issues=%s",
                        idx + 1, r.get("issues"),
                    )
                    so.font_scale = round(so.font_scale * 0.85, 2)
                    changed = True
        except Exception:  # noqa: BLE001
            continue
    return changed


def _visual_review_loop(
    slides: list[SlideOutline],
    theme_obj: Theme,
    out_path: Path,
    *,
    title: str,
    author: str,
    result: dict[str, Any],
    max_rounds: int = 2,
) -> None:
    """视觉评估闭环: 渲染截图 → 多模态评审 → 应用修复 → 重渲染,最多
    max_rounds 轮评审。任何失败静默跳过(主结果零影响)。从不抛异常。"""
    try:
        if _in_pytest():
            return
        if not _ppt_visual_review_enabled():
            return
        if not any(so.layout == "image_full" for so in slides):
            return  # 仅 AI 图文 deck 走视觉闭环(文本/模板deck不必烧 vision)
        renderer = _get_ppt_renderer()
        if not renderer.com_render_available():
            return
        render_fn = getattr(renderer, "render_pptx_to_pngs_safe", None)
        if not callable(render_fn):
            return
        from .ppt_visual_review import review_slides

        rounds_meta: list[dict[str, Any]] = []
        for rnd in range(1, max_rounds + 1):
            shot_dir = Path(str(out_path)).with_suffix(f".review{rnd}")
            pngs = render_fn(str(out_path), str(shot_dir), timeout=150.0)
            if not pngs:
                break
            meta = [
                {
                    "title": so.title,
                    "variant": so.image_variant or so.layout,
                    "n_bullets": len([b for b in (so.bullets or []) if str(b).strip()]),
                }
                for so in slides
            ]
            reviews = review_slides(pngs, meta)
            if not reviews:
                break
            n_bad = sum(1 for r in reviews if not r.get("ok", True))
            rounds_meta.append({
                "round": rnd,
                "pages": len(reviews),
                "issues": n_bad,
                "detail": [r for r in reviews if not r.get("ok", True)],
            })
            if n_bad == 0:
                break
            if not _apply_review_actions(slides, reviews):
                break  # 有问题但没有可自动修的动作 → 停(报告里留痕)
            _render_fromscratch(slides, theme_obj, out_path, title=title, author=author)
        if rounds_meta:
            result["visual_review"] = rounds_meta
            log.info(
                "visual_review_loop done rounds=%d final_issues=%d",
                len(rounds_meta),
                rounds_meta[-1]["issues"] if rounds_meta else -1,
            )
    except Exception as exc:  # noqa: BLE001
        log.warning("visual review loop failed: %s", str(exc)[:200])


def _apply_review_actions_template(
    slides: list[SlideOutline],
    reviews: list[dict[str, Any]],
    page_map: list[int],
    banned_pages: dict[int, set[int]],
) -> bool:
    """模板版动作应用: shrink_text 缩该页填充文字;change_page 把该页上一轮
    用的设计页 ban 掉(重渲染时换页重填)。返回是否需要重渲染。"""
    changed = False
    for r in reviews:
        try:
            idx = int(r.get("page") or 0) - 1
            if idx < 0 or idx >= len(slides) or r.get("ok", True):
                continue
            so = slides[idx]
            action = str(r.get("action") or "ok")
            if action == "shrink_text":
                if so.font_scale > 0.72:
                    log.info(
                        "visual_review(template) fix page=%d shrink_text issues=%s",
                        idx + 1, r.get("issues"),
                    )
                    so.font_scale = round(so.font_scale * 0.85, 2)
                    changed = True
            elif action == "change_page":
                if idx < len(page_map):
                    log.info(
                        "visual_review(template) fix page=%d change_page(ban %d) issues=%s",
                        idx + 1, page_map[idx], r.get("issues"),
                    )
                    banned_pages.setdefault(idx, set()).add(page_map[idx])
                    changed = True
        except Exception:  # noqa: BLE001
            continue
    return changed


def _visual_review_loop_template(
    slides: list[SlideOutline],
    template_path: str,
    out_path: Path,
    *,
    title: str,
    author: str,
    result: dict[str, Any],
    max_rounds: int = 2,
) -> None:
    """模板(design-fill)deck 的视觉评估闭环: 截图 → 评审(模板动作集) →
    缩字号/换设计页 → 重渲染。失败静默,主结果零影响。从不抛异常。"""
    try:
        if _in_pytest():
            return
        if not _ppt_visual_review_enabled():
            return
        renderer = _get_ppt_renderer()
        if not renderer.com_render_available():
            return
        render_fn = getattr(renderer, "render_pptx_to_pngs_safe", None)
        if not callable(render_fn):
            return
        from .ppt_visual_review import review_slides

        banned_pages: dict[int, set[int]] = {}
        rounds_meta: list[dict[str, Any]] = []
        for rnd in range(1, max_rounds + 1):
            shot_dir = Path(str(out_path)).with_suffix(f".review{rnd}")
            pngs = render_fn(str(out_path), str(shot_dir), timeout=150.0)
            if not pngs:
                break
            meta = [
                {
                    "title": so.title,
                    "variant": "template",
                    "n_bullets": len([b for b in (so.bullets or []) if str(b).strip()]),
                }
                for so in slides
            ]
            reviews = review_slides(pngs, meta, mode="template")
            if not reviews:
                break
            n_bad = sum(1 for r in reviews if not r.get("ok", True))
            rounds_meta.append({
                "round": rnd,
                "pages": len(reviews),
                "issues": n_bad,
                "detail": [r for r in reviews if not r.get("ok", True)],
            })
            if n_bad == 0:
                break
            page_map = list(result.get("page_map") or [])
            if not _apply_review_actions_template(slides, reviews, page_map, banned_pages):
                break
            new_result = _render_with_design_pages(
                slides, template_path, title=title, author=author,
                out_path=out_path, banned_pages=banned_pages,
            )
            if not new_result:
                break
            result["page_map"] = new_result.get("page_map")
            result["slide_count"] = new_result.get("slide_count", result.get("slide_count"))
        if rounds_meta:
            result["visual_review"] = rounds_meta
            log.info(
                "visual_review_loop(template) done rounds=%d final_issues=%d",
                len(rounds_meta),
                rounds_meta[-1]["issues"] if rounds_meta else -1,
            )
    except Exception as exc:  # noqa: BLE001
        log.warning("template visual review loop failed: %s", str(exc)[:200])


def _resolve_output_path(p: Optional[str]) -> Path:
    if p:
        path = Path(p).expanduser().resolve()
        path.parent.mkdir(parents=True, exist_ok=True)
        return path
    ts = int(time.time())
    fname = f"deskpet-ppt-{ts}.pptx"
    # 默认落 <user_data>/OutPut/PPT/(用户好找,系统 temp 没人翻得到);
    # paths 不可用(独立脚本等)再回退 temp。
    try:
        from paths import output_dir  # type: ignore[import-not-found]

        return output_dir("PPT") / fname
    except Exception:  # noqa: BLE001
        return Path(tempfile.gettempdir()) / fname


# ---------------------------------------------------------------------
# Tool registry wiring
# ---------------------------------------------------------------------


_PPT_SCHEMA = {
    "name": "ppt_create",
    "description": (
        "【注意：若用户只给主题就要做 PPT（含「惊艳/AI配图」需求），请改用 ppt_pro —— "
        "它会自动调研+大纲确认+优先生图/连不上回退模板。本工具仅用于：用户已给出完整 outline / "
        "明确要跳过调研和确认 / 只要朴素快出 / 要一次做模板与AI各一份。】\n"
        "Generate a professional .pptx presentation locally from an outline. "
        "Returns the file path on success; falls back to a Markdown outline "
        "when python-pptx is unavailable. 生成成功后【必须】把返回的 path 完整"
        "路径告诉用户(用户要知道文件存哪了)。Use this AFTER you've decided on a "
        "structured slide outline. Do not stuff long paragraphs into bullets "
        "— bullets are cues, not scripts.\n"
        "两种视觉风格(二选一,别混用):\n"
        "① AI 整页生图(最惊艳,适合'用AI配图/惊艳/视觉冲击/封面海报感'的需求): "
        "每页 layout='image_full' + 写一段英文 image_prompt,DeskPet 会用 "
        "AI 出图生成电影感全屏背景图铺满整页、标题压在底部暗带。"
        "这种模式【不要】传 template 参数。每张图约 1~2 分钟,4 页请耐心等。\n"
        "② 模板填充(可编辑/正式商务): 传 template=模板名,用模板的设计页填文字。"
        "这种模式【不要】给页面写 image_prompt(模板自带配图,AI 图用不上会白生成)。"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "outline": {
                "description": (
                    "List of slide dicts, OR a JSON string of such a list. "
                    "Each slide: {layout, title, subtitle?, bullets?, left?, "
                    "right?, left_title?, right_title?, image_path?, image_prompt?, caption?, "
                    "quote?, cite?, notes?}. layout ∈ {title, section, bullet, "
                    "two_column, image, image_full, quote, toc}.\n"
                    "image_full = AI 全幅生图页(惊艳+丰富,首选): 配 image_prompt"
                    "(英文,电影感全屏画面,深色调、主体偏右、no text no watermark)。"
                    "★封面页只给 title+caption → 全幅图+底部暗带;★内容页给 title+"
                    "bullets(3-4 条,每条≤24字精炼) → 全幅图+左侧深色面板放标题和要点"
                    "(我方自控排版,不挤)。这样既惊艳又有内容,做整份 AI 视觉 PPT 时"
                    "每页都用 image_full(封面无 bullets,内容页带 bullets)。\n"
                    "image_prompt = 该页要 AI 生成的图。image_full 页必配。"
                    "传了 template 时也可写 image_prompt: 会把 AI 图换进模板的图片位"
                    "(模板专业排版 + 定制 AI 视觉)。"
                ),
                "type": ["array", "string"],
            },
            "theme": {
                "description": "Visual theme. minimal=clean business; dark=tech demos; playful=marketing.",
                "type": "string",
                "enum": list(VALID_THEMES),
                "default": "minimal",
            },
            "title": {"type": "string", "description": "Document title (core properties)."},
            "author": {"type": "string", "description": "Author name. Defaults to DeskPet."},
            "output_path": {
                "type": "string",
                "description": "Absolute output path. Defaults to a temp file.",
            },
            "template": {
                "type": "string",
                "description": (
                    "Optional .pptx template path. When provided and valid, DeskPet "
                    "loads it, adds slides from its layouts, fills placeholders, and "
                    "inherits editable formatting from the template."
                ),
            },
            "dry_run": {
                "type": "boolean",
                "description": (
                    "WI-T1.6 outline 预览模式。True → 不写 .pptx，仅返回 "
                    "outline markdown 作为 text artifact 供用户确认。建议 ≥ 5 "
                    "张幻灯片的 deck 先 dry_run=true 让用户审查 outline，再 "
                    "dry_run=false 实际生成。"
                ),
                "default": False,
            },
        },
        "required": ["outline"],
    },
}


def _build_template_description() -> str:
    """工具 schema 里 ``template`` 参数的说明 —— 列出可选【大类】(不是单个模板)。

    模板库有上百套,无法逐一列名。LLM 只挑一个大类名,引擎再按预览图视觉
    选该类里最贴主题的一套设计。也支持直传一个 .pptx 绝对路径。
    """
    try:
        from . import ppt_template_picker as _tpl

        cats = _tpl.list_categories()
    except Exception:  # noqa: BLE001
        cats = []
    if not cats:
        return (
            "Optional .pptx absolute path. 当前模板库不可用,要模板填充请传"
            "一个 .pptx 绝对路径。"
        )
    lines = []
    for cat in cats:
        key = cat.get("key", "")
        hint = ""
        try:
            hint = _tpl.CATEGORY_STYLE_HINTS.get(key, "")
        except Exception:  # noqa: BLE001
            hint = ""
        lines.append(f"「{key}」{hint}(约 {cat.get('count', '?')} 套)".rstrip())
    listing = "；".join(lines)
    return (
        "可编辑模板填充。用户要『正式/专业/精美/可编辑』PPT 时优先用模板:"
        "从下列【大类】里【按名字精确】选一个传入(只传大类名,别自己编名字),"
        "桌宠会按预览图为你的主题挑该类里最合适的一套设计。也可直传 .pptx 绝对路径。"
        "按主题选最贴的风格：\n"
        f"{listing}。"
    )


def _build_ppt_schema() -> dict[str, Any]:
    schema = {
        **_PPT_SCHEMA,
        "parameters": {
            **_PPT_SCHEMA["parameters"],
            "properties": {
                **_PPT_SCHEMA["parameters"]["properties"],
            },
        },
    }
    schema["parameters"]["properties"]["template"] = {
        **schema["parameters"]["properties"]["template"],
        "description": _build_template_description(),
    }
    return schema


_PPT_SCHEMA = _build_ppt_schema()


def _ppt_async_enabled() -> bool:
    """config ``[ppt].async_enabled``(默认 True)。带 AI 生图的 deck 走后台。"""
    try:
        return _cfg_bool(_ppt_config_section(), "async_enabled", True)
    except Exception:  # noqa: BLE001
        return True


def _count_image_prompt_pages(outline: Any) -> int:
    try:
        return sum(1 for so in parse_outline(outline) if so.image_prompt)
    except Exception:  # noqa: BLE001
        return 0


def _handle_ppt_create(args: dict, task_id: str) -> str:
    """Sync handler wired into the tool registry.

    纯文本/模板填充是同步快路径。但带 image_prompt 的 deck 要串行调
    gpt-image-2(每张 70~180s,N 页可达数分钟) —— 同步会把 agent 回合
    卡死。这种 deck 走后台异步: 秒回「制作中」,ImageGenerationWorker
    的事件循环上跑完整生成,做好用 notifier 推回桌宠 + 自动打开。
    """
    kwargs = dict(
        theme=str(args.get("theme") or "minimal"),
        title=str(args.get("title") or ""),
        author=str(args.get("author") or "DeskPet"),
        output_path=(str(args["output_path"]) if args.get("output_path") else None),
        template=(str(args["template"]) if args.get("template") else None),
        dry_run=bool(args.get("dry_run", False)),
    )
    outline = args.get("outline")

    worker = args.get("_image_worker")
    sid = str(args.get("_session_id") or "default")
    n_imgs = _count_image_prompt_pages(outline)
    can_async = (
        n_imgs >= 1
        and not kwargs["dry_run"]
        and worker is not None
        and getattr(worker, "alive", lambda: False)()
        and _ppt_async_enabled()
    )

    if can_async:
        async def _bg_job() -> None:
            loop = asyncio.get_running_loop()
            try:
                result = await loop.run_in_executor(
                    None, lambda: ppt_create(outline, **kwargs)
                )
            except Exception as exc:  # noqa: BLE001
                result = {"ok": False, "error": f"{type(exc).__name__}: {exc}"}
            try:
                path = result.get("path")
                if result.get("ok") and path:
                    _open_image_file(path)
                    await worker.notifier(
                        sid,
                        f"✨ 图文 PPT 做好啦！已自动打开～\n"
                        f"📁 {Path(str(path)).name}（{result.get('slide_count', '?')} 页）\n"
                        f"📂 保存在：{path}",
                    )
                else:
                    await worker.notifier(
                        sid, f"😿 PPT 没做成：{result.get('error') or '未知错误'}"
                    )
            except Exception as exc:  # noqa: BLE001
                log.warning("ppt async notify failed: %s", exc)

        if worker.submit_background(lambda: _bg_job()):
            return json.dumps(
                {
                    "ok": True,
                    "status": "generating",
                    "message": (
                        f"🎨 在做图文 PPT 啦～要给 {n_imgs} 页各画一张 AI 图，"
                        f"约几分钟，做好我自动打开发你，这会儿可以先聊点别的~"
                    ),
                },
                ensure_ascii=False,
            )
        # submit 失败 → 落同步路径(慢但出图)
        log.warning("ppt async submit failed; falling back to sync")

    result = ppt_create(outline, **kwargs)
    return json.dumps(result, ensure_ascii=False)


def _open_image_file(path: str) -> bool:
    """复用 image_tools 的 OS 打开(桌宠场景,故意打开)。从不抛异常。"""
    try:
        from .image_tools import _open_file as _of  # type: ignore

        return bool(_of(Path(str(path))))
    except Exception:  # noqa: BLE001
        return False


async def _maybe_await(value):
    if inspect.isawaitable(value):
        return await value
    return value


async def _notify(notifier, session_id: str, message: str) -> None:
    if notifier is None:
        return
    try:
        await _maybe_await(notifier(session_id, message))
    except Exception as exc:  # noqa: BLE001
        log.warning("ppt_pro notify failed: %s", exc)


def _sync_notify(notifier, session_id: str, loop: asyncio.AbstractEventLoop):
    def _send(message: str) -> None:
        try:
            asyncio.run_coroutine_threadsafe(_notify(notifier, session_id, message), loop)
        except Exception as exc:  # noqa: BLE001
            log.warning("ppt_pro sync notify failed: %s", exc)

    return _send


def set_ppt_pro_services(
    *,
    outline_propose=None,
    notifier=None,
    run_blocking=None,
    artifact_pusher=None,
    receipt_reporter=None,
) -> None:
    _PPT_PRO_CTX.update({
        "outline_propose": outline_propose,
        "notifier": notifier,
        "run_blocking": run_blocking,
        "artifact_pusher": artifact_pusher,
        "receipt_reporter": receipt_reporter,
    })


def _ppt_pro_cancel(sid: str) -> None:
    task = _PPT_PRO_RUNNING.get(str(sid or "default"))
    if task is not None and not task.done():
        task.cancel()


async def _run_blocking_call(run_blocking, fn):
    return await _maybe_await(run_blocking(fn))


async def _ppt_pro_report_done(result: dict, *, notifier, session_id: str) -> None:
    artifact_pusher = _PPT_PRO_CTX.get("artifact_pusher")
    receipt_reporter = _PPT_PRO_CTX.get("receipt_reporter")
    try:
        if result.get("ok"):
            path = result.get("path")
            artifacts = list(result.get("artifacts") or [])
            text = f"PPT 已生成：{path}" if path else "PPT 已生成。"
            if artifact_pusher is not None:
                try:
                    await _maybe_await(artifact_pusher(session_id, artifacts, text))
                except Exception as exc:  # noqa: BLE001
                    log.warning("ppt_pro artifact push failed: %s", exc)
            if receipt_reporter is not None:
                try:
                    await _maybe_await(receipt_reporter(session_id, outcome="ok", path=path))
                except Exception as exc:  # noqa: BLE001
                    log.warning("ppt_pro receipt report failed: %s", exc)
            if path:
                opened = _open_image_file(str(path))
                if opened:
                    await _notify(notifier, session_id, f"✨ PPT 做好啦，已自动打开：{path}")
                else:
                    # No associated app (e.g. no WPS/PowerPoint installed) or
                    # the open call was refused — don't claim we opened it.
                    await _notify(
                        notifier,
                        session_id,
                        f"✨ PPT 做好啦，已保存到：{path}\n（没能自动打开，麻烦你手动打开看看～）",
                    )
            else:
                await _notify(notifier, session_id, "✨ PPT 做好啦。")
            return

        error = result.get("error") or "未知错误"
        if receipt_reporter is not None:
            try:
                await _maybe_await(receipt_reporter(session_id, outcome="failed", path=result.get("path")))
            except Exception as exc:  # noqa: BLE001
                log.warning("ppt_pro failed receipt report failed: %s", exc)
        await _notify(notifier, session_id, f"PPT 没做成：{error}")
    except Exception as exc:  # noqa: BLE001
        log.warning("ppt_pro report_done failed: %s", exc)


async def _ppt_pro_orchestrate(
    *,
    topic: str,
    pages: int,
    depth: str,
    theme: str,
    image_mode: bool,
    title: str,
    author: str,
    output_path: Optional[str],
    outline_propose,
    notifier,
    run_blocking,
    session_id: str,
) -> None:
    cfg = _ppt_pro_cfg()
    try:
        await _notify(notifier, session_id, "🔍 正在围绕主题做深度调研…")
        report = await _maybe_await(
            _research_topic_for_ppt(topic, depth=depth, timeout_s=cfg.research_timeout_s)
        )
        if cfg.save_research and report is not None:
            _save_and_index_research(topic, report)
        sources = getattr(report, "citations", None) if report is not None else None
        n_src = len(sources or [])
        if report is not None:
            await _notify(notifier, session_id, f"📚 调研完成（{n_src} 个来源），正在拟大纲…")
        else:
            await _notify(notifier, session_id, "📝 调研未取得来源，按通用知识拟大纲…")

        from deskpet.tools.research_tools import _resolve_default_llm_call

        llm = await _resolve_default_llm_call()
        slides = await _draft_outline_from_research(
            topic,
            report,
            pages=pages,
            theme=theme,
            image_mode=image_mode,
            llm_call=llm,
        )

        confirmed = False
        for _ in range(int(cfg.max_revisions) + 1):
            decision = await _maybe_await(
                outline_propose(
                    session_id,
                    topic=topic,
                    slides=slides,
                    sources_count=n_src,
                    outline_md=_outline_to_markdown(slides),
                    no_research=(report is None),
                )
            )
            action = str((decision or {}).get("action") or "modify").strip().lower()
            if action == "accept":
                confirmed = True
                break
            if action == "reuse":
                row = ppt_outline_store.get_outline((decision or {}).get("reuse_id"))
                slides = parse_outline(json.loads(row["slides_json"]))
                confirmed = True
                break
            if action == "cancel":
                await _notify(notifier, session_id, "好的，已取消，没有生成 PPT。")
                return
            slides = await _draft_outline_from_research(
                topic,
                report,
                pages=pages,
                theme=theme,
                image_mode=image_mode,
                llm_call=llm,
                feedback=str((decision or {}).get("feedback") or ""),
                prev_slides=slides,
            )
        if not confirmed:
            await _notify(notifier, session_id, "大纲改了好几轮还没定，先暂停，需要再叫我。")
            return

        await _notify(notifier, session_id, "✅ 大纲已确认，开始生成…")
        render_to = float(cfg.render_timeout_s or max(600, int(pages) * 120))
        loop = asyncio.get_running_loop()
        result = await asyncio.wait_for(
            _run_blocking_call(
                run_blocking,
                lambda: _render_pro(
                    slides,
                    theme=theme,
                    title=title,
                    author=author,
                    output_path=output_path,
                    image_mode=image_mode,
                    probe_timeout_s=cfg.image_probe_timeout_s,
                    notify=_sync_notify(notifier, session_id, loop),
                ),
            ),
            timeout=render_to,
        )
        await _ppt_pro_report_done(result, notifier=notifier, session_id=session_id)
    except asyncio.CancelledError:
        await _notify(notifier, session_id, "PPT 任务已停止。")
        raise
    except Exception as exc:  # noqa: BLE001
        log.exception("ppt_pro orchestrate failed")
        await _ppt_pro_report_done(
            {"ok": False, "error": f"{type(exc).__name__}: {exc}"},
            notifier=notifier,
            session_id=session_id,
        )


def _coerce_ppt_pro_args(args: dict[str, Any]) -> dict[str, Any]:
    cfg = _ppt_pro_cfg()
    topic = str(args.get("topic") or "").strip()
    pages = max(3, min(20, _cfg_int(args, "pages", 8)))
    # 调研档位**锁死为配置值(默认 deep),忽略 LLM 传入的 depth** —— 用户要求
    # "充分调研"，不让 LLM 降档到 standard/light。pro_default_depth 见 [ppt] 配置。
    depth = str(cfg.default_depth or "deep").strip().lower()
    if depth not in {"light", "standard", "deep"}:
        depth = "deep"
    theme = str(args.get("theme") or "minimal").strip().lower()
    if theme not in VALID_THEMES:
        theme = "minimal"
    return {
        "topic": topic,
        "pages": pages,
        "depth": depth,
        "theme": theme,
        "image_mode": bool(args.get("image_mode", True)),
        "title": str(args.get("title") or topic),
        "author": str(args.get("author") or "DeskPet"),
        "output_path": str(args["output_path"]) if args.get("output_path") else None,
    }


async def _handle_ppt_pro(args: dict, task_id: str = "") -> dict[str, Any]:  # noqa: ARG001
    sid = str(args.get("_session_id") or "default")
    outline_propose = _PPT_PRO_CTX.get("outline_propose")
    notifier = _PPT_PRO_CTX.get("notifier")
    run_blocking = _PPT_PRO_CTX.get("run_blocking")
    if outline_propose is None or notifier is None or run_blocking is None:
        return {
            "ok": False,
            "error": "ppt_pro 接线缺失（outline_propose/notifier/run_blocking 未注入）",
            "fallback": "请用 ppt_create",
        }

    kwargs = _coerce_ppt_pro_args(dict(args or {}))
    topic = kwargs["topic"]
    if not topic:
        return {"ok": False, "error": "topic is required"}

    cur = _PPT_PRO_RUNNING.get(sid)
    if cur is not None and not cur.done():
        if topic == _PPT_PRO_RUNNING_TOPIC.get(sid):
            return {
                "ok": True,
                "status": "already_running",
                "message": "这份 PPT 已经在做了，等我把大纲或成品推回来。",
            }
        _ppt_pro_cancel(sid)
        await _notify(notifier, sid, "好，换个主题，我重新来过。")

    async def _runner() -> None:
        try:
            await _ppt_pro_orchestrate(
                **kwargs,
                outline_propose=outline_propose,
                notifier=notifier,
                run_blocking=run_blocking,
                session_id=sid,
            )
        except asyncio.CancelledError:
            await _notify(notifier, sid, "已停止当前 PPT 任务。")
            raise
        except Exception as exc:  # noqa: BLE001
            log.exception("ppt_pro runner failed")
            await _notify(notifier, sid, f"PPT 没做成：{exc}")
        finally:
            me = asyncio.current_task()
            if _PPT_PRO_RUNNING.get(sid) is me:
                _PPT_PRO_RUNNING.pop(sid, None)
                _PPT_PRO_RUNNING_TOPIC.pop(sid, None)

    task = asyncio.create_task(_runner())
    _PPT_PRO_RUNNING[sid] = task
    _PPT_PRO_RUNNING_TOPIC[sid] = topic
    _PPT_PRO_TASKS.add(task)
    task.add_done_callback(_PPT_PRO_TASKS.discard)
    return {
        "ok": True,
        "status": "researching",
        "message": "收到，我先围绕这个主题做调研，拟好大纲会给你确认。",
    }


_PPT_PRO_SCHEMA = {
    "name": "ppt_pro",
    "description": (
        "【做 PPT 首选工具】用户只给一个主题就要做（正式/调研型/惊艳/AI配图）PPT 时，"
        "用这个，不要用 ppt_create、也不要自己先 web_search/todo_write 调研。"
        "它会自动一条龙：① 用 deepresearch 深度调研该主题；② 基于调研拟大纲；"
        "③ 弹「大纲确认卡」给用户确认/修改/取消/复用历史；④ 用户确认后优先用 AI "
        "出图生成惊艳整页配图 PPT，连不上出图服务则自动回退精美模板。只需调一次、传 topic 即可。"
        "返回 status='researching'（或 'already_running'）表示已在后台进行——"
        "【不要】重复调用、不要因没看到成品就重试，安静等后台推进度/成品。"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "topic": {
                "type": "string",
                "description": "PPT topic or brief.",
            },
            "pages": {
                "type": "integer",
                "minimum": 3,
                "maximum": 20,
                "default": 8,
            },
            # depth 不暴露给 LLM —— 调研档位由后端锁死为 [ppt].pro_default_depth(默认 deep)，
            # 不让 LLM 降档，保证"充分调研"。
            "theme": {
                "type": "string",
                "enum": list(VALID_THEMES),
                "default": "minimal",
            },
            "image_mode": {
                "type": "boolean",
                "default": True,
                "description": "True uses AI full-slide images when reachable; false uses templates.",
            },
            "title": {"type": "string"},
            "author": {"type": "string", "default": "DeskPet"},
            "output_path": {
                "type": "string",
                "description": "Optional absolute output .pptx path.",
            },
        },
        "required": ["topic"],
    },
}


def _register_ppt_tool() -> None:
    """Module-import side effect: register ppt_create with the registry.

    Wrapped in a try/except so import-cycles or missing registry don't
    break test collection — every test in this repo imports the tool
    module directly without needing the registry.
    """
    try:
        from .registry import registry  # type: ignore
        registry.register(
            "ppt_create",
            "ppt",
            _build_ppt_schema(),
            _handle_ppt_create,
            permission_category="write_file",
            # 纯文本/模板填充约 1~3s;但带 image_prompt 的整页生图(B-2)会
            # 在 _autofill_image_prompts 里串行调 AI 出图(每张 70~190s,
            # 偶发等满 300s),N 页 deck 可达数分钟。给足预算避免 registry
            # 在生图跑完前杀掉 handler。注:同步阻塞 UX 由后续异步化改善。
            timeout_seconds=1200.0,
            concurrency_safe=False,  # G3: writes .pptx to disk
        )
    except Exception as exc:  # noqa: BLE001
        log.debug("ppt tool registration skipped: %s", exc)


def _register_ppt_pro_tool() -> None:
    try:
        if not _ppt_pro_cfg().enabled:
            return
        from .registry import registry  # type: ignore

        registry.register(
            "ppt_pro",
            "ppt",
            _PPT_PRO_SCHEMA,
            _handle_ppt_pro,
            permission_category="write_file",
            timeout_seconds=60.0,
            concurrency_safe=False,
        )
    except Exception as exc:  # noqa: BLE001
        log.debug("ppt_pro tool registration skipped: %s", exc)


_register_ppt_tool()
_register_ppt_pro_tool()
