# SPDX-FileCopyrightText: 2026 DennyWanye
# SPDX-License-Identifier: BUSL-1.1

"""Manual E2E for PPT + DeepResearch.

What it does
------------
1. Runs ``research_run`` against a fixed topic with a deterministic
   FakeLLM (no live network needed) and prints the resulting report.
2. Pipes the report's section titles into ``ppt_create`` to produce a
   real ``.pptx`` file on disk. Asserts the file exists, slide count,
   and dumps the file path.

Why
---
This is the "smoke test" that proves the modules cooperate end-to-end.
It's separate from pytest so we can run it locally and inspect the
.pptx in PowerPoint / Keynote manually.

Run
---
    cd backend
    python scripts/e2e_ppt_deepresearch.py
"""
from __future__ import annotations

import asyncio
import json
import sys
import tempfile
import time
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from deskpet.tools.research_tools import research_run
from deskpet.tools.ppt_tools import ppt_create


# --------------------------- fake plumbing ----------------------------


class FakeLLM:
    """Two-step canned LLM: plan + synth."""

    def __init__(self):
        self._step = 0

    async def __call__(self, prompt: str) -> str:
        self._step += 1
        if self._step == 1:
            return json.dumps([
                "What is structured logging?",
                "Why does structured logging matter for production?",
                "Which Python libraries support it best?",
                "What are the common pitfalls?",
            ])
        # Synth — keep [^n] count <= n_citations
        return (
            "# Structured logging — a 2026 primer\n\n"
            "## TL;DR\n\n"
            "Structured logging treats log lines as machine-parseable records "
            "rather than free-form strings. Production systems benefit from "
            "filtering, alerting, and observability when logs are JSON or "
            "key-value pairs.\n\n"
            "## Background\n\n"
            "Traditional ``print``-style logs make it hard to grep across "
            "fields like ``request_id`` or ``user_id`` [^1]. Structured "
            "logging fixes this by emitting structured payloads [^2].\n\n"
            "## Current state\n\n"
            "Python's ecosystem has matured: ``structlog`` (this project "
            "uses it) and ``loguru`` are the de-facto picks [^2]. The "
            "stdlib's ``logging`` plus a JSON formatter also works [^1].\n\n"
            "## Open questions\n\n"
            "Verbosity vs cost: structured logs are 3-5× larger than plain "
            "text, which matters at scale [^3].\n\n"
            "## What's next\n\n"
            "OpenTelemetry's logs spec is converging the industry on a "
            "common schema across languages [^4].\n"
        )


async def fake_search(query: str, *, max_results: int = 5):
    canned = {
        "What is structured logging?": [
            {"url": "https://en.wikipedia.org/wiki/Structured_logging",
             "title": "Structured logging — Wikipedia", "snippet": ""},
        ],
        "Why does structured logging matter for production?": [
            {"url": "https://www.structlog.org/en/stable/why.html",
             "title": "Why structlog?", "snippet": ""},
        ],
        "Which Python libraries support it best?": [
            {"url": "https://docs.python.org/3/library/logging.html",
             "title": "logging — Python docs", "snippet": ""},
        ],
        "What are the common pitfalls?": [
            {"url": "https://blog.example.com/log-cost",
             "title": "The cost of verbose logs", "snippet": ""},
        ],
    }
    return canned.get(query, [])[:max_results]


async def fake_extract(url: str):
    now = time.time()
    bodies = {
        "https://en.wikipedia.org/wiki/Structured_logging": (
            "Structured logging is a method of recording log messages "
            "with key-value metadata so they can be filtered, indexed, "
            "and queried programmatically. " * 8
        ),
        "https://www.structlog.org/en/stable/why.html": (
            "structlog binds key-value pairs to the log record so each "
            "message carries enough context for production debugging. "
            "It also supports asyncio and standard logging adapters. " * 6
        ),
        "https://docs.python.org/3/library/logging.html": (
            "Python's built-in logging module supports structured logging "
            "via extra= and custom formatters; combined with json-logging "
            "you can emit JSON lines. " * 5
        ),
        "https://blog.example.com/log-cost": (
            "Verbose structured logs cost roughly 3-5x more storage and "
            "ingestion than plain text equivalents. Sampling and field "
            "selection are mitigations. " * 4
        ),
    }
    if url not in bodies:
        return {"ok": False, "error": "no fixture", "url": url}
    return {
        "ok": True, "url": url,
        "title": url.rsplit("/", 1)[-1].replace("_", " ").title() or url,
        "text": bodies[url],
        "fetched_at": now,
    }


# --------------------------- e2e driver -------------------------------


async def main() -> int:
    topic = "Structured logging in Python — 2026 state of the art"
    print(f"\n===== Step 1: research_run({topic!r}) =====\n")
    report = await research_run(
        topic,
        llm_call=FakeLLM(),
        search=fake_search,
        extract=fake_extract,
        max_sub_questions=4,
        max_urls_per_query=1,
    )
    print(f"sub_questions: {len(report.sub_questions)}")
    for q in report.sub_questions:
        print(f"  - {q}")
    print(f"\ncitations: {len(report.citations)}")
    for c in report.citations:
        print(f"  [{c.n}] {c.title}  ← {c.url}")
    print(f"\ncoverage: {report.coverage}")
    print(f"errors: {report.errors}")
    if not report.coverage.get("cite_check_ok"):
        print("⚠️  cite_check failed — this would normally trigger a retry")
        return 1
    print(f"\n--- report_md (first 400 chars) ---\n{report.report_md[:400]}\n…")

    print("\n===== Step 2: turn report into PPT outline =====\n")
    # Naive section-to-slide conversion (a real skill would let LLM do this)
    outline = [{
        "layout": "title",
        "title": topic,
        "subtitle": f"自动生成 / {len(report.citations)} 个引用",
    }, {
        "layout": "toc", "title": "目录",
        "bullets": ["TL;DR", "Background", "Current state",
                    "Open questions", "What's next"],
    }, {
        "layout": "bullet", "title": "TL;DR",
        "bullets": [report.summary[:60] + "…" if len(report.summary) > 60 else report.summary],
    }]
    # One bullet slide per section heading from report_md
    import re
    for section in re.findall(r"##\s+([^\n]+)\n+([^\n]+(?:\n[^\n#]+)*)", report.report_md):
        heading, body = section
        if heading.strip().lower() == "tl;dr":
            continue
        if heading.strip() == "引用":
            continue
        # Strip footnote markers + take the first 2-3 sentences as bullets
        plain = re.sub(r"\[\^\d+\]", "", body).strip()
        sents = re.split(r"(?<=[。.!?])\s+", plain)
        bullets = [s.strip()[:80] for s in sents[:3] if s.strip()]
        outline.append({
            "layout": "bullet",
            "title": heading.strip(),
            "bullets": bullets,
            "notes": body[:500],
        })
    outline.append({
        "layout": "section",
        "title": "引用",
        "subtitle": f"{len(report.citations)} 个来源",
    })
    # Sources slide
    outline.append({
        "layout": "bullet", "title": "Sources",
        "bullets": [f"[{c.n}] {c.title} — {c.url}" for c in report.citations],
    })

    print(f"outline slides: {len(outline)}")
    out_dir = Path(tempfile.gettempdir()) / "deskpet-e2e"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "structured-logging-2026.pptx"
    result = ppt_create(
        outline,
        theme="minimal",
        title=topic,
        author="DeskPet",
        output_path=str(out_path),
    )
    print(f"\nppt_create → {result}")
    if not result.get("ok"):
        print("PPT failed; markdown fallback printed below:")
        print(result.get("markdown_fallback", ""))
        return 1
    pptx_path = Path(result["path"])
    print(f"\n✅ PPT generated: {pptx_path}")
    print(f"   size: {pptx_path.stat().st_size:,} bytes")
    print(f"   slides: {result['slide_count']}")

    # Round-trip read with python-pptx to prove it's valid
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        assert len(prs.slides) == result["slide_count"]
        first_slide_text = " ".join(
            sh.text_frame.text for sh in prs.slides[0].shapes
            if sh.has_text_frame
        )
        print(f"   round-trip OK — first slide text includes topic? "
              f"{topic in first_slide_text}")
    except Exception as exc:  # noqa: BLE001
        print(f"   ⚠️  round-trip read failed: {exc}")
        return 1

    print(f"\n===== Done. Open the file in PowerPoint:\n    {pptx_path}\n")
    return 0


if __name__ == "__main__":
    sys.exit(asyncio.run(main()))
