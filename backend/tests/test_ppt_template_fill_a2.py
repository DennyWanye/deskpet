# SPDX-FileCopyrightText: 2026 DennyWanye
# SPDX-License-Identifier: BUSL-1.1

"""A-2 PPT template placeholder fill tests."""
from __future__ import annotations

import base64
from pathlib import Path

import pytest

from deskpet.tools.ppt_tools import _HAS_PPTX, ppt_create


pytestmark = pytest.mark.skipif(not _HAS_PPTX, reason="python-pptx not installed")


PNG_1X1 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)


@pytest.fixture
def template_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation

    path = tmp_path / "tpl.pptx"
    Presentation().save(path)
    return path


@pytest.fixture
def tiny_png(tmp_path: Path) -> Path:
    path = tmp_path / "tiny.png"
    path.write_bytes(base64.b64decode(PNG_1X1))
    return path


def _slide_text(slide) -> str:
    return "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))


def _has_picture(slide) -> bool:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return any(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image")
        for shape in slide.shapes
    )


def test_template_fills_title_and_bullets(tmp_path: Path, template_pptx: Path) -> None:
    from pptx import Presentation

    out = tmp_path / "filled.pptx"
    outline = [
        {"layout": "title", "title": "封面", "subtitle": "副标题"},
        {"layout": "bullet", "title": "要点", "bullets": ["A", "B", "C"]},
    ]

    result = ppt_create(outline, output_path=str(out), template=str(template_pptx))

    assert result["ok"] is True
    assert Path(result["path"]).is_file()
    prs = Presentation(result["path"])
    assert "封面" in _slide_text(prs.slides[0])
    slide1_text = _slide_text(prs.slides[1])
    assert "要点" in slide1_text
    assert "A" in slide1_text
    assert "B" in slide1_text
    assert "C" in slide1_text


def test_template_none_unchanged(tmp_path: Path) -> None:
    out = tmp_path / "scratch.pptx"
    outline = [
        {"layout": "title", "title": "封面", "subtitle": "副标题"},
        {"layout": "bullet", "title": "要点", "bullets": ["A"]},
    ]

    result = ppt_create(outline, output_path=str(out), template=None)

    assert result["ok"] is True
    assert result["theme"] == "minimal"
    assert result["theme"] != "template"
    assert result["slide_count"] == 2


def test_template_bad_path_falls_back(tmp_path: Path) -> None:
    out = tmp_path / "fallback.pptx"

    result = ppt_create(
        [{"layout": "bullet", "title": "仍然生成", "bullets": ["A"]}],
        output_path=str(out),
        template="G:\\no\\such.pptx",
    )

    assert result["ok"] is True
    assert result["theme"] == "minimal"
    assert out.is_file()


def test_template_picture_placeholder(tmp_path: Path, template_pptx: Path, tiny_png: Path) -> None:
    from pptx import Presentation

    out = tmp_path / "picture.pptx"

    result = ppt_create(
        [{"layout": "image", "title": "图", "image_path": str(tiny_png)}],
        output_path=str(out),
        template=str(template_pptx),
    )

    assert result["ok"] is True
    prs = Presentation(result["path"])
    assert _has_picture(prs.slides[0])


def test_template_mode_generates_for_swap(
    monkeypatch, tmp_path: Path, template_pptx: Path, tiny_png: Path,
) -> None:
    """模板模式 + image_prompt → 生图(供换进设计页图片位:丰富内容+定制视觉)。"""
    from deskpet.tools import image_tools

    calls: list[list[str]] = []

    def fake_generate_images(prompts, **kwargs):
        calls.append(list(prompts))
        return [{"prompt": prompts[0], "path": str(tiny_png), "error": None}]

    monkeypatch.setattr(image_tools, "generate_images", fake_generate_images)
    out = tmp_path / "prompt-picture.pptx"

    result = ppt_create(
        [{"layout": "bullet", "title": "要点", "bullets": ["A", "B"],
          "image_prompt": "a clean desk pet hero"}],
        output_path=str(out),
        template=str(template_pptx),
    )

    assert result["ok"] is True
    # 模板模式带 image_prompt → 真生图(给 _swap_design_picture 用)
    assert len(calls) == 1 and len(calls[0]) == 1
    assert calls[0][0].startswith("a clean desk pet hero")


def test_template_mode_no_prompt_no_gen(
    monkeypatch, tmp_path: Path, template_pptx: Path, tiny_png: Path,
) -> None:
    """模板模式但没写 image_prompt → 不生图(省钱;用模板自带配图)。"""
    from deskpet.tools import image_tools

    calls: list[list[str]] = []

    def fake_generate_images(prompts, **kwargs):
        calls.append(list(prompts))
        return [{"prompt": prompts[0], "path": str(tiny_png), "error": None}]

    monkeypatch.setattr(image_tools, "generate_images", fake_generate_images)
    out = tmp_path / "no-prompt.pptx"

    result = ppt_create(
        [{"layout": "bullet", "title": "要点", "bullets": ["A", "B"]}],
        output_path=str(out),
        template=str(template_pptx),
    )

    assert result["ok"] is True
    assert calls == []


def test_image_full_autofills_without_template(
    monkeypatch, tmp_path: Path, tiny_png: Path,
) -> None:
    """无 template 的 image_full + image_prompt → 真触发 AI 生图。"""
    from deskpet.tools import image_tools

    calls: list[list[str]] = []

    def fake_generate_images(prompts, **kwargs):
        calls.append(list(prompts))
        return [{"prompt": prompts[0], "path": str(tiny_png), "error": None}]

    monkeypatch.setattr(image_tools, "generate_images", fake_generate_images)
    out = tmp_path / "imagefull.pptx"

    result = ppt_create(
        [{"layout": "image_full", "title": "封面", "image_prompt": "a cinematic ai city"}],
        output_path=str(out),
    )

    assert result["ok"] is True
    assert len(calls) == 1 and len(calls[0]) == 1
    assert calls[0][0].startswith("a cinematic ai city")
