"""Tests for deskpet.tools.ppt_tools — PPT generation.

Coverage:
  * parse_outline tolerates JSON string, code-fenced JSON, list, dict
  * SlideOutline.normalize coerces unknown layout → bullet, strips whitespace
  * ppt_create happy path: file exists, slide_count matches, file
    readable by python-pptx, layouts applied per outline
  * Each of the 7 layouts renders without error
  * Each of the 3 themes renders without error
  * Image layout with missing file falls back to bullet without raising
  * Markdown fallback renders cleanly when outline is non-empty
  * Empty outline → ok:False with error + empty markdown_fallback
  * Garbage outline → ok:False with error
  * Speaker notes survive round-trip
  * Document core properties (title, author) round-trip
"""
from __future__ import annotations

import json
import os
from pathlib import Path

import pytest

from deskpet.tools import ppt_tools as ppt
from deskpet.tools.ppt_tools import (
    SlideOutline,
    parse_outline,
    render_markdown_fallback,
    ppt_create,
    get_theme,
    VALID_LAYOUTS,
    VALID_THEMES,
    _HAS_PPTX,
)


# ----------------------------------------------------------------------
# parse_outline / SlideOutline
# ----------------------------------------------------------------------


def test_parse_outline_plain_json_list() -> None:
    raw = json.dumps([
        {"layout": "title", "title": "Hello", "subtitle": "world"},
        {"layout": "bullet", "title": "T", "bullets": ["a", "b"]},
    ])
    out = parse_outline(raw)
    assert len(out) == 2
    assert out[0].layout == "title"
    assert out[0].subtitle == "world"
    assert out[1].bullets == ["a", "b"]


def test_parse_outline_fenced_json() -> None:
    raw = "```json\n" + json.dumps([{"layout": "section", "title": "X"}]) + "\n```"
    out = parse_outline(raw)
    assert len(out) == 1
    assert out[0].layout == "section"


def test_parse_outline_python_list_of_dicts() -> None:
    out = parse_outline([{"layout": "toc", "bullets": ["a"]}])
    assert len(out) == 1
    assert out[0].layout == "toc"


def test_parse_outline_single_dict_promoted() -> None:
    out = parse_outline({"layout": "title", "title": "X"})
    assert len(out) == 1
    assert out[0].title == "X"


def test_parse_outline_empty_inputs() -> None:
    assert parse_outline(None) == []
    assert parse_outline("") == []
    assert parse_outline("   ") == []
    assert parse_outline("garbage no json") == []


def test_parse_outline_strips_unknown_keys() -> None:
    raw = json.dumps([{"layout": "bullet", "title": "X", "foo": "bar"}])
    out = parse_outline(raw)
    assert len(out) == 1
    assert out[0].title == "X"


def test_outline_normalize_unknown_layout_falls_back_to_bullet() -> None:
    so = SlideOutline(layout="weird", title="x").normalize()
    assert so.layout == "bullet"


def test_outline_normalize_strips_and_drops_empty_bullets() -> None:
    so = SlideOutline(
        layout="bullet", title="  T  ",
        bullets=["a", "", "  ", "b "],
    ).normalize()
    assert so.title == "T"
    assert so.bullets == ["a", "b"]


def test_outline_normalize_handles_missing_fields() -> None:
    so = SlideOutline().normalize()
    assert so.layout == "bullet"
    assert so.title == ""
    assert so.bullets == []


# ----------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------


def test_get_theme_known() -> None:
    for name in VALID_THEMES:
        t = get_theme(name)
        assert t.name == name
        assert len(t.background_rgb) == 3


def test_get_theme_unknown_falls_back_to_minimal() -> None:
    t = get_theme("nonexistent")
    assert t.name == "minimal"


# ----------------------------------------------------------------------
# Markdown fallback
# ----------------------------------------------------------------------


def test_markdown_fallback_renders_layouts() -> None:
    slides = parse_outline([
        {"layout": "title", "title": "Deck", "subtitle": "subtitle"},
        {"layout": "bullet", "title": "Points", "bullets": ["one", "two"]},
        {"layout": "quote", "quote": "Be brave.", "cite": "Anon"},
        {"layout": "two_column", "title": "Cmp",
         "left_title": "L", "left": ["a"],
         "right_title": "R", "right": ["b"]},
    ])
    md = render_markdown_fallback(slides, title="My Deck", author="Me")
    assert "# My Deck" in md
    assert "_作者: Me_" in md
    assert "## 第 1 页 — Deck" in md
    assert "- one" in md
    assert "> Be brave." in md
    assert "**左 — L**" in md
    assert "**右 — R**" in md


def test_markdown_fallback_empty_outline_returns_minimal() -> None:
    md = render_markdown_fallback([], title="X")
    assert md.startswith("# X")


# ----------------------------------------------------------------------
# ppt_create — happy paths
# ----------------------------------------------------------------------


pytestmark_pptx = pytest.mark.skipif(not _HAS_PPTX, reason="python-pptx not installed")


@pytest.fixture
def out_path(tmp_path: Path) -> Path:
    return tmp_path / "deck.pptx"


@pytestmark_pptx
def test_ppt_create_happy_path(out_path: Path) -> None:
    outline = [
        {"layout": "title", "title": "Demo Deck", "subtitle": "by DeskPet"},
        {"layout": "toc", "title": "Agenda", "bullets": ["Intro", "Body", "End"]},
        {"layout": "bullet", "title": "Why", "bullets": ["Save time", "Stay focused"]},
        {"layout": "two_column", "title": "Compare",
         "left_title": "Pros", "left": ["fast", "cheap"],
         "right_title": "Cons", "right": ["limited"]},
        {"layout": "quote", "quote": "Less is more.", "cite": "Mies"},
        {"layout": "section", "title": "Wrap-up"},
    ]
    result = ppt_create(outline, output_path=str(out_path), title="Demo", author="Tester")
    assert result["ok"] is True
    assert result["slide_count"] == 6
    assert result["theme"] == "minimal"
    assert out_path.is_file()
    assert out_path.stat().st_size > 5000  # non-trivially sized .pptx


@pytestmark_pptx
def test_ppt_create_each_theme(out_path: Path) -> None:
    outline = [{"layout": "title", "title": "T"}, {"layout": "bullet", "title": "B", "bullets": ["x"]}]
    for theme in VALID_THEMES:
        p = out_path.with_stem(f"deck-{theme}")
        result = ppt_create(outline, theme=theme, output_path=str(p))
        assert result["ok"] is True
        assert result["theme"] == theme
        assert p.is_file()


@pytestmark_pptx
def test_ppt_create_each_layout(out_path: Path) -> None:
    """All 7 layouts must render without raising."""
    outline = [
        {"layout": "title", "title": "Title slide"},
        {"layout": "section", "title": "Section divider"},
        {"layout": "bullet", "title": "Bullets", "bullets": ["a", "b"]},
        {"layout": "two_column", "title": "Cmp", "left": ["L"], "right": ["R"]},
        {"layout": "image", "title": "Image", "caption": "missing"},
        {"layout": "quote", "quote": "q", "cite": "c"},
        {"layout": "toc", "title": "TOC", "bullets": ["1", "2"]},
    ]
    result = ppt_create(outline, output_path=str(out_path))
    assert result["ok"] is True
    assert result["slide_count"] == 7


@pytestmark_pptx
def test_ppt_create_round_trip_titles_and_count(out_path: Path) -> None:
    """python-pptx can read back the file and titles survive."""
    outline = [
        {"layout": "title", "title": "Round Trip"},
        {"layout": "bullet", "title": "Items", "bullets": ["x", "y"]},
    ]
    result = ppt_create(outline, output_path=str(out_path))
    assert result["ok"] is True

    from pptx import Presentation
    prs = Presentation(str(out_path))
    assert len(prs.slides) == 2
    # Walk first slide text and find our title
    all_text = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            all_text.append(shape.text_frame.text)
    assert any("Round Trip" in t for t in all_text)


@pytestmark_pptx
def test_ppt_create_speaker_notes_round_trip(out_path: Path) -> None:
    outline = [
        {"layout": "title", "title": "X", "notes": "remember to smile"},
    ]
    result = ppt_create(outline, output_path=str(out_path))
    assert result["ok"] is True
    from pptx import Presentation
    prs = Presentation(str(out_path))
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert "remember to smile" in notes


@pytestmark_pptx
def test_ppt_create_core_properties_round_trip(out_path: Path) -> None:
    outline = [{"layout": "title", "title": "X"}]
    ppt_create(outline, output_path=str(out_path), title="Doc T", author="Alice")
    from pptx import Presentation
    prs = Presentation(str(out_path))
    assert prs.core_properties.title == "Doc T"
    assert prs.core_properties.author == "Alice"


# ----------------------------------------------------------------------
# ppt_create — failure / degradation
# ----------------------------------------------------------------------


def test_ppt_create_empty_outline_returns_error_and_fallback() -> None:
    result = ppt_create([], output_path=None)
    assert result["ok"] is False
    assert "error" in result
    assert "markdown_fallback" in result


def test_ppt_create_garbage_outline_returns_error() -> None:
    result = ppt_create("not json at all")
    assert result["ok"] is False
    assert "error" in result
    assert "markdown_fallback" in result


@pytestmark_pptx
def test_ppt_create_image_missing_falls_back_gracefully(out_path: Path) -> None:
    outline = [
        {"layout": "image", "title": "x",
         "image_path": "/nonexistent/path.png",
         "caption": "this should still render"},
    ]
    result = ppt_create(outline, output_path=str(out_path))
    assert result["ok"] is True
    # File should exist + be a valid .pptx
    from pptx import Presentation
    prs = Presentation(str(out_path))
    # text trace mentions the caption + missing marker
    text = " ".join(
        sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame
    )
    assert "image missing" in text or "should still render" in text


@pytestmark_pptx
def test_ppt_create_unknown_theme_uses_minimal(out_path: Path) -> None:
    outline = [{"layout": "title", "title": "X"}]
    result = ppt_create(outline, theme="nonexistent", output_path=str(out_path))
    assert result["ok"] is True
    assert result["theme"] == "minimal"


@pytestmark_pptx
def test_ppt_create_default_output_path_lands_in_temp(tmp_path: Path, monkeypatch) -> None:
    # Force tempdir into our pytest tmp_path
    monkeypatch.setattr("tempfile.gettempdir", lambda: str(tmp_path))
    outline = [{"layout": "title", "title": "X"}]
    result = ppt_create(outline)
    assert result["ok"] is True
    assert Path(result["path"]).is_file()
    assert str(tmp_path) in result["path"]


@pytestmark_pptx
def test_ppt_create_json_string_outline(out_path: Path) -> None:
    raw = json.dumps([
        {"layout": "title", "title": "JSON String"},
        {"layout": "bullet", "title": "B", "bullets": ["only one"]},
    ])
    result = ppt_create(raw, output_path=str(out_path))
    assert result["ok"] is True
    assert result["slide_count"] == 2


# --- T7: chart layout enhancement (beta builtin skills) ----------------

@pytestmark_pptx
def test_t7_2_chart_layout_renders(out_path: Path) -> None:
    outline = [{
        "layout": "chart",
        "title": "季度营收",
        "chart": {
            "type": "bar",
            "categories": ["Q1", "Q2", "Q3"],
            "series": [{"name": "营收", "values": [10, 20, 15]}],
        },
    }]
    result = ppt_create(outline, output_path=str(out_path))
    assert result["ok"] is True
    prs = ppt._Presentation(str(out_path))
    slide = prs.slides[0]
    assert any(shape.has_chart for shape in slide.shapes)


@pytestmark_pptx
def test_t7_3_chart_line_and_pie(out_path: Path) -> None:
    for ctype in ("line", "pie"):
        outline = [{
            "layout": "chart",
            "title": ctype,
            "chart": {
                "type": ctype,
                "categories": ["A", "B"],
                "series": [{"name": "s", "values": [1, 2]}],
            },
        }]
        result = ppt_create(outline, output_path=str(out_path))
        assert result["ok"] is True, ctype


@pytestmark_pptx
def test_t7_4_chart_missing_data_degrades(out_path: Path) -> None:
    # No categories/series → must degrade to bullets, not crash.
    outline = [{"layout": "chart", "title": "空图表", "bullets": ["要点A"]}]
    result = ppt_create(outline, output_path=str(out_path))
    assert result["ok"] is True
    assert result["slide_count"] == 1
