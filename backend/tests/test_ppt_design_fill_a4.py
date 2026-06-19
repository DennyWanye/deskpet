# SPDX-FileCopyrightText: 2026 DennyWanye
# SPDX-License-Identifier: BUSL-1.1

"""A-4 PPT design-page reuse tests."""
from __future__ import annotations

from pathlib import Path

import pytest

from deskpet.tools.ppt_tools import _HAS_PPTX, ppt_create


pytestmark = pytest.mark.skipif(not _HAS_PPTX, reason="python-pptx not installed")


LOREM = "LOREM Presentations are communication tools for clear decisions."


def _add_textbox(slide, text: str, *, left: float, top: float, width: float, height: float, pt: int):
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.text = ""
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(pt)
    return shape


def _make_design_template(path: Path, *, slide_count: int = 4) -> Path:
    from pptx import Presentation

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for idx in range(slide_count):
        slide = prs.slides.add_slide(blank)
        _add_textbox(
            slide,
            f"DESIGN TITLE {idx}",
            left=0.6,
            top=0.4,
            width=8.0,
            height=1.0,
            pt=60,
        )
        _add_textbox(
            slide,
            LOREM,
            left=0.8,
            top=1.8,
            width=3.6,
            height=1.2,
            pt=20,
        )
        _add_textbox(
            slide,
            LOREM,
            left=5.0,
            top=1.8,
            width=3.6,
            height=1.2,
            pt=20,
        )
    prs.save(path)
    return path


def _slide_text(slide) -> str:
    return "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))


def _deck_text(prs) -> str:
    return "\n".join(_slide_text(slide) for slide in prs.slides)


def test_design_fill_replaces_title_and_bullets(tmp_path: Path) -> None:
    from pptx import Presentation

    template = _make_design_template(tmp_path / "design.pptx")
    out = tmp_path / "out.pptx"

    result = ppt_create(
        [
            {"layout": "title", "title": "AI PPT 商业方案", "subtitle": "复用设计页"},
            {"layout": "bullet", "title": "核心卖点", "bullets": ["保留视觉", "替换文字槽", "自动删未用页"]},
        ],
        output_path=str(out),
        template=str(template),
    )

    assert result["ok"] is True
    assert result["theme"] == "template-design"
    prs = Presentation(result["path"])
    assert "AI PPT 商业方案" in _slide_text(prs.slides[0])
    slide1_text = _slide_text(prs.slides[1])
    assert "核心卖点" in slide1_text
    assert "保留视觉" in slide1_text
    assert "替换文字槽" in slide1_text
    assert "自动删未用页" in slide1_text
    assert "LOREM" not in _deck_text(prs)


def test_design_fill_keeps_style(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    template = _make_design_template(tmp_path / "design.pptx")
    out = tmp_path / "styled.pptx"

    result = ppt_create(
        [{"layout": "bullet", "title": "样式保留", "bullets": ["字号不变"]}],
        output_path=str(out),
        template=str(template),
    )

    assert result["ok"] is True
    prs = Presentation(result["path"])
    title_shape = next(shape for shape in prs.slides[0].shapes if "样式保留" in getattr(shape, "text", ""))
    assert title_shape.text_frame.paragraphs[0].runs[0].font.size == Pt(60)


def test_design_fill_page_count(tmp_path: Path) -> None:
    from pptx import Presentation

    template = _make_design_template(tmp_path / "design.pptx", slide_count=4)
    out = tmp_path / "count.pptx"

    result = ppt_create(
        [
            {"layout": "title", "title": "第一页"},
            {"layout": "bullet", "title": "第二页", "bullets": ["A"]},
        ],
        output_path=str(out),
        template=str(template),
    )

    assert result["ok"] is True
    prs = Presentation(result["path"])
    assert len(prs.slides) == 2


@pytest.mark.parametrize("slide_count", [0, 1, 2])
def test_design_fill_falls_back_on_tiny_template(tmp_path: Path, slide_count: int) -> None:
    from pptx import Presentation

    template = tmp_path / f"tiny-{slide_count}.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for _ in range(slide_count):
        prs.slides.add_slide(blank)
    prs.save(template)

    out = tmp_path / f"tiny-{slide_count}-out.pptx"
    result = ppt_create(
        [{"layout": "bullet", "title": "仍然生成", "bullets": ["A"]}],
        output_path=str(out),
        template=str(template),
    )

    assert result["ok"] is True
    assert Path(result["path"]).is_file()


def test_design_fill_two_column(tmp_path: Path) -> None:
    from pptx import Presentation

    template = _make_design_template(tmp_path / "design.pptx")
    out = tmp_path / "two-column.pptx"

    result = ppt_create(
        [{
            "layout": "two_column",
            "title": "左右对比",
            "left_title": "左侧",
            "left": ["左一", "左二"],
            "right_title": "右侧",
            "right": ["右一", "右二"],
        }],
        output_path=str(out),
        template=str(template),
    )

    assert result["ok"] is True
    prs = Presentation(result["path"])
    text = _slide_text(prs.slides[0])
    assert "左侧" in text
    assert "左一" in text
    assert "右侧" in text
    assert "右一" in text
