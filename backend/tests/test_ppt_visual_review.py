# SPDX-FileCopyrightText: 2026 DennyWanye
# SPDX-License-Identifier: BUSL-1.1

"""问题1 视觉评估闭环测试: 渲染→评审→应用修复(换版式/缩字号)→重渲染。
评审 LLM 与 COM 渲染全 mock,验证闭环编排与动作应用。"""
from __future__ import annotations

import json
from pathlib import Path
from types import SimpleNamespace

import pytest

from deskpet.tools import ppt_tools
from deskpet.tools.ppt_tools import (
    _HAS_PPTX,
    _apply_review_actions,
    _visual_review_loop,
    get_theme,
    parse_outline,
)
from deskpet.tools.ppt_visual_review import _parse_review_json

pytestmark = pytest.mark.skipif(not _HAS_PPTX, reason="python-pptx not installed")


# ── 评审 JSON 解析 ──────────────────────────────────────────────

def test_parse_review_json_tolerates_fences():
    text = """好的,以下是质检结果:
```json
[{"page":1,"ok":true,"issues":[],"action":"ok"},
 {"page":2,"ok":false,"issues":["标题压脸"],"action":"change_variant","variant":"split_left"}]
```"""
    out = _parse_review_json(text)
    assert len(out) == 2
    assert out[1]["action"] == "change_variant"
    assert out[1]["variant"] == "split_left"


def test_parse_review_json_garbage_returns_empty():
    assert _parse_review_json("不是 json") == []
    assert _parse_review_json("") == []


# ── 动作应用 ────────────────────────────────────────────────────

def _mk_slides():
    return parse_outline([
        {"layout": "image_full", "title": "封面", "image_variant": "cover"},
        {"layout": "image_full", "title": "内容", "bullets": ["A", "B"],
         "image_variant": "split_right"},
    ])


def test_apply_change_variant():
    slides = _mk_slides()
    changed = _apply_review_actions(slides, [
        {"page": 2, "ok": False, "issues": ["压脸"],
         "action": "change_variant", "variant": "split_left"},
    ])
    assert changed is True
    assert slides[1].image_variant == "split_left"


def test_apply_shrink_text_with_floor():
    slides = _mk_slides()
    r = [{"page": 2, "ok": False, "issues": ["溢出"], "action": "shrink_text"}]
    assert _apply_review_actions(slides, r) is True
    assert slides[1].font_scale == 0.85
    # 反复 shrink 有下限(0.72), 不会缩没
    for _ in range(10):
        _apply_review_actions(slides, r)
    assert slides[1].font_scale > 0.6


def test_apply_ok_or_invalid_no_change():
    slides = _mk_slides()
    assert _apply_review_actions(slides, [
        {"page": 1, "ok": True, "action": "ok"},
        {"page": 99, "ok": False, "action": "change_variant", "variant": "top"},
        {"page": 2, "ok": False, "action": "change_variant", "variant": "bogus"},
    ]) is False
    assert slides[1].image_variant == "split_right"


# ── 闭环编排 ────────────────────────────────────────────────────

def _loop_env(monkeypatch, tmp_path, reviews_rounds):
    """搭闭环测试环境: mock 渲染器(写假png) + mock review_slides(按轮出)。"""
    monkeypatch.setattr(ppt_tools, "_in_pytest", lambda: False)

    shots: list[str] = []

    def fake_render(pptx, out_dir, **kw):
        d = Path(out_dir)
        d.mkdir(parents=True, exist_ok=True)
        p = d / "slide1.png"
        p.write_bytes(b"\x89PNG" + b"x" * 2048)
        shots.append(str(d))
        return [str(p)]

    monkeypatch.setattr(
        ppt_tools, "ppt_render",
        SimpleNamespace(
            com_render_available=lambda: True,
            render_pptx_to_pngs_safe=fake_render,
        ),
        raising=False,
    )

    calls = {"n": 0}

    def fake_review(pngs, meta, **kw):
        i = min(calls["n"], len(reviews_rounds) - 1)
        calls["n"] += 1
        return reviews_rounds[i]

    import deskpet.tools.ppt_visual_review as vr
    monkeypatch.setattr(vr, "review_slides", fake_review)

    rerenders = {"n": 0}
    real_render = ppt_tools._render_fromscratch

    def counting_render(*a, **kw):
        rerenders["n"] += 1
        return real_render(*a, **kw)

    monkeypatch.setattr(ppt_tools, "_render_fromscratch", counting_render)
    return shots, calls, rerenders


def test_loop_fix_then_clean(monkeypatch, tmp_path):
    """第1轮发现问题→修→重渲染;第2轮 clean→停。报告进 result。"""
    slides = _mk_slides()
    out = tmp_path / "deck.pptx"
    ppt_tools._render_fromscratch(slides, get_theme("dark"), out, title="t", author="a")
    reviews_rounds = [
        [{"page": 2, "ok": False, "issues": ["压脸"],
          "action": "change_variant", "variant": "top"}],
        [{"page": 1, "ok": True, "action": "ok"},
         {"page": 2, "ok": True, "action": "ok"}],
    ]
    shots, calls, rerenders = _loop_env(monkeypatch, tmp_path, reviews_rounds)
    result: dict = {"ok": True}
    _visual_review_loop(slides, get_theme("dark"), out,
                        title="t", author="a", result=result)
    assert calls["n"] == 2              # 两轮评审
    assert rerenders["n"] == 1          # 修复后重渲染一次
    assert slides[1].image_variant == "top"
    vrep = result.get("visual_review")
    assert vrep and vrep[0]["issues"] == 1 and vrep[1]["issues"] == 0


def test_loop_clean_first_round_stops(monkeypatch, tmp_path):
    slides = _mk_slides()
    out = tmp_path / "deck.pptx"
    ppt_tools._render_fromscratch(slides, get_theme("dark"), out, title="t", author="a")
    reviews_rounds = [[{"page": 1, "ok": True, "action": "ok"}]]
    shots, calls, rerenders = _loop_env(monkeypatch, tmp_path, reviews_rounds)
    result: dict = {"ok": True}
    _visual_review_loop(slides, get_theme("dark"), out,
                        title="t", author="a", result=result)
    assert calls["n"] == 1
    assert rerenders["n"] == 0


def test_loop_review_unavailable_silent(monkeypatch, tmp_path):
    """评审返回 [](vision 不可用) → 闭环静默跳过,不重渲染不写报告。"""
    slides = _mk_slides()
    out = tmp_path / "deck.pptx"
    ppt_tools._render_fromscratch(slides, get_theme("dark"), out, title="t", author="a")
    shots, calls, rerenders = _loop_env(monkeypatch, tmp_path, [[]])
    result: dict = {"ok": True}
    _visual_review_loop(slides, get_theme("dark"), out,
                        title="t", author="a", result=result)
    assert rerenders["n"] == 0
    assert "visual_review" not in result


def test_loop_skips_non_image_deck(monkeypatch, tmp_path):
    """纯文本 deck(无 image_full) → 不烧 vision。"""
    slides = parse_outline([{"layout": "bullet", "title": "t", "bullets": ["a"]}])
    out = tmp_path / "deck.pptx"
    ppt_tools._render_fromscratch(slides, get_theme("dark"), out, title="t", author="a")
    shots, calls, rerenders = _loop_env(
        monkeypatch, tmp_path,
        [[{"page": 1, "ok": False, "action": "shrink_text"}]],
    )
    result: dict = {"ok": True}
    _visual_review_loop(slides, get_theme("dark"), out,
                        title="t", author="a", result=result)
    assert calls["n"] == 0


def test_loop_config_off(monkeypatch, tmp_path):
    slides = _mk_slides()
    out = tmp_path / "deck.pptx"
    ppt_tools._render_fromscratch(slides, get_theme("dark"), out, title="t", author="a")
    shots, calls, rerenders = _loop_env(
        monkeypatch, tmp_path, [[{"page": 1, "ok": False, "action": "shrink_text"}]],
    )
    monkeypatch.setattr(ppt_tools, "_ppt_visual_review_enabled", lambda: False)
    result: dict = {"ok": True}
    _visual_review_loop(slides, get_theme("dark"), out,
                        title="t", author="a", result=result)
    assert calls["n"] == 0


# ── 模板版闭环(design-fill) ─────────────────────────────────────

def test_template_apply_shrink_and_change_page():
    from deskpet.tools.ppt_tools import _apply_review_actions_template

    slides = _mk_slides()
    banned: dict = {}
    page_map = [3, 7]
    changed = _apply_review_actions_template(slides, [
        {"page": 1, "ok": False, "issues": ["溢出"], "action": "shrink_text"},
        {"page": 2, "ok": False, "issues": ["数字与文字重叠"], "action": "change_page"},
    ], page_map, banned)
    assert changed is True
    assert slides[0].font_scale == 0.85
    assert banned == {1: {7}}  # 第2页(idx1) ban 掉设计页7


def test_template_scale_shape_runs_floor():
    """_scale_shape_runs 有 9pt 下限,无显式字号的 run 不动。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from deskpet.tools.ppt_tools import _scale_shape_runs

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "测试"
    r.font.size = Pt(20)
    _scale_shape_runs(box, 0.85)
    assert box.text_frame.paragraphs[0].runs[0].font.size == Pt(17)
    for _ in range(20):
        _scale_shape_runs(box, 0.5)
    assert box.text_frame.paragraphs[0].runs[0].font.size.pt >= 9


def test_template_loop_change_page_rerenders(monkeypatch, tmp_path):
    """模板闭环: 第1轮 change_page → ban+重渲染(banned 传入);第2轮 clean。"""
    from deskpet.tools.ppt_tools import _visual_review_loop_template

    monkeypatch.setattr(ppt_tools, "_in_pytest", lambda: False)

    def fake_render_pngs(pptx, out_dir, **kw):
        d = Path(out_dir); d.mkdir(parents=True, exist_ok=True)
        p = d / "slide1.png"; p.write_bytes(b"\x89PNG" + b"x" * 2048)
        return [str(p)]

    monkeypatch.setattr(
        ppt_tools, "ppt_render",
        SimpleNamespace(com_render_available=lambda: True,
                        render_pptx_to_pngs_safe=fake_render_pngs),
        raising=False,
    )

    calls = {"n": 0}
    rounds = [
        [{"page": 1, "ok": False, "issues": ["重叠"], "action": "change_page"}],
        [{"page": 1, "ok": True, "action": "ok"}],
    ]

    def fake_review(pngs, meta, **kw):
        assert kw.get("mode") == "template"
        i = min(calls["n"], len(rounds) - 1); calls["n"] += 1
        return rounds[i]

    import deskpet.tools.ppt_visual_review as vr
    monkeypatch.setattr(vr, "review_slides", fake_review)

    rerenders = {"n": 0, "banned": None}

    def fake_design_render(slides, template_path, **kw):
        rerenders["n"] += 1
        rerenders["banned"] = kw.get("banned_pages")
        return {"ok": True, "path": kw.get("out_path"), "page_map": [9],
                "slide_count": 1}

    monkeypatch.setattr(ppt_tools, "_render_with_design_pages", fake_design_render)

    slides = _mk_slides()[:1]
    out = tmp_path / "deck.pptx"
    out.write_bytes(b"PK fake")
    result: dict = {"ok": True, "page_map": [5]}
    _visual_review_loop_template(slides, "tpl.pptx", out,
                                 title="t", author="a", result=result)
    assert calls["n"] == 2
    assert rerenders["n"] == 1
    assert rerenders["banned"] == {0: {5}}   # 第1页 ban 掉设计页5
    assert result["page_map"] == [9]          # 重渲染后映射更新
    assert result["visual_review"][0]["issues"] == 1
    assert result["visual_review"][1]["issues"] == 0
